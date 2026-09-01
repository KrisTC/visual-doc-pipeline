#!/usr/bin/env python3
"""Synthetic regression tests for PPTX folder replacement."""

from __future__ import annotations

from contextlib import redirect_stderr, redirect_stdout
from io import BytesIO, StringIO
import json
from pathlib import Path
import re
from tempfile import TemporaryDirectory
from typing import cast
import unittest
from unittest.mock import patch
from uuid import UUID
import xml.etree.ElementTree as ElementTree
from zipfile import ZIP_DEFLATED, ZipFile

from PIL import Image
from docx import Document
from openpyxl import Workbook, load_workbook
from openpyxl.styles import Alignment, Font, Protection
from openpyxl.worksheet.table import Table
from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
# pypdf does not publish PEP 561 metadata for its generic object model.
from pypdf import PdfReader, PdfWriter
from pypdf.generic import ArrayObject, ByteStringObject, ContentStream, DecodedStreamObject, DictionaryObject, FloatObject, NameObject, NumberObject, TextStringObject
# skia-python does not publish PEP 561 stubs; this is the native rendering boundary.
import skia  # type: ignore[import-not-found]

from pipeline.folder_replacement import (
    FolderReplacementResult,
    parse_include_patterns,
    replace_input_folder,
)
from pipeline.folder_replacement.processor import ProgressFactory, ProgressReporter
from pipeline.folder_replacement.docx import _docx_ocr_backgrounds
from pipeline.folder_replacement.pptx import _pptx_ocr_backgrounds
from pipeline.folder_replacement.pdf import (
    _PdfPaintSpan,
    _PdfShownText,
    _PdfVisualRegion,
    _PdfReplacementSerializationError,
    _pdf_apply_paint_span_bullet_overrides,
    _pdf_apply_legacy_bullet_override,
    _pdf_content_has_annotations,
    _pdf_decode_composite_bytes,
    _pdf_expansion_geometry_is_known,
    _pdf_fitted_region_operations,
    _pdf_is_candidate_bullet_error,
    _pdf_text_advance,
)
from pipeline.folder_replacement.xlsx import _replace_drawing
from pipeline.folder_replacement.docx import _validate_docx_embedded_fonts
from pipeline.bounded_text_layout import (
    BoundedTextBox,
    BoundedTextParagraph,
    BoundedTextRun,
    PortableTextUnsupportedError,
    noto_typefaces,
)
from pipeline.portable_bullet_overrides import LegacyBulletOverride
from pipeline.portable_fonts import static_noto_bytes
from pipeline.ocr import BoundingPolygon, OcrRequest, OcrResult, OcrText, PixelPoint
from pipeline.ocr.provider import LocalContractTestSkip
from pipeline.text_replacement import (
    TextReplacementProvider,
    TextReplacementRequest,
    TextReplacementResult,
)
from pipeline.text_replacement_plugins.character_mask import CharacterMaskProvider



from folder_replacement_test_support import (
    FolderReplacementTestCase,
    _CountingOcrProvider,
    _EmptyOcrProvider,
    _FailingOcrProvider,
    _LowConfidenceOcrProvider,
    _RecordedProgress,
    _RecordingReplacementProvider,
    _VectorOutlineOcrProvider,
    _synthetic_pdf_visual_region,
)

class FolderReplacementPptxTests(FolderReplacementTestCase):
    # Verifies FR-2026-08-27-01.
    def test_uses_an_unambiguous_direct_slide_background_for_embedded_image_ocr(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            presentation = Path(temporary_directory) / "input.pptx"
            with ZipFile(presentation, "w", ZIP_DEFLATED) as archive:
                archive.writestr(
                    "ppt/slides/slide1.xml",
                    """<?xml version=\"1.0\"?>
                    <p:sld xmlns:p=\"http://schemas.openxmlformats.org/presentationml/2006/main\"
                        xmlns:a=\"http://schemas.openxmlformats.org/drawingml/2006/main\">
                      <p:cSld><p:bg><p:bgPr><a:solidFill><a:srgbClr val=\"102030\"/>
                      </a:solidFill></p:bgPr></p:bg></p:cSld>
                    </p:sld>""",
                )
                archive.writestr(
                    "ppt/slides/_rels/slide1.xml.rels",
                    """<?xml version=\"1.0\"?>
                    <Relationships xmlns=\"http://schemas.openxmlformats.org/package/2006/relationships\">
                      <Relationship Id=\"rId1\"
                        Type=\"http://schemas.openxmlformats.org/officeDocument/2006/relationships/image\"
                        Target=\"../media/image1.png\"/>
                    </Relationships>""",
                )
                archive.writestr("ppt/media/image1.png", b"synthetic")

            self.assertEqual({"ppt/media/image1.png": (16, 32, 48)}, _pptx_ocr_backgrounds(presentation))

    # Verifies FR-2026-08-03-14, FR-2026-08-03-15, and FR-2026-08-03-16.
    def test_pptx_basic_layout_replaces_and_explicitly_fits_text_frames(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"
            output_root = root / "output"
            input_root.mkdir()
            source = input_root / "deck.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = "Placeholder text"
            text_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(2), Inches(0.5))
            text_frame = text_box.text_frame
            text_frame.paragraphs[0].text = "Short text"
            text_frame.paragraphs[0].runs[0].font.size = Pt(36)
            paragraph_properties = text_frame.paragraphs[0]._p.get_or_add_pPr()
            default_run_properties = paragraph_properties.get_or_add_defRPr()
            paragraph_properties.remove(default_run_properties)
            paragraph_properties.append(
                parse_xml(
                    '<a:buAutoNum xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
                    'type="arabicPeriod"/>'
                )
            )
            paragraph_properties.append(default_run_properties)
            text_frame.paragraphs[0].line_spacing = Pt(14)
            text_frame.add_paragraph()
            group = slide.shapes.add_group_shape()
            group.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(0.5)).text = "Grouped text"
            table_frame = slide.shapes.add_table(2, 3, Inches(1), Inches(3), Inches(4), Inches(1))
            table = table_frame.table
            table.cell(0, 0).text = "Table text"
            table.cell(0, 1).text = "Fallback text"
            fallback_run = table.cell(0, 1).text_frame.paragraphs[0].runs[0]
            fallback_run.font.name = "Source Fallback Font"
            fallback_run.font.size = Pt(31)
            table.cell(1, 0).text = "Merged table text"
            table.cell(1, 0).merge(table.cell(1, 1))
            table.columns[1].width = 0
            expected_table_widths = tuple(int(column.width) for column in table.columns)
            expected_table_heights = tuple(int(row.height) for row in table.rows)
            presentation.save(str(source))

            result = self._run(
                input_root,
                output_root,
                _EmptyOcrProvider(),
                _RecordingReplacementProvider(replacement_text="Long replacement text " * 500),
                document_text_layout="preserve-basic-layout",
            )

            self.assertEqual(1, result.processed_files)
            self.assertEqual(6, result.replaced_native_text_items)
            output = output_root / "deck.pptx"
            output_presentation = Presentation(str(output))
            preserved_blank_paragraph_shape = next(
                shape
                for shape in output_presentation.slides[0].shapes
                if shape.has_text_frame and len(shape.text_frame.paragraphs) == 2
            )
            self.assertEqual("", preserved_blank_paragraph_shape.text_frame.paragraphs[1].text)
            output_table = next(
                shape.table for shape in output_presentation.slides[0].shapes if shape.has_table
            )
            self.assertEqual(
                expected_table_widths,
                tuple(int(column.width) for column in output_table.columns),
            )
            self.assertEqual(
                expected_table_heights,
                tuple(int(row.height) for row in output_table.rows),
            )
            self.assertIn("Long replacement text", output_table.cell(0, 0).text)
            self.assertIn("Long replacement text", output_table.cell(0, 1).text)
            self.assertTrue(output_table.cell(1, 0).is_merge_origin)
            self.assertIn("Long replacement text", output_table.cell(1, 0).text)
            fallback_output_run = output_table.cell(0, 1).text_frame.paragraphs[0].runs[0]
            self.assertEqual("Source Fallback Font", fallback_output_run.font.name)
            self.assertEqual(31.0, fallback_output_run.font.size.pt)
            with ZipFile(output) as archive:
                slide_xml = archive.read("ppt/slides/slide1.xml")
                package_parts = archive.namelist()
            self.assertIn(b"noAutofit", slide_xml)
            self.assertIn(b"Noto Sans JP", slide_xml)
            self.assertNotIn(b"Placeholder text", slide_xml)
            self.assertNotIn(b"Short text", slide_xml)
            self.assertNotIn(b"Grouped text", slide_xml)
            self.assertFalse(any(part.startswith("ppt/fonts/") for part in package_parts))
            self.assertTrue(
                any(int(size) < 3600 for size in re.findall(rb' sz="(\d+)"', slide_xml))
            )
            self.assertIn(b'sz="100"', slide_xml)
            self._assert_valid_drawingml_font_sizes(slide_xml)
            self._assert_drawingml_paragraph_property_order(slide_xml)

    # Verifies FR-2026-08-04-11 and FR-2026-08-05-01.
    def test_pptx_replaces_reachable_smartart_data_and_editable_wordart(self) -> None:
        for layout_mode in (
            "preserve-source-formatting",
            "preserve-basic-layout",
            "preserve-basic-layout-source-font",
        ):
            with self.subTest(layout_mode=layout_mode), TemporaryDirectory() as temporary_directory:
                root = Path(temporary_directory)
                input_root = root / "input"
                output_root = root / "output"
                input_root.mkdir()
                source = input_root / "deck.pptx"
                presentation = Presentation()
                slide = presentation.slides.add_slide(presentation.slide_layouts[6])
                wordart = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
                wordart.name = "editable-wordart"
                wordart.text = "WordArt source"
                wordart.text_frame.paragraphs[0].runs[0].font.name = "Source WordArt Font"
                wordart.text_frame.paragraphs[0].runs[0].font.size = Pt(24)
                wordart.text_frame._element.bodyPr.set("anchor", "ctr")
                wordart.text_frame._element.bodyPr.append(
                    parse_xml(
                        '<a:prstTxWarp xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
                        'prst="textNoShape"><a:avLst/></a:prstTxWarp>'
                    )
                )
                wordart.text_frame.paragraphs[0].runs[0]._r.get_or_add_rPr().append(
                    parse_xml(
                        '<a:pattFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
                        'prst="pct10"><a:fgClr><a:srgbClr val="112233"/></a:fgClr>'
                        "<a:bgClr><a:srgbClr val=\"FFFFFF\"/></a:bgClr></a:pattFill>"
                    )
                )
                wordart._element.spPr.append(
                    parse_xml(
                        '<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                        '<a:outerShdw blurRad="40000" dist="20000" dir="5400000"/>'
                        "</a:effectLst>"
                    )
                )
                coloured_text = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
                coloured_text.name = "coloured-text"
                coloured_text.text = "Coloured source"
                coloured_run = coloured_text.text_frame.paragraphs[0].runs[0]
                coloured_run.font.name = "Source Coloured Font"
                coloured_run.font.size = Pt(24)
                coloured_run.hyperlink.address = "https://example.invalid"
                run_properties = coloured_run._r.get_or_add_rPr()
                for markup in (
                    '<a:solidFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                    '<a:srgbClr val="112233"/></a:solidFill>',
                    '<a:ln xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" w="12700">'
                    '<a:solidFill><a:srgbClr val="445566"/></a:solidFill></a:ln>',
                    '<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                    '<a:outerShdw blurRad="40000" dist="20000" dir="5400000"/></a:effectLst>',
                    '<a:highlight xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                    '<a:srgbClr val="778899"/></a:highlight>',
                    '<a:uLn xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                    '<a:solidFill><a:srgbClr val="AABBCC"/></a:solidFill></a:uLn>',
                    '<a:uFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                    '<a:solidFill><a:srgbClr val="DDEEFF"/></a:solidFill></a:uFill>',
                ):
                    run_properties.append(parse_xml(markup))
                run_properties.set("lang", "en-GB")
                end_properties = coloured_text.text_frame.paragraphs[0]._p.get_or_add_endParaRPr()
                end_properties.set("lang", "en-GB")
                end_properties.append(
                    parse_xml(
                        '<a:highlight xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                        '<a:srgbClr val="AABBCC"/></a:highlight>'
                    )
                )
                presentation.save(str(source))
                self._add_reachable_smartart_data_part(source)

                provider = _RecordingReplacementProvider(replacement_text="Replaced")
                result = self._run(
                    input_root,
                    output_root,
                    _EmptyOcrProvider(),
                    provider,
                    document_text_layout=layout_mode,
                )

                self.assertEqual(1, result.processed_files)
                self.assertGreaterEqual(result.replaced_native_text_items, 4)
                requested_texts = [request.text for request in provider.requests]
                self.assertEqual(1, requested_texts.count("SmartArt first"))
                self.assertEqual(1, requested_texts.count("SmartArt second"))
                self.assertEqual(1, requested_texts.count("WordArt source"))
                self.assertEqual(1, requested_texts.count("Coloured source"))
                output = output_root / "deck.pptx"
                loaded = Presentation(str(output))
                output_wordart = next(
                    shape
                    for shape in loaded.slides[0].shapes
                    if shape.name == "editable-wordart"
                )
                output_coloured_text = next(
                    shape
                    for shape in loaded.slides[0].shapes
                    if shape.name == "coloured-text"
                )
                self.assertEqual("Replaced", output_wordart.text)
                self.assertEqual("Replaced", output_coloured_text.text)
                expected_wordart_font = (
                    "Source WordArt Font"
                    if layout_mode == "preserve-source-formatting"
                    else "Noto Sans JP"
                )
                expected_coloured_font = (
                    "Source Coloured Font"
                    if layout_mode == "preserve-source-formatting"
                    else "Noto Sans JP"
                )
                self.assertEqual(
                    expected_wordart_font,
                    output_wordart.text_frame.paragraphs[0].runs[0].font.name,
                )
                self.assertEqual(
                    expected_coloured_font,
                    output_coloured_text.text_frame.paragraphs[0].runs[0].font.name,
                )
                if layout_mode != "preserve-source-formatting":
                    self.assertEqual(MSO_AUTO_SIZE.NONE, output_wordart.text_frame.auto_size)
                    self.assertEqual(MSO_AUTO_SIZE.NONE, output_coloured_text.text_frame.auto_size)
                output_end_properties = (
                    output_coloured_text.text_frame.paragraphs[0]._p.endParaRPr
                )
                self.assertIsNotNone(output_end_properties.find(qn("a:highlight")))
                self.assertEqual("en-GB", output_end_properties.get("lang"))
                with ZipFile(output) as archive:
                    smartart_data = archive.read("ppt/diagrams/data1.xml")
                    slide_xml = archive.read("ppt/slides/slide1.xml")
                self.assertEqual(2, smartart_data.count(b"Replaced"))
                self.assertNotIn(b"SmartArt first", smartart_data)
                self.assertNotIn(b"SmartArt second", smartart_data)
                self.assertIn(b"outerShdw", slide_xml)
                self.assertIn(b"prstTxWarp", slide_xml)
                self.assertIn(b"pattFill", slide_xml)
                self.assertIn(b'anchor="ctr"', slide_xml)
                for expected_markup in (
                    b"solidFill",
                    b"ln w=\"12700\"",
                    b"effectLst",
                    b"highlight",
                    b"uLn",
                    b"uFill",
                    b'lang="en-GB"',
                    b"hlinkClick",
                ):
                    self.assertIn(expected_markup, slide_xml)

    # Verifies FR-2026-08-04-05.
    def test_pptx_no_autofit_uses_source_width_and_natural_height(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"
            output_root = root / "output"
            input_root.mkdir()
            source = input_root / "deck.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            no_autofit = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(0.2))
            no_autofit.name = "no-autofit"
            no_autofit.text = "Source text " * 40
            no_autofit.text_frame.paragraphs[0].runs[0].font.size = Pt(24)
            no_autofit.text_frame.auto_size = MSO_AUTO_SIZE.NONE
            no_autofit_dimensions = (int(no_autofit.width), int(no_autofit.height))
            text_to_fit_shape = slide.shapes.add_textbox(
                Inches(1), Inches(2), Inches(6), Inches(0.2)
            )
            text_to_fit_shape.name = "text-to-fit-shape"
            text_to_fit_shape.text = "Source text " * 40
            text_to_fit_shape.text_frame.paragraphs[0].runs[0].font.size = Pt(24)
            text_to_fit_shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            table = slide.shapes.add_table(1, 1, Inches(1), Inches(3), Inches(6), Inches(0.2)).table
            table.cell(0, 0).text = "Source text " * 40
            table.cell(0, 0).text_frame.paragraphs[0].runs[0].font.size = Pt(24)
            presentation.save(str(source))

            self._run(
                input_root,
                output_root,
                _EmptyOcrProvider(),
                _RecordingReplacementProvider(replacement_text="Replacement text " * 25),
                document_text_layout="preserve-basic-layout",
            )

            output_presentation = Presentation(str(output_root / "deck.pptx"))
            output_slide = output_presentation.slides[0]
            output_no_autofit = next(
                shape for shape in output_slide.shapes if shape.name == "no-autofit"
            )
            output_text_to_fit_shape = next(
                shape for shape in output_slide.shapes if shape.name == "text-to-fit-shape"
            )
            output_table = next(shape.table for shape in output_slide.shapes if shape.has_table)
            no_autofit_size = output_no_autofit.text_frame.paragraphs[0].runs[0].font.size.pt
            text_to_fit_shape_size = (
                output_text_to_fit_shape.text_frame.paragraphs[0].runs[0].font.size.pt
            )
            table_size = output_table.cell(0, 0).text_frame.paragraphs[0].runs[0].font.size.pt
            self.assertEqual(no_autofit_dimensions, (int(output_no_autofit.width), int(output_no_autofit.height)))
            self.assertGreater(no_autofit_size, text_to_fit_shape_size)
            self.assertGreater(no_autofit_size, table_size)

    # Verifies FR-2026-08-03-15.
    def test_pptx_basic_layout_retains_vertical_alignment_xml(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"
            output_root = root / "output"
            input_root.mkdir()
            source = input_root / "deck.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            middle = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
            middle.name = "middle-aligned"
            middle.text = "Middle aligned"
            middle.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            default = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(3), Inches(1))
            default.name = "default-aligned"
            default.text = "Default aligned"
            presentation.save(str(source))

            self._run(
                input_root,
                output_root,
                _EmptyOcrProvider(),
                _RecordingReplacementProvider(replacement_text="Replacement"),
                document_text_layout="preserve-basic-layout",
            )

            output_presentation = Presentation(str(output_root / "deck.pptx"))
            output_shapes = {shape.name: shape for shape in output_presentation.slides[0].shapes}
            self.assertEqual("ctr", output_shapes["middle-aligned"].text_frame._element.bodyPr.get("anchor"))
            self.assertIsNone(output_shapes["default-aligned"].text_frame._element.bodyPr.get("anchor"))

    # Verifies FR-2026-08-27-02.
    def test_pptx_source_font_mode_falls_back_when_source_faces_are_unavailable(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"
            output_root = root / "output"
            input_root.mkdir()
            source = input_root / "deck.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            named_source = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
            named_source.name = "named-source"
            named_source.text = "Named source font"
            named_run = named_source.text_frame.paragraphs[0].runs[0]
            named_run.font.name = "Source Presentation Font"
            named_run.font.size = Pt(24)
            unnamed_source = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(3), Inches(1))
            unnamed_source.name = "unnamed-source"
            unnamed_source.text = "Fallback font"
            unnamed_source.text_frame.paragraphs[0].runs[0].font.size = Pt(24)
            themed_source = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(3), Inches(1))
            themed_source.name = "themed-source"
            themed_source.text = "Theme source font"
            themed_run = themed_source.text_frame.paragraphs[0].runs[0]
            themed_run.font.name = "+mj-lt"
            themed_run.font.size = Pt(24)
            presentation.save(str(source))

            self._run(
                input_root,
                output_root,
                _EmptyOcrProvider(),
                _RecordingReplacementProvider(replacement_text="Replacement text " * 20),
                document_text_layout="preserve-basic-layout-source-font",
            )

            output_presentation = Presentation(str(output_root / "deck.pptx"))
            output_shapes = output_presentation.slides[0].shapes
            output_named = next(shape for shape in output_shapes if shape.name == "named-source")
            output_unnamed = next(shape for shape in output_shapes if shape.name == "unnamed-source")
            output_themed = next(shape for shape in output_shapes if shape.name == "themed-source")
            self.assertEqual(
                "Noto Sans JP",
                output_named.text_frame.paragraphs[0].runs[0].font.name,
            )
            self.assertEqual(
                "Noto Sans JP",
                output_unnamed.text_frame.paragraphs[0].runs[0].font.name,
            )
            self.assertEqual(
                "Noto Sans JP",
                output_themed.text_frame.paragraphs[0].runs[0].font.name,
            )
            self.assertEqual(MSO_AUTO_SIZE.NONE, output_named.text_frame.auto_size)

    # Verifies FR-2026-08-04-13.
    def test_replaces_pptx_speaker_notes_in_every_document_text_layout_mode(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"
            output_root = root / "output"
            input_root.mkdir()
            source = input_root / "notes.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = "Slide title"
            presentation.save(str(source))
            self._add_speaker_note_part(source)

            for layout_mode in (
                "preserve-source-formatting",
                "preserve-basic-layout",
                "preserve-basic-layout-source-font",
            ):
                with self.subTest(document_text_layout=layout_mode):
                    destination_root = output_root / layout_mode
                    result = self._run(
                        input_root,
                        destination_root,
                        _EmptyOcrProvider(),
                        _RecordingReplacementProvider(replacement_text="Translated speaker note"),
                        document_text_layout=layout_mode,
                    )

                    self.assertEqual(1, result.processed_files)
                    self.assertGreaterEqual(result.replaced_native_text_items, 1)
                    with ZipFile(destination_root / "notes.pptx") as archive:
                        note_xml = archive.read("ppt/notesSlides/notesSlide1.xml")
                        slide_relationships = archive.read("ppt/slides/_rels/slide1.xml.rels")
                    self.assertIn(b"Translated speaker note", note_xml)
                    self.assertNotIn(b"Speaker note source", note_xml)
                    self.assertIn(b'uri="keep-note-extension"', note_xml)
                    self.assertIn(b"../notesSlides/notesSlide1.xml", slide_relationships)
                    Presentation(str(destination_root / "notes.pptx"))



if __name__ == "__main__":
    unittest.main()
