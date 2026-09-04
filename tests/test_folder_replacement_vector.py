#!/usr/bin/env python3
"""Synthetic regression tests for VECTOR folder replacement."""

from __future__ import annotations

from contextlib import redirect_stderr, redirect_stdout
from io import BytesIO, StringIO
import json
from pathlib import Path
import re
import struct
from tempfile import TemporaryDirectory
from typing import cast
import unittest
from unittest.mock import patch
from uuid import UUID
import xml.etree.ElementTree as ElementTree
from zipfile import ZIP_DEFLATED, ZipFile

from PIL import Image
from docx import Document
from openpyxl import Workbook, load_workbook
from openpyxl.styles import Alignment, Font, Protection
from openpyxl.worksheet.table import Table
from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
# pypdf does not publish PEP 561 metadata for its generic object model.
from pypdf import PdfReader, PdfWriter
from pypdf.generic import ArrayObject, ByteStringObject, ContentStream, DecodedStreamObject, DictionaryObject, FloatObject, NameObject, NumberObject, TextStringObject
# skia-python does not publish PEP 561 stubs; this is the native rendering boundary.
import skia  # type: ignore[import-not-found]

from pipeline.folder_replacement import (
    FolderReplacementResult,
    parse_include_patterns,
    replace_input_folder,
)
from pipeline.folder_replacement.processor import ProgressFactory, ProgressReporter
from pipeline.folder_replacement.docx import _docx_ocr_backgrounds
from pipeline.folder_replacement.pptx import _pptx_ocr_backgrounds
from pipeline.folder_replacement.pdf import (
    _PdfPaintSpan,
    _PdfShownText,
    _PdfVisualRegion,
    _PdfReplacementSerializationError,
    _pdf_apply_paint_span_bullet_overrides,
    _pdf_apply_legacy_bullet_override,
    _pdf_content_has_annotations,
    _pdf_decode_composite_bytes,
    _pdf_expansion_geometry_is_known,
    _pdf_fitted_region_operations,
    _pdf_is_candidate_bullet_error,
    _pdf_text_advance,
)
from pipeline.folder_replacement.xlsx import _replace_drawing
from pipeline.folder_replacement.docx import _validate_docx_embedded_fonts
from pipeline.bounded_text_layout import (
    BoundedTextBox,
    BoundedTextParagraph,
    BoundedTextRun,
    PortableTextUnsupportedError,
    noto_typefaces,
)
from pipeline.portable_bullet_overrides import LegacyBulletOverride
from pipeline.portable_fonts import static_noto_bytes
from pipeline.ocr import BoundingPolygon, OcrRequest, OcrResult, OcrText, PixelPoint
from pipeline.ocr.provider import LocalContractTestSkip
from pipeline.text_replacement import (
    TextReplacementProvider,
    TextReplacementRequest,
    TextReplacementResult,
)
from pipeline.text_replacement_plugins.character_mask import CharacterMaskProvider



from folder_replacement_test_support import (
    FolderReplacementTestCase,
    _CountingOcrProvider,
    _EmptyOcrProvider,
    _FailingOcrProvider,
    _LowConfidenceOcrProvider,
    _RecordedProgress,
    _RecordingReplacementProvider,
    _VectorOutlineOcrProvider,
    _synthetic_pdf_visual_region,
)

class FolderReplacementVectorTests(FolderReplacementTestCase):
    # Verifies FR-2026-08-03-05.
    def test_replaces_editable_svg_vector_text_without_ocr(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"
            output_root = root / "output"
            input_root.mkdir()
            with ZipFile(input_root / "document.docx", "w", ZIP_DEFLATED) as archive:
                archive.writestr(
                    "word/media/vector.svg",
                    b'<svg xmlns="http://www.w3.org/2000/svg"><text>Mask</text></svg>',
                )
            ocr_provider = _CountingOcrProvider()

            result = self._run(
                input_root, output_root, ocr_provider, _RecordingReplacementProvider()
            )

            self.assertEqual(1, result.processed_files)
            self.assertEqual(1, result.replaced_native_text_items)
            self.assertEqual(0, result.replaced_image_regions)
            self.assertEqual(0, result.retained_vector_graphics)
            self.assertEqual(0, ocr_provider.calls)
            with ZipFile(output_root / "document.docx") as archive:
                vector_data = archive.read("word/media/vector.svg")
            self.assertIn(b"####", vector_data)
            self.assertNotIn(b"Mask", vector_data)

    # Verifies FR-2026-08-03-05.
    def test_retains_noneditable_vector_and_reports_it(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"
            output_root = root / "output"
            input_root.mkdir()
            vector_data = b'<svg xmlns="http://www.w3.org/2000/svg"><path d="M0 0"/></svg>'
            with ZipFile(input_root / "document.docx", "w", ZIP_DEFLATED) as archive:
                archive.writestr("word/media/vector.svg", vector_data)

            standard_error = StringIO()
            with redirect_stderr(standard_error):
                result = self._run(
                    input_root, output_root, _EmptyOcrProvider(), _RecordingReplacementProvider()
                )

            self.assertEqual(1, result.processed_files)
            self.assertEqual(1, result.retained_vector_graphics)
            self.assertIn("Retained vector", standard_error.getvalue())
            with ZipFile(output_root / "document.docx") as archive:
                self.assertEqual(vector_data, archive.read("word/media/vector.svg"))

    # Verifies FR-2026-09-04-01.
    def test_embedded_emf_receives_requested_basic_layout_mode(self) -> None:
        """The office-package route must preserve EMF fitting options."""
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"
            output_root = root / "output"
            input_root.mkdir()
            source = input_root / "document.pptx"
            presentation = Presentation()
            presentation.slides.add_slide(presentation.slide_layouts[6])
            presentation.save(str(source))
            with ZipFile(source) as archive:
                parts = {entry.filename: archive.read(entry) for entry in archive.infolist()}
            parts["ppt/media/vector.emf"] = _unclipped_emf_with_vertical_rule()
            _add_pptx_emf_relationship(parts)
            with ZipFile(source, "w", ZIP_DEFLATED) as archive:
                for name, data in parts.items():
                    archive.writestr(name, data)

            result = self._run(
                input_root,
                output_root,
                _EmptyOcrProvider(),
                _LongReplacementProvider(),
                document_text_layout="preserve-basic-layout",
            )
            self.assertEqual([], result.failures)

            output_files = list(output_root.glob("*.pptx"))
            self.assertEqual(1, len(output_files))
            with ZipFile(output_files[0]) as archive:
                emf_data = archive.read("ppt/media/vector.emf")
            self.assertEqual((0.0, 0.0), _emf_text_scales(emf_data)[0])
            font_heights = _emf_font_heights(emf_data)
            self.assertLess(abs(font_heights[0]), abs(font_heights[1]))

class _LongReplacementProvider:
    def replace(self, request: TextReplacementRequest) -> TextReplacementResult:
        return TextReplacementResult("A considerably longer heading", 1.0)


def _unclipped_emf_with_vertical_rule() -> bytes:
    records = (
        _emf_font_record("Noto Sans JP"),
        struct.pack("<III", 37, 12, 1),
        struct.pack("<IIii", 27, 16, 60, -10),
        struct.pack("<IIii", 54, 16, 60, 30),
        _emf_exttextout_record("Old", 0, 0, 40, 20),
    )
    header = bytearray(88)
    struct.pack_into("<II", header, 0, 1, len(header))
    eof = struct.pack("<IIIII", 14, 20, 0, 0, 0)
    result = bytearray(header + b"".join(records) + eof)
    struct.pack_into("<I", result, 48, len(result))
    struct.pack_into("<I", result, 52, len(records) + 2)
    return bytes(result)


def _emf_font_record(family: str) -> bytes:
    record = bytearray(104)
    struct.pack_into("<II", record, 0, 82, len(record))
    struct.pack_into("<I", record, 8, 1)
    struct.pack_into("<iiiii", record, 12, -20, 0, 0, 0, 400)
    record[40:104] = family.encode("utf-16-le").ljust(64, b"\0")
    return bytes(record)


def _emf_exttextout_record(text: str, left: int, top: int, right: int, bottom: int) -> bytes:
    text_bytes = text.encode("utf-16-le")
    record = bytearray(76 + len(text_bytes))
    struct.pack_into("<II", record, 0, 84, len(record))
    struct.pack_into("<iiii", record, 8, left, top, right, bottom)
    struct.pack_into("<ii", record, 36, left, top)
    struct.pack_into("<I", record, 44, len(text))
    struct.pack_into("<I", record, 48, 76)
    record[76:] = text_bytes
    record.extend(b"\0" * ((-len(record)) % 4))
    struct.pack_into("<I", record, 4, len(record))
    return bytes(record)


def _emf_text_scales(data: bytes) -> list[tuple[float, float]]:
    scales: list[tuple[float, float]] = []
    offset = 0
    while offset < len(data):
        record_type, record_size = struct.unpack_from("<II", data, offset)
        if record_type == 84:
            scales.append(struct.unpack_from("<ff", data, offset + 28))
        offset += record_size
    return scales


def _emf_font_heights(data: bytes) -> list[int]:
    heights: list[int] = []
    offset = 0
    while offset < len(data):
        record_type, record_size = struct.unpack_from("<II", data, offset)
        if record_type == 82:
            heights.append(struct.unpack_from("<i", data, offset + 12)[0])
        offset += record_size
    return heights


def _add_pptx_emf_relationship(parts: dict[str, bytes]) -> None:
    relationship_part = "ppt/slides/_rels/slide1.xml.rels"
    relationships = ElementTree.fromstring(parts[relationship_part])
    ElementTree.SubElement(
        relationships,
        "{http://schemas.openxmlformats.org/package/2006/relationships}Relationship",
        Id="rId99",
        Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/image",
        Target="../media/vector.emf",
    )
    parts[relationship_part] = ElementTree.tostring(
        relationships, encoding="utf-8", xml_declaration=True
    )
    content_types = ElementTree.fromstring(parts["[Content_Types].xml"])
    ElementTree.SubElement(
        content_types,
        "{http://schemas.openxmlformats.org/package/2006/content-types}Default",
        Extension="emf",
        ContentType="image/x-emf",
    )
    parts["[Content_Types].xml"] = ElementTree.tostring(
        content_types, encoding="utf-8", xml_declaration=True
    )


if __name__ == "__main__":
    unittest.main()
