#!/usr/bin/env python3
"""Synthetic regression tests for the folder replacement command."""

from __future__ import annotations

from contextlib import redirect_stderr, redirect_stdout
from io import BytesIO, StringIO
import json
from pathlib import Path
import re
from tempfile import TemporaryDirectory
from typing import cast
import unittest
from unittest.mock import patch
from uuid import UUID
import xml.etree.ElementTree as ElementTree
from zipfile import ZIP_DEFLATED, ZipFile

from PIL import Image
from docx import Document
from openpyxl import Workbook, load_workbook
from openpyxl.styles import Alignment, Font, Protection
from openpyxl.worksheet.table import Table
from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
# pypdf does not publish PEP 561 metadata for its generic object model.
from pypdf import PdfReader, PdfWriter
from pypdf.generic import ArrayObject, ByteStringObject, ContentStream, DecodedStreamObject, DictionaryObject, FloatObject, NameObject, NumberObject, TextStringObject
# skia-python does not publish PEP 561 stubs; this is the native rendering boundary.
import skia  # type: ignore[import-not-found]

from pipeline.folder_replacement import (
    FolderReplacementResult,
    parse_include_patterns,
    replace_input_folder,
)
from pipeline.folder_replacement.processor import ProgressFactory, ProgressReporter
from pipeline.folder_replacement.docx import _docx_ocr_backgrounds
from pipeline.folder_replacement.pptx import _pptx_ocr_backgrounds
from pipeline.folder_replacement.pdf import (
    _PdfReplacementSerializationError,
    _pdf_decode_composite_bytes,
    _pdf_text_advance,
)
from pipeline.folder_replacement.xlsx import _replace_drawing
from pipeline.folder_replacement.docx import _validate_docx_embedded_fonts
from pipeline.bounded_text_layout import PortableTextUnsupportedError, noto_typefaces
from pipeline.portable_fonts import static_noto_bytes
from pipeline.ocr import BoundingPolygon, OcrRequest, OcrResult, OcrText, PixelPoint
from pipeline.ocr.provider import LocalContractTestSkip
from pipeline.text_replacement import (
    TextReplacementProvider,
    TextReplacementRequest,
    TextReplacementResult,
)
from pipeline.text_replacement_plugins.character_mask import CharacterMaskProvider


FONT_PATH = Path(__file__).parent / "assets" / "fonts" / "NotoSansJP[wght].ttf"


class _EmptyOcrProvider:
    supported_languages = frozenset({"en"})
    supports_local_contract_test = False
    skipped_local_contract_angles: frozenset[int] = frozenset()
    skipped_local_contract_cases: frozenset[LocalContractTestSkip] = frozenset()

    def recognize(self, request: OcrRequest) -> OcrResult:
        return OcrResult(())


class _CountingOcrProvider(_EmptyOcrProvider):
    def __init__(self) -> None:
        self.calls = 0

    def recognize(self, request: OcrRequest) -> OcrResult:
        self.calls += 1
        return OcrResult(())


class _FailingOcrProvider(_EmptyOcrProvider):
    def recognize(self, request: OcrRequest) -> OcrResult:
        raise RuntimeError("Synthetic OCR failed.") from ValueError(
            "Synthetic OCR chained detail must not be recorded."
        )


class _LowConfidenceOcrProvider(_EmptyOcrProvider):
    def recognize(self, request: OcrRequest) -> OcrResult:
        return OcrResult(
            (
                OcrText(
                    "skip me",
                    0.64,
                    BoundingPolygon(
                        (
                            PixelPoint(1, 1),
                            PixelPoint(20, 1),
                            PixelPoint(20, 10),
                            PixelPoint(1, 10),
                        )
                    ),
                ),
            )
        )


class _RecordingReplacementProvider:
    def __init__(self, filename: str | None = None, replacement_text: str | None = None) -> None:
        self.filename = filename
        self.replacement_text = replacement_text
        self.requests: list[TextReplacementRequest] = []

    def replace(self, request: TextReplacementRequest) -> TextReplacementResult:
        self.requests.append(request)
        if request.is_filename and self.filename is not None:
            return TextReplacementResult(self.filename, 1.0)
        if request.is_filename:
            return TextReplacementResult(request.text, 1.0)
        return TextReplacementResult(self.replacement_text or "#" * len(request.text), 1.0)


class _FailingReplacementProvider(_RecordingReplacementProvider):
    def replace(self, request: TextReplacementRequest) -> TextReplacementResult:
        self.requests.append(request)
        if request.is_filename:
            return TextReplacementResult(request.text, 1.0)
        raise RuntimeError("Synthetic text replacement failed.") from ValueError(
            "Synthetic chained detail must not be recorded."
        )


class _RecordedProgress:
    def __init__(self, total: int, label: str) -> None:
        self.total = total
        self.label = label
        self.postfixes: list[str] = []
        self.updates = 0
        self.closed = False

    def set_postfix_str(self, text: str) -> None:
        self.postfixes.append(text)

    def update(self, count: float | None = None) -> bool | None:
        self.updates += 1 if count is None else int(count)
        return None

    def close(self) -> None:
        self.closed = True


class FolderReplacementTests(unittest.TestCase):
    # Verifies FR-2026-08-27-01.
    def test_uses_a_direct_document_background_for_embedded_word_image_ocr(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            document = Path(temporary_directory) / "input.docx"
            with ZipFile(document, "w", ZIP_DEFLATED) as archive:
                archive.writestr(
                    "word/document.xml",
                    """<?xml version=\"1.0\"?>
                    <w:document xmlns:w=\"http://schemas.openxmlformats.org/wordprocessingml/2006/main\">
                      <w:background w:color=\"102030\"/>
                    </w:document>""",
                )
                archive.writestr("word/media/image1.png", b"synthetic")

            self.assertEqual({"word/media/image1.png": (16, 32, 48)}, _docx_ocr_backgrounds(document))

    # Verifies FR-2026-08-27-01.
    def test_uses_an_unambiguous_direct_slide_background_for_embedded_image_ocr(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            presentation = Path(temporary_directory) / "input.pptx"
            with ZipFile(presentation, "w", ZIP_DEFLATED) as archive:
                archive.writestr(
                    "ppt/slides/slide1.xml",
                    """<?xml version=\"1.0\"?>
                    <p:sld xmlns:p=\"http://schemas.openxmlformats.org/presentationml/2006/main\"
                        xmlns:a=\"http://schemas.openxmlformats.org/drawingml/2006/main\">
                      <p:cSld><p:bg><p:bgPr><a:solidFill><a:srgbClr val=\"102030\"/>
                      </a:solidFill></p:bgPr></p:bg></p:cSld>
                    </p:sld>""",
                )
                archive.writestr(
                    "ppt/slides/_rels/slide1.xml.rels",
                    """<?xml version=\"1.0\"?>
                    <Relationships xmlns=\"http://schemas.openxmlformats.org/package/2006/relationships\">
                      <Relationship Id=\"rId1\"
                        Type=\"http://schemas.openxmlformats.org/officeDocument/2006/relationships/image\"
                        Target=\"../media/image1.png\"/>
                    </Relationships>""",
                )
                archive.writestr("ppt/media/image1.png", b"synthetic")

            self.assertEqual({"ppt/media/image1.png": (16, 32, 48)}, _pptx_ocr_backgrounds(presentation))

    # Verifies FR-2026-08-22-02.
    def test_include_patterns_select_relative_supported_paths(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"
            output_root = root / "output"
            selected_directory = input_root / "selected"
            selected_directory.mkdir(parents=True)
            self._write_png(selected_directory / "keep.png")
            self._write_png(input_root / "skip.png")

            result = self._run(
                input_root,
                output_root,
                _EmptyOcrProvider(),
                _RecordingReplacementProvider(),
                include_patterns=("selected/*.png",),
            )

            self.assertEqual(1, result.processed_files)
            self.assertEqual(1, result.ignored_files)
            self.assertTrue((output_root / "selected" / "keep.png").is_file())
            self.assertFalse((output_root / "skip.png").exists())
            self.assertFalse((output_root / "skip.png.diagnostics.json").exists())

    # Verifies FR-2026-08-22-02.
    def test_include_patterns_allow_comma_separated_repeated_values_and_zero_matches(self) -> None:
        self.assertEqual(
            ("*.png", "nested/*.pdf", "*.docx"),
            parse_include_patterns(("*.png,nested/*.pdf", "*.docx")),
        )
        with self.assertRaisesRegex(ValueError, "must not be empty"):
            parse_include_patterns(("*.png,",))

        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"
            output_root = root / "output"
            input_root.mkdir()
            self._write_png(input_root / "source.png")

            result = self._run(
                input_root,
                output_root,
                _EmptyOcrProvider(),
                _RecordingReplacementProvider(),
                include_patterns=("*.pdf",),
            )

            self.assertEqual(0, result.processed_files)
            self.assertEqual(1, result.ignored_files)
            self.assertFalse(output_root.exists())

    # Verifies FR-2026-08-04-07.
    def test_xlsx_drawing_layout_writes_schema_ordered_run_properties(self) -> None:
        drawing = b'''<xdr:wsDr xmlns:xdr="http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"><xdr:sp><xdr:spPr><a:xfrm><a:ext cx="914400" cy="457200"/></a:xfrm></xdr:spPr><xdr:txBody><a:bodyPr/><a:lstStyle/><a:p><a:r><a:rPr sz="2400"><a:latin typeface="Source Sans"/></a:rPr><a:t>Original</a:t></a:r><a:endParaRPr/></a:p></xdr:txBody></xdr:sp></xdr:wsDr>'''
        updated, count = _replace_drawing(
            drawing, _RecordingReplacementProvider(replacement_text="A longer replacement"), "en", "en",
            noto_typefaces(),
            True,
        )
        self.assertEqual(1, count)
        root = ElementTree.fromstring(updated)
        paragraph = next(item for item in root.iter() if item.tag.endswith("}p"))
        children = list(paragraph)
        self.assertLess(next(index for index, item in enumerate(children) if item.tag.endswith("}r")), next(index for index, item in enumerate(children) if item.tag.endswith("}endParaRPr")))
        run_properties = next(item for item in root.iter() if item.tag.endswith("}rPr"))
        self.assertNotIn("typeface", run_properties.attrib)
        self.assertTrue(any(item.tag.endswith("}latin") for item in run_properties))

    # Verifies FR-2026-08-04-07.
    def test_pdf_basic_layout_embeds_a_fitted_font_for_bounded_freetext(self) -> None:
        with TemporaryDirectory() as temporary:
            input_root = Path(temporary) / "input"; input_root.mkdir()
            output_root = Path(temporary) / "output"
            source = input_root / "annotation.pdf"
            writer = PdfWriter(); page = writer.add_blank_page(200, 100)
            annotation = DictionaryObject({
                NameObject("/Type"): NameObject("/Annot"), NameObject("/Subtype"): NameObject("/FreeText"),
                NameObject("/Rect"): ArrayObject([NumberObject(10), NumberObject(10), NumberObject(100), NumberObject(30)]),
                NameObject("/Contents"): TextStringObject("short"), NameObject("/DA"): TextStringObject("/Helv 16 Tf 0 g"),
            })
            page[NameObject("/Annots")] = ArrayObject([writer._add_object(annotation)])
            with source.open("wb") as handle: writer.write(handle)

            result = self._run(input_root, output_root, _EmptyOcrProvider(), _RecordingReplacementProvider(replacement_text="A longer fitted annotation"), document_text_layout="preserve-basic-layout")

            self.assertEqual(1, result.processed_files)
            from pypdf import PdfReader
            output = PdfReader(output_root / "annotation.pdf")
            annotations = cast(ArrayObject, output.pages[0].get("/Annots"))
            annotation_output = cast(DictionaryObject, annotations[0].get_object())
            self.assertIn("longer fitted", str(annotation_output.get("/Contents")))
            annotation_appearance = annotation_output.get("/AP")
            assert annotation_appearance is not None
            appearance_dictionary = cast(DictionaryObject, annotation_appearance.get_object())
            normal_appearance = appearance_dictionary.get("/N")
            assert normal_appearance is not None
            appearance = cast(DictionaryObject, normal_appearance.get_object())
            resources = cast(DictionaryObject, appearance.get("/Resources"))
            fonts = cast(DictionaryObject, resources.get("/Font"))
            self.assertIn("/PipelineNoto", str(resources))
            embedded_font = fonts.get("/PipelineNoto")
            assert embedded_font is not None
            type_zero = cast(DictionaryObject, embedded_font.get_object())
            descendants = type_zero.get("/DescendantFonts")
            assert isinstance(descendants, ArrayObject)
            descendant = cast(DictionaryObject, descendants[0].get_object())
            font_descriptor = descendant.get("/FontDescriptor")
            assert font_descriptor is not None
            descriptor = cast(DictionaryObject, font_descriptor.get_object())
            self.assertIsNotNone(descriptor.get("/FontFile2"))

    # Verifies FR-2026-08-03-03.
    def test_processes_supported_files_ignores_others_and_resolves_filename_collisions(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"
            output_root = root / "output"
            input_root.mkdir()
            self._write_png(input_root / "first.png")
            self._write_png(input_root / "second.png")
            (input_root / "ignore.txt").write_text("ignore", encoding="utf-8")
            provider = _RecordingReplacementProvider("same.png")

            result = self._run(input_root, output_root, _EmptyOcrProvider(), provider)

            self.assertEqual(2, result.processed_files)
            self.assertEqual(1, result.ignored_files)
            self.assertEqual(0, result.failed_files)
            self.assertTrue((output_root / "same.png").is_file())
            self.assertTrue((output_root / "same (2).png").is_file())
            self.assertFalse((output_root / "ignore.txt.diagnostics.json").exists())

    # Verifies FR-2026-08-03-03.
    def test_replaces_office_native_text_in_word_drawing_and_spreadsheet_parts(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"
            output_root = root / "output"
            input_root.mkdir()
            source = input_root / "document.docx"
            with ZipFile(source, "w", ZIP_DEFLATED) as archive:
                archive.writestr(
                    "word/document.xml",
                    """<?xml version=\"1.0\"?><w:document xmlns:w=\"http://schemas.openxmlformats.org/wordprocessingml/2006/main\"><w:body><w:p><w:r><w:t>Word</w:t></w:r></w:p></w:body></w:document>""",
                )
                archive.writestr(
                    "ppt/slides/slide1.xml",
                    """<?xml version=\"1.0\"?><p:sld xmlns:p=\"http://schemas.openxmlformats.org/presentationml/2006/main\" xmlns:a=\"http://schemas.openxmlformats.org/drawingml/2006/main\"><a:t>Slide</a:t></p:sld>""",
                )
                archive.writestr(
                    "xl/sharedStrings.xml",
                    """<?xml version=\"1.0\"?><sst xmlns=\"http://schemas.openxmlformats.org/spreadsheetml/2006/main\"><si><t>Cell</t></si></sst>""",
                )
            result = self._run(
                input_root, output_root, _EmptyOcrProvider(), _RecordingReplacementProvider()
            )

            self.assertEqual(1, result.processed_files)
            self.assertEqual(3, result.replaced_native_text_items)
            with ZipFile(output_root / "document.docx") as archive:
                payload = "\n".join(
                    archive.read(part).decode("utf-8")
                    for part in (
                        "word/document.xml",
                        "ppt/slides/slide1.xml",
                        "xl/sharedStrings.xml",
                    )
                )
            self.assertNotIn("Word", payload)
            self.assertNotIn("Slide", payload)
            self.assertNotIn("Cell", payload)
            self.assertIn("####", payload)
            self.assertIn("#####", payload)

    # Verifies FR-2026-08-03-14, FR-2026-08-03-15, and FR-2026-08-03-16.
    def test_pptx_basic_layout_replaces_and_explicitly_fits_text_frames(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"
            output_root = root / "output"
            input_root.mkdir()
            source = input_root / "deck.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = "Placeholder text"
            text_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(2), Inches(0.5))
            text_frame = text_box.text_frame
            text_frame.paragraphs[0].text = "Short text"
            text_frame.paragraphs[0].runs[0].font.size = Pt(36)
            paragraph_properties = text_frame.paragraphs[0]._p.get_or_add_pPr()
            default_run_properties = paragraph_properties.get_or_add_defRPr()
            paragraph_properties.remove(default_run_properties)
            paragraph_properties.append(
                parse_xml(
                    '<a:buAutoNum xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
                    'type="arabicPeriod"/>'
                )
            )
            paragraph_properties.append(default_run_properties)
            text_frame.paragraphs[0].line_spacing = Pt(14)
            text_frame.add_paragraph()
            group = slide.shapes.add_group_shape()
            group.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(0.5)).text = "Grouped text"
            table_frame = slide.shapes.add_table(2, 3, Inches(1), Inches(3), Inches(4), Inches(1))
            table = table_frame.table
            table.cell(0, 0).text = "Table text"
            table.cell(0, 1).text = "Fallback text"
            fallback_run = table.cell(0, 1).text_frame.paragraphs[0].runs[0]
            fallback_run.font.name = "Source Fallback Font"
            fallback_run.font.size = Pt(31)
            table.cell(1, 0).text = "Merged table text"
            table.cell(1, 0).merge(table.cell(1, 1))
            table.columns[1].width = 0
            expected_table_widths = tuple(int(column.width) for column in table.columns)
            expected_table_heights = tuple(int(row.height) for row in table.rows)
            presentation.save(str(source))

            result = self._run(
                input_root,
                output_root,
                _EmptyOcrProvider(),
                _RecordingReplacementProvider(replacement_text="Long replacement text " * 500),
                document_text_layout="preserve-basic-layout",
            )

            self.assertEqual(1, result.processed_files)
            self.assertEqual(6, result.replaced_native_text_items)
            output = output_root / "deck.pptx"
            output_presentation = Presentation(str(output))
            preserved_blank_paragraph_shape = next(
                shape
                for shape in output_presentation.slides[0].shapes
                if shape.has_text_frame and len(shape.text_frame.paragraphs) == 2
            )
            self.assertEqual("", preserved_blank_paragraph_shape.text_frame.paragraphs[1].text)
            output_table = next(
                shape.table for shape in output_presentation.slides[0].shapes if shape.has_table
            )
            self.assertEqual(
                expected_table_widths,
                tuple(int(column.width) for column in output_table.columns),
            )
            self.assertEqual(
                expected_table_heights,
                tuple(int(row.height) for row in output_table.rows),
            )
            self.assertIn("Long replacement text", output_table.cell(0, 0).text)
            self.assertIn("Long replacement text", output_table.cell(0, 1).text)
            self.assertTrue(output_table.cell(1, 0).is_merge_origin)
            self.assertIn("Long replacement text", output_table.cell(1, 0).text)
            fallback_output_run = output_table.cell(0, 1).text_frame.paragraphs[0].runs[0]
            self.assertEqual("Source Fallback Font", fallback_output_run.font.name)
            self.assertEqual(31.0, fallback_output_run.font.size.pt)
            with ZipFile(output) as archive:
                slide_xml = archive.read("ppt/slides/slide1.xml")
                package_parts = archive.namelist()
            self.assertIn(b"noAutofit", slide_xml)
            self.assertIn(b"Noto Sans JP", slide_xml)
            self.assertNotIn(b"Placeholder text", slide_xml)
            self.assertNotIn(b"Short text", slide_xml)
            self.assertNotIn(b"Grouped text", slide_xml)
            self.assertFalse(any(part.startswith("ppt/fonts/") for part in package_parts))
            self.assertTrue(
                any(int(size) < 3600 for size in re.findall(rb' sz="(\d+)"', slide_xml))
            )
            self.assertIn(b'sz="100"', slide_xml)
            self._assert_valid_drawingml_font_sizes(slide_xml)
            self._assert_drawingml_paragraph_property_order(slide_xml)

    # Verifies FR-2026-08-04-11 and FR-2026-08-05-01.
    def test_pptx_replaces_reachable_smartart_data_and_editable_wordart(self) -> None:
        for layout_mode in (
            "preserve-source-formatting",
            "preserve-basic-layout",
            "preserve-basic-layout-source-font",
        ):
            with self.subTest(layout_mode=layout_mode), TemporaryDirectory() as temporary_directory:
                root = Path(temporary_directory)
                input_root = root / "input"
                output_root = root / "output"
                input_root.mkdir()
                source = input_root / "deck.pptx"
                presentation = Presentation()
                slide = presentation.slides.add_slide(presentation.slide_layouts[6])
                wordart = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
                wordart.name = "editable-wordart"
                wordart.text = "WordArt source"
                wordart.text_frame.paragraphs[0].runs[0].font.name = "Source WordArt Font"
                wordart.text_frame.paragraphs[0].runs[0].font.size = Pt(24)
                wordart.text_frame._element.bodyPr.set("anchor", "ctr")
                wordart.text_frame._element.bodyPr.append(
                    parse_xml(
                        '<a:prstTxWarp xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
                        'prst="textNoShape"><a:avLst/></a:prstTxWarp>'
                    )
                )
                wordart.text_frame.paragraphs[0].runs[0]._r.get_or_add_rPr().append(
                    parse_xml(
                        '<a:pattFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
                        'prst="pct10"><a:fgClr><a:srgbClr val="112233"/></a:fgClr>'
                        "<a:bgClr><a:srgbClr val=\"FFFFFF\"/></a:bgClr></a:pattFill>"
                    )
                )
                wordart._element.spPr.append(
                    parse_xml(
                        '<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                        '<a:outerShdw blurRad="40000" dist="20000" dir="5400000"/>'
                        "</a:effectLst>"
                    )
                )
                coloured_text = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
                coloured_text.name = "coloured-text"
                coloured_text.text = "Coloured source"
                coloured_run = coloured_text.text_frame.paragraphs[0].runs[0]
                coloured_run.font.name = "Source Coloured Font"
                coloured_run.font.size = Pt(24)
                coloured_run.hyperlink.address = "https://example.invalid"
                run_properties = coloured_run._r.get_or_add_rPr()
                for markup in (
                    '<a:solidFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                    '<a:srgbClr val="112233"/></a:solidFill>',
                    '<a:ln xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" w="12700">'
                    '<a:solidFill><a:srgbClr val="445566"/></a:solidFill></a:ln>',
                    '<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                    '<a:outerShdw blurRad="40000" dist="20000" dir="5400000"/></a:effectLst>',
                    '<a:highlight xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                    '<a:srgbClr val="778899"/></a:highlight>',
                    '<a:uLn xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                    '<a:solidFill><a:srgbClr val="AABBCC"/></a:solidFill></a:uLn>',
                    '<a:uFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                    '<a:solidFill><a:srgbClr val="DDEEFF"/></a:solidFill></a:uFill>',
                ):
                    run_properties.append(parse_xml(markup))
                run_properties.set("lang", "en-GB")
                end_properties = coloured_text.text_frame.paragraphs[0]._p.get_or_add_endParaRPr()
                end_properties.set("lang", "en-GB")
                end_properties.append(
                    parse_xml(
                        '<a:highlight xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                        '<a:srgbClr val="AABBCC"/></a:highlight>'
                    )
                )
                presentation.save(str(source))
                self._add_reachable_smartart_data_part(source)

                provider = _RecordingReplacementProvider(replacement_text="Replaced")
                result = self._run(
                    input_root,
                    output_root,
                    _EmptyOcrProvider(),
                    provider,
                    document_text_layout=layout_mode,
                )

                self.assertEqual(1, result.processed_files)
                self.assertGreaterEqual(result.replaced_native_text_items, 4)
                requested_texts = [request.text for request in provider.requests]
                self.assertEqual(1, requested_texts.count("SmartArt first"))
                self.assertEqual(1, requested_texts.count("SmartArt second"))
                self.assertEqual(1, requested_texts.count("WordArt source"))
                self.assertEqual(1, requested_texts.count("Coloured source"))
                output = output_root / "deck.pptx"
                loaded = Presentation(str(output))
                output_wordart = next(
                    shape
                    for shape in loaded.slides[0].shapes
                    if shape.name == "editable-wordart"
                )
                output_coloured_text = next(
                    shape
                    for shape in loaded.slides[0].shapes
                    if shape.name == "coloured-text"
                )
                self.assertEqual("Replaced", output_wordart.text)
                self.assertEqual("Replaced", output_coloured_text.text)
                expected_wordart_font = (
                    "Source WordArt Font"
                    if layout_mode == "preserve-source-formatting"
                    else "Noto Sans JP"
                )
                expected_coloured_font = (
                    "Source Coloured Font"
                    if layout_mode == "preserve-source-formatting"
                    else "Noto Sans JP"
                )
                self.assertEqual(
                    expected_wordart_font,
                    output_wordart.text_frame.paragraphs[0].runs[0].font.name,
                )
                self.assertEqual(
                    expected_coloured_font,
                    output_coloured_text.text_frame.paragraphs[0].runs[0].font.name,
                )
                if layout_mode != "preserve-source-formatting":
                    self.assertEqual(MSO_AUTO_SIZE.NONE, output_wordart.text_frame.auto_size)
                    self.assertEqual(MSO_AUTO_SIZE.NONE, output_coloured_text.text_frame.auto_size)
                output_end_properties = (
                    output_coloured_text.text_frame.paragraphs[0]._p.endParaRPr
                )
                self.assertIsNotNone(output_end_properties.find(qn("a:highlight")))
                self.assertEqual("en-GB", output_end_properties.get("lang"))
                with ZipFile(output) as archive:
                    smartart_data = archive.read("ppt/diagrams/data1.xml")
                    slide_xml = archive.read("ppt/slides/slide1.xml")
                self.assertEqual(2, smartart_data.count(b"Replaced"))
                self.assertNotIn(b"SmartArt first", smartart_data)
                self.assertNotIn(b"SmartArt second", smartart_data)
                self.assertIn(b"outerShdw", slide_xml)
                self.assertIn(b"prstTxWarp", slide_xml)
                self.assertIn(b"pattFill", slide_xml)
                self.assertIn(b'anchor="ctr"', slide_xml)
                for expected_markup in (
                    b"solidFill",
                    b"ln w=\"12700\"",
                    b"effectLst",
                    b"highlight",
                    b"uLn",
                    b"uFill",
                    b'lang="en-GB"',
                    b"hlinkClick",
                ):
                    self.assertIn(expected_markup, slide_xml)

    # Verifies FR-2026-08-04-05.
    def test_pptx_no_autofit_uses_source_width_and_natural_height(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"
            output_root = root / "output"
            input_root.mkdir()
            source = input_root / "deck.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            no_autofit = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(0.2))
            no_autofit.name = "no-autofit"
            no_autofit.text = "Source text " * 40
            no_autofit.text_frame.paragraphs[0].runs[0].font.size = Pt(24)
            no_autofit.text_frame.auto_size = MSO_AUTO_SIZE.NONE
            no_autofit_dimensions = (int(no_autofit.width), int(no_autofit.height))
            text_to_fit_shape = slide.shapes.add_textbox(
                Inches(1), Inches(2), Inches(6), Inches(0.2)
            )
            text_to_fit_shape.name = "text-to-fit-shape"
            text_to_fit_shape.text = "Source text " * 40
            text_to_fit_shape.text_frame.paragraphs[0].runs[0].font.size = Pt(24)
            text_to_fit_shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            table = slide.shapes.add_table(1, 1, Inches(1), Inches(3), Inches(6), Inches(0.2)).table
            table.cell(0, 0).text = "Source text " * 40
            table.cell(0, 0).text_frame.paragraphs[0].runs[0].font.size = Pt(24)
            presentation.save(str(source))

            self._run(
                input_root,
                output_root,
                _EmptyOcrProvider(),
                _RecordingReplacementProvider(replacement_text="Replacement text " * 25),
                document_text_layout="preserve-basic-layout",
            )

            output_presentation = Presentation(str(output_root / "deck.pptx"))
            output_slide = output_presentation.slides[0]
            output_no_autofit = next(
                shape for shape in output_slide.shapes if shape.name == "no-autofit"
            )
            output_text_to_fit_shape = next(
                shape for shape in output_slide.shapes if shape.name == "text-to-fit-shape"
            )
            output_table = next(shape.table for shape in output_slide.shapes if shape.has_table)
            no_autofit_size = output_no_autofit.text_frame.paragraphs[0].runs[0].font.size.pt
            text_to_fit_shape_size = (
                output_text_to_fit_shape.text_frame.paragraphs[0].runs[0].font.size.pt
            )
            table_size = output_table.cell(0, 0).text_frame.paragraphs[0].runs[0].font.size.pt
            self.assertEqual(no_autofit_dimensions, (int(output_no_autofit.width), int(output_no_autofit.height)))
            self.assertGreater(no_autofit_size, text_to_fit_shape_size)
            self.assertGreater(no_autofit_size, table_size)

    # Verifies FR-2026-08-03-15.
    def test_pptx_basic_layout_retains_vertical_alignment_xml(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"
            output_root = root / "output"
            input_root.mkdir()
            source = input_root / "deck.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            middle = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
            middle.name = "middle-aligned"
            middle.text = "Middle aligned"
            middle.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            default = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(3), Inches(1))
            default.name = "default-aligned"
            default.text = "Default aligned"
            presentation.save(str(source))

            self._run(
                input_root,
                output_root,
                _EmptyOcrProvider(),
                _RecordingReplacementProvider(replacement_text="Replacement"),
                document_text_layout="preserve-basic-layout",
            )

            output_presentation = Presentation(str(output_root / "deck.pptx"))
            output_shapes = {shape.name: shape for shape in output_presentation.slides[0].shapes}
            self.assertEqual("ctr", output_shapes["middle-aligned"].text_frame._element.bodyPr.get("anchor"))
            self.assertIsNone(output_shapes["default-aligned"].text_frame._element.bodyPr.get("anchor"))

    # Verifies FR-2026-08-27-02.
    def test_pptx_source_font_mode_falls_back_when_source_faces_are_unavailable(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"
            output_root = root / "output"
            input_root.mkdir()
            source = input_root / "deck.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            named_source = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
            named_source.name = "named-source"
            named_source.text = "Named source font"
            named_run = named_source.text_frame.paragraphs[0].runs[0]
            named_run.font.name = "Source Presentation Font"
            named_run.font.size = Pt(24)
            unnamed_source = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(3), Inches(1))
            unnamed_source.name = "unnamed-source"
            unnamed_source.text = "Fallback font"
            unnamed_source.text_frame.paragraphs[0].runs[0].font.size = Pt(24)
            themed_source = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(3), Inches(1))
            themed_source.name = "themed-source"
            themed_source.text = "Theme source font"
            themed_run = themed_source.text_frame.paragraphs[0].runs[0]
            themed_run.font.name = "+mj-lt"
            themed_run.font.size = Pt(24)
            presentation.save(str(source))

            self._run(
                input_root,
                output_root,
                _EmptyOcrProvider(),
                _RecordingReplacementProvider(replacement_text="Replacement text " * 20),
                document_text_layout="preserve-basic-layout-source-font",
            )

            output_presentation = Presentation(str(output_root / "deck.pptx"))
            output_shapes = output_presentation.slides[0].shapes
            output_named = next(shape for shape in output_shapes if shape.name == "named-source")
            output_unnamed = next(shape for shape in output_shapes if shape.name == "unnamed-source")
            output_themed = next(shape for shape in output_shapes if shape.name == "themed-source")
            self.assertEqual(
                "Noto Sans JP",
                output_named.text_frame.paragraphs[0].runs[0].font.name,
            )
            self.assertEqual(
                "Noto Sans JP",
                output_unnamed.text_frame.paragraphs[0].runs[0].font.name,
            )
            self.assertEqual(
                "Noto Sans JP",
                output_themed.text_frame.paragraphs[0].runs[0].font.name,
            )
            self.assertEqual(MSO_AUTO_SIZE.NONE, output_named.text_frame.auto_size)

    # Verifies FR-2026-08-03-14.
    def test_xlsx_basic_layout_fits_only_explicitly_bounded_cells(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"
            output_root = root / "output"
            input_root.mkdir()
            source = input_root / "workbook.xlsx"
            workbook = Workbook()
            worksheet = workbook.active
            self.assertIsNotNone(worksheet)
            assert worksheet is not None
            worksheet.column_dimensions["A"].width = 12
            worksheet.column_dimensions["B"].width = 12
            worksheet.row_dimensions[1].height = 30
            worksheet.row_dimensions[2].height = 30
            worksheet["A1"] = "Bounded source"
            worksheet["A1"].font = Font(name="Source Sans", sz=24, bold=True)
            worksheet["A1"].alignment = Alignment(horizontal="center", wrap_text=False)
            worksheet["A1"].protection = Protection(locked=False)
            worksheet.merge_cells("A2:B2")
            worksheet["A2"] = "Merged source"
            worksheet["C1"] = "Unbounded source"
            workbook.save(source)

            result = self._run(
                input_root,
                output_root,
                _EmptyOcrProvider(),
                _RecordingReplacementProvider(replacement_text="Long replacement text " * 30),
                document_text_layout="preserve-basic-layout",
            )

            self.assertEqual(1, result.processed_files)
            self.assertEqual(3, result.replaced_native_text_items)
            output = load_workbook(output_root / "workbook.xlsx")
            output_sheet = output.active
            self.assertIsNotNone(output_sheet)
            assert output_sheet is not None
            self.assertIn("Long replacement text", output_sheet["A1"].value)
            # FR-2026-08-04-07: XLSX has no portable embedded-font path, so the
            # bounded cell keeps its source face while using the shared fitted size.
            self.assertEqual("Source Sans", output_sheet["A1"].font.name)
            self.assertLess(output_sheet["A1"].font.sz, 24)
            self.assertTrue(output_sheet["A1"].alignment.wrap_text)
            self.assertFalse(output_sheet["A1"].alignment.shrink_to_fit)
            with ZipFile(output_root / "workbook.xlsx") as archive:
                styles = ElementTree.fromstring(archive.read("xl/styles.xml"))
            namespace = "{http://schemas.openxmlformats.org/spreadsheetml/2006/main}"
            cell_xfs = styles.find(namespace + "cellXfs")
            assert cell_xfs is not None
            self.assertFalse(any(
                [child.tag for child in xf].index(namespace + "protection")
                < [child.tag for child in xf].index(namespace + "alignment")
                for xf in cell_xfs
                if namespace + "protection" in [child.tag for child in xf]
                and namespace + "alignment" in [child.tag for child in xf]
            ))
            self.assertIn("Long replacement text", output_sheet["A2"].value)
            self.assertEqual("Calibri", output_sheet["A2"].font.name)
            self.assertIn("Long replacement text", output_sheet["C1"].value)
            self.assertEqual("Calibri", output_sheet["C1"].font.name)

    # Verifies FR-2026-08-03-14.
    def test_xlsx_basic_layout_retains_structured_table_headers(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"
            output_root = root / "output"
            input_root.mkdir()
            source = input_root / "workbook.xlsx"
            workbook = Workbook()
            worksheet = workbook.active
            self.assertIsNotNone(worksheet)
            assert worksheet is not None
            worksheet.column_dimensions["A"].width = 16
            worksheet.column_dimensions["B"].width = 16
            worksheet.row_dimensions[1].height = 24
            worksheet.row_dimensions[2].height = 24
            worksheet.append(["Synthetic Header One", "Synthetic Header Two"])
            worksheet.append(["First body value", "Second body value"])
            worksheet.add_table(Table(displayName="SyntheticTable", ref="A1:B2"))
            workbook.save(source)

            result = self._run(
                input_root,
                output_root,
                _EmptyOcrProvider(),
                _RecordingReplacementProvider(replacement_text="Replacement body value"),
                document_text_layout="preserve-basic-layout",
            )

            self.assertEqual(2, result.replaced_native_text_items)
            output = load_workbook(output_root / "workbook.xlsx")
            output_sheet = output.active
            self.assertIsNotNone(output_sheet)
            assert output_sheet is not None
            self.assertEqual("Synthetic Header One", output_sheet["A1"].value)
            self.assertEqual("Synthetic Header Two", output_sheet["B1"].value)
            self.assertEqual("Replacement body value", output_sheet["A2"].value)
            self.assertEqual("Replacement body value", output_sheet["B2"].value)
            with ZipFile(source) as source_archive, ZipFile(output_root / "workbook.xlsx") as output_archive:
                self.assertEqual(
                    source_archive.read("xl/tables/table1.xml"),
                    output_archive.read("xl/tables/table1.xml"),
                )

    # Verifies FR-2026-08-04-07 and FR-2026-08-22-03.
    def test_docx_basic_layout_fits_drawing_text_and_embeds_conformant_fonts(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"
            output_root = root / "output"
            input_root.mkdir()
            source = input_root / "document.docx"
            self._write_complete_docx(source)
            self._run(
                input_root, output_root, _EmptyOcrProvider(),
                _RecordingReplacementProvider(replacement_text="Long replacement text " * 20),
                document_text_layout="preserve-basic-layout",
            )
            output = output_root / "document.docx"
            Document(str(output))
            with ZipFile(output) as archive:
                data = archive.read("word/document.xml")
                font_table = ElementTree.fromstring(archive.read("word/fontTable.xml"))
                font_relationships = ElementTree.fromstring(
                    archive.read("word/_rels/fontTable.xml.rels")
                )
                content_types = ElementTree.fromstring(archive.read("[Content_Types].xml"))
                settings = ElementTree.fromstring(archive.read("word/settings.xml"))
                self._assert_word_run_property_order(ElementTree.fromstring(data))
                targets = {item.get("Id"): item.get("Target") for item in font_relationships}
                embedded_font_ids: set[str] = set()
                namespace = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"
                relationships_namespace = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
                font_sources = {
                    "Noto Sans JP": "sans-serif",
                    "Noto Serif JP": "serif",
                    "Noto Sans Mono": "fixed-width",
                }
                for family, classification in font_sources.items():
                    entries = [
                        item for item in font_table.findall(f"{namespace}font")
                        if item.get(f"{namespace}name") == family
                    ]
                    self.assertEqual(1, len(entries))
                    if family != "Noto Sans JP":
                        expected_metadata = self._sfnt_font_metadata(
                            static_noto_bytes(classification, False)
                        )
                        metadata_names = [
                            child.tag.rsplit("}", 1)[-1]
                            for child in entries[0]
                        ]
                        self.assertEqual(
                            [
                                "panose1", "charset", "family", "pitch", "sig",
                                "embedRegular", "embedBold",
                            ],
                            metadata_names,
                        )
                        for child in entries[0][:5]:
                            self.assertEqual(
                                expected_metadata[child.tag.rsplit("}", 1)[-1]],
                                {
                                    key.rsplit("}", 1)[-1]: value
                                    for key, value in child.attrib.items()
                                },
                            )
                    for style in ("embedRegular", "embedBold"):
                        embedding = entries[0].find(f"{namespace}{style}")
                        self.assertIsNotNone(embedding)
                        assert embedding is not None
                        relationship_id = embedding.get(f"{relationships_namespace}id")
                        font_key = embedding.get(f"{namespace}fontKey")
                        self.assertIsNotNone(relationship_id)
                        self.assertIsNotNone(font_key)
                        assert relationship_id is not None and font_key is not None
                        self.assertNotIn(relationship_id, embedded_font_ids)
                        embedded_font_ids.add(relationship_id)
                        target = targets[relationship_id]
                        recovered = bytearray(archive.read(f"word/{target}"))
                        standard_key = UUID(font_key.strip("{}")).bytes[::-1]
                        for offset in range(min(32, len(recovered))):
                            recovered[offset] ^= standard_key[offset % 16]
                        self.assertEqual(
                            bytes(recovered),
                            static_noto_bytes(classification, style == "embedBold"),
                        )
                        self.assertIsNotNone(
                            skia.Typeface.MakeFromData(skia.Data.MakeWithCopy(bytes(recovered)))
                        )
                self.assertEqual(6, len(embedded_font_ids))
                self.assertEqual(
                    1,
                    sum(item.get("Extension") == "odttf" for item in content_types),
                )
                self.assertEqual(
                    1,
                    sum(item.get("PartName") == "/word/fontTable.xml" for item in content_types),
                )
                content_type_names = [item.tag.rsplit("}", 1)[-1] for item in content_types]
                first_override = content_type_names.index("Override")
                self.assertNotIn("Default", content_type_names[first_override:])
                setting_names = [item.tag.rsplit("}", 1)[-1] for item in settings]
                self.assertEqual(
                    [
                        "writeProtection", "zoom", "embedTrueTypeFonts",
                        "bordersDoNotSurroundHeader",
                    ],
                    setting_names,
                )
            self.assertIn(b"Noto Sans JP", data)
            self.assertIn(b"Long replacement text", data)
            generated_property_names = {
                child.tag.rsplit("}", 1)[-1]
                for properties in ElementTree.fromstring(data).findall(
                    ".//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}txbxContent"
                    "//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}rPr"
                )
                for child in properties
            }
            self.assertTrue({"b", "i", "u"}.issubset(generated_property_names))

            corrupted = root / "corrupted.docx"
            with ZipFile(output) as source_archive, ZipFile(corrupted, "w", ZIP_DEFLATED) as corrupted_archive:
                corrupted_part = "word/fonts/pipeline-sans-serif-regular.odttf"
                for entry in source_archive.infolist():
                    payload = source_archive.read(entry.filename)
                    if entry.filename == corrupted_part:
                        payload = b"bad font" + payload[8:]
                    corrupted_archive.writestr(entry, payload)
            with ZipFile(corrupted) as corrupted_archive:
                parts = {
                    entry.filename: corrupted_archive.read(entry.filename)
                    for entry in corrupted_archive.infolist()
                }
            with self.assertRaisesRegex(ValueError, "not a loadable OpenType font"):
                _validate_docx_embedded_fonts(parts)

            missing_setting_parts = parts.copy()
            settings = ElementTree.fromstring(missing_setting_parts["word/settings.xml"])
            setting = settings.find("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}embedTrueTypeFonts")
            self.assertIsNotNone(setting)
            assert setting is not None
            settings.remove(setting)
            missing_setting_parts["word/settings.xml"] = ElementTree.tostring(settings)
            with self.assertRaisesRegex(ValueError, "embedded-font setting is missing"):
                _validate_docx_embedded_fonts(missing_setting_parts)

            misplaced_setting_parts = parts.copy()
            settings = ElementTree.fromstring(misplaced_setting_parts["word/settings.xml"])
            setting = settings.find("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}embedTrueTypeFonts")
            self.assertIsNotNone(setting)
            assert setting is not None
            settings.remove(setting)
            settings.append(setting)
            misplaced_setting_parts["word/settings.xml"] = ElementTree.tostring(settings)
            with self.assertRaisesRegex(ValueError, "not in CT_Settings order"):
                _validate_docx_embedded_fonts(misplaced_setting_parts)

    # Verifies FR-2026-08-03-08.
    def test_preserves_compatibility_prefix_required_by_rewritten_office_xml(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"
            output_root = root / "output"
            input_root.mkdir()
            with ZipFile(input_root / "document.pptx", "w", ZIP_DEFLATED) as archive:
                archive.writestr(
                    "ppt/slides/slide1.xml",
                    """<?xml version=\"1.0\"?>
                    <p:sld xmlns:p=\"http://schemas.openxmlformats.org/presentationml/2006/main\"
                      xmlns:a=\"http://schemas.openxmlformats.org/drawingml/2006/main\"
                      xmlns:mc=\"http://schemas.openxmlformats.org/markup-compatibility/2006\"
                      xmlns:p14=\"http://schemas.microsoft.com/office/powerpoint/2010/main\">
                      <a:t>Mask</a:t>
                      <mc:AlternateContent><mc:Choice Requires=\"p14\"><p14:content/></mc:Choice></mc:AlternateContent>
                    </p:sld>""",
                )

            self._run(input_root, output_root, _EmptyOcrProvider(), _RecordingReplacementProvider())

            with ZipFile(output_root / "document.pptx") as archive:
                xml = archive.read("ppt/slides/slide1.xml")
            self.assertIn(b'Requires="p14"', xml)
            self.assertIn(
                b'xmlns:p14="http://schemas.microsoft.com/office/powerpoint/2010/main"',
                xml,
            )

    # Verifies FR-2026-08-03-03.
    def test_processes_every_embedded_office_bitmap(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"
            output_root = root / "output"
            input_root.mkdir()
            image_data = BytesIO()
            Image.new("RGB", (20, 20), "white").save(image_data, "PNG")
            with ZipFile(input_root / "document.docx", "w", ZIP_DEFLATED) as archive:
                archive.writestr("word/media/image1.png", image_data.getvalue())
            ocr_provider = _CountingOcrProvider()

            result = self._run(
                input_root, output_root, ocr_provider, _RecordingReplacementProvider()
            )

            self.assertEqual(1, result.processed_files)
            self.assertEqual(1, ocr_provider.calls)
            with ZipFile(output_root / "document.docx") as archive:
                self.assertTrue(archive.read("word/media/image1.png"))

    # Verifies FR-2026-08-03-04.
    def test_reports_file_name_and_document_work_progress(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"
            output_root = root / "output"
            input_root.mkdir()
            image_data = BytesIO()
            Image.new("RGB", (20, 20), "white").save(image_data, "PNG")
            with ZipFile(input_root / "document.docx", "w", ZIP_DEFLATED) as archive:
                archive.writestr("word/media/image1.png", image_data.getvalue())
            progress_bars: list[_RecordedProgress] = []

            def make_progress(total: int, label: str) -> ProgressReporter:
                progress = _RecordedProgress(total, label)
                progress_bars.append(progress)
                return progress

            standard_output = StringIO()
            with redirect_stdout(standard_output):
                result = self._run(
                    input_root,
                    output_root,
                    _CountingOcrProvider(),
                    _RecordingReplacementProvider(),
                    show_progress=True,
                    progress_factory=make_progress,
                )

            self.assertEqual(1, result.processed_files)
            self.assertEqual("Processing: document.docx\n", standard_output.getvalue())
            self.assertEqual(1, len(progress_bars))
            self.assertEqual(2, progress_bars[0].total)
            self.assertEqual("document.docx", progress_bars[0].label)
            self.assertEqual(["embedded image 1", "native text"], progress_bars[0].postfixes)
            self.assertEqual(2, progress_bars[0].updates)
            self.assertTrue(progress_bars[0].closed)

    # Verifies FR-2026-08-27-11.
    def test_reports_pdf_native_text_progress_for_each_page_and_form_pass(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"
            output_root = root / "output"
            input_root.mkdir()
            source = input_root / "document.pdf"
            writer = PdfWriter()
            writer.add_blank_page(100, 100)
            writer.add_blank_page(100, 100)
            with source.open("wb") as output_file:
                writer.write(output_file)
            progress_bars: list[_RecordedProgress] = []

            def make_progress(total: int, label: str) -> ProgressReporter:
                progress = _RecordedProgress(total, label)
                progress_bars.append(progress)
                return progress

            result = self._run(
                input_root,
                output_root,
                _EmptyOcrProvider(),
                _RecordingReplacementProvider(),
                show_progress=True,
                progress_factory=make_progress,
            )

            self.assertEqual(1, result.processed_files)
            self.assertEqual(1, len(progress_bars))
            self.assertEqual(3, progress_bars[0].total)
            self.assertEqual(
                ["native text page 1/2", "native text page 2/2", "native form fields"],
                progress_bars[0].postfixes,
            )
            self.assertEqual(3, progress_bars[0].updates)

    # Verifies FR-2026-08-03-05.
    def test_replaces_editable_svg_vector_text_without_ocr(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"
            output_root = root / "output"
            input_root.mkdir()
            with ZipFile(input_root / "document.docx", "w", ZIP_DEFLATED) as archive:
                archive.writestr(
                    "word/media/vector.svg",
                    b'<svg xmlns="http://www.w3.org/2000/svg"><text>Mask</text></svg>',
                )
            ocr_provider = _CountingOcrProvider()

            result = self._run(
                input_root, output_root, ocr_provider, _RecordingReplacementProvider()
            )

            self.assertEqual(1, result.processed_files)
            self.assertEqual(1, result.replaced_native_text_items)
            self.assertEqual(0, result.replaced_image_regions)
            self.assertEqual(0, result.retained_vector_graphics)
            self.assertEqual(0, ocr_provider.calls)
            with ZipFile(output_root / "document.docx") as archive:
                vector_data = archive.read("word/media/vector.svg")
            self.assertIn(b"####", vector_data)
            self.assertNotIn(b"Mask", vector_data)

    # Verifies FR-2026-08-03-05.
    def test_retains_noneditable_vector_and_reports_it(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"
            output_root = root / "output"
            input_root.mkdir()
            vector_data = b'<svg xmlns="http://www.w3.org/2000/svg"><path d="M0 0"/></svg>'
            with ZipFile(input_root / "document.docx", "w", ZIP_DEFLATED) as archive:
                archive.writestr("word/media/vector.svg", vector_data)

            standard_error = StringIO()
            with redirect_stderr(standard_error):
                result = self._run(
                    input_root, output_root, _EmptyOcrProvider(), _RecordingReplacementProvider()
                )

            self.assertEqual(1, result.processed_files)
            self.assertEqual(1, result.retained_vector_graphics)
            self.assertIn("Retained vector", standard_error.getvalue())
            with ZipFile(output_root / "document.docx") as archive:
                self.assertEqual(vector_data, archive.read("word/media/vector.svg"))

    # Verifies FR-2026-08-03-03.
    def test_skips_low_confidence_ocr_text_and_continues_after_a_file_failure(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"
            output_root = root / "output"
            input_root.mkdir()
            self._write_png(input_root / "good.png")
            (input_root / "broken.png").write_bytes(b"not an image")
            provider = _RecordingReplacementProvider()

            failure_output = StringIO()
            with redirect_stderr(failure_output):
                result = self._run(
                    input_root,
                    output_root,
                    _LowConfidenceOcrProvider(),
                    provider,
                    diagnostics_enabled=True,
                )

            self.assertEqual(1, result.processed_files)
            self.assertEqual(1, result.failed_files)
            self.assertTrue((output_root / "good.png").is_file())
            diagnostic = output_root / "broken.png.diagnostics.json"
            self.assertTrue(diagnostic.is_file())
            self.assertEqual(
                "file_processing_failed",
                json.loads(diagnostic.read_text(encoding="utf-8"))["entries"][0]["reason_code"],
            )
            self.assertFalse(any(not request.is_filename for request in provider.requests))
            self.assertIn("broken.png", failure_output.getvalue())

    # Verifies FR-2026-08-28-03.
    def test_word_failure_records_safe_native_text_context_and_continues(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"
            output_root = root / "output"
            input_root.mkdir()
            with ZipFile(input_root / "document.docx", "w", ZIP_DEFLATED) as archive:
                archive.writestr(
                    "word/document.xml",
                    """<?xml version="1.0"?>
                    <w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
                      <w:body><w:p><w:r><w:t>synthetic request text</w:t></w:r></w:p></w:body>
                    </w:document>""",
                )
            self._write_png(input_root / "good.png")

            result = self._run(
                input_root,
                output_root,
                _EmptyOcrProvider(),
                _FailingReplacementProvider(),
                diagnostics_enabled=True,
            )

            self.assertEqual(1, result.failed_files)
            self.assertEqual(1, result.processed_files)
            self.assertFalse((output_root / "document.docx").exists())
            self.assertTrue((output_root / "good.png").exists())
            diagnostic = json.loads(
                (output_root / "document.docx.diagnostics.json").read_text(encoding="utf-8")
            )
            entry = diagnostic["entries"][0]
            self.assertEqual("RuntimeError", entry["exception_type"])
            self.assertEqual(["ValueError"], entry["exception_cause_types"])
            self.assertEqual(
                {
                    "stage": "native_xml",
                    "container_kind": "office_xml",
                    "operation": "text_replacement",
                    "location": {"package_part": "word/document.xml"},
                    "request": {
                        "kind": "text_replacement",
                        "source_language": "en",
                        "target_language": "en",
                        "is_filename": False,
                        "input_character_count": 22,
                    },
                },
                entry["failure_context"],
            )
            serialized = json.dumps(diagnostic)
            self.assertNotIn("synthetic request text", serialized)
            self.assertNotIn("Synthetic chained detail", serialized)

    # Verifies FR-2026-08-28-03.
    def test_word_embedded_image_failure_records_safe_ocr_context(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"
            output_root = root / "output"
            input_root.mkdir()
            image_bytes = BytesIO()
            Image.new("RGB", (7, 5), (16, 32, 48)).save(image_bytes, format="PNG")
            with ZipFile(input_root / "document.docx", "w", ZIP_DEFLATED) as archive:
                archive.writestr(
                    "word/document.xml",
                    """<?xml version="1.0"?>
                    <w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"/>""",
                )
                archive.writestr("word/media/image1.png", image_bytes.getvalue())

            result = self._run(
                input_root,
                output_root,
                _FailingOcrProvider(),
                _RecordingReplacementProvider(),
                diagnostics_enabled=True,
            )

            self.assertEqual(1, result.failed_files)
            entry = json.loads(
                (output_root / "document.docx.diagnostics.json").read_text(encoding="utf-8")
            )["entries"][0]
            self.assertEqual("RuntimeError", entry["exception_type"])
            self.assertEqual(["ValueError"], entry["exception_cause_types"])
            self.assertEqual(
                {
                    "stage": "embedded_bitmap",
                    "container_kind": "embedded_bitmap",
                    "operation": "ocr",
                    "location": {
                        "package_part": "word/media/image1.png",
                        "item_index": 1,
                    },
                    "request": {
                        "kind": "ocr",
                        "language": "en",
                        "image_width": 7,
                        "image_height": 5,
                        "image_mode": "RGB",
                    },
                },
                entry["failure_context"],
            )
            self.assertNotIn("Synthetic OCR chained detail", json.dumps(entry))

    # Verifies FR-2026-08-03-03.
    def test_replaces_native_pdf_text_content(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"
            output_root = root / "output"
            input_root.mkdir()
            source = input_root / "document.pdf"
            writer = PdfWriter()
            page = writer.add_blank_page(100, 100)
            contents = DecodedStreamObject()
            contents.set_data(b"BT /F1 12 Tf 10 10 Td (Hello) Tj ET")
            page.replace_contents(ContentStream(contents, writer))
            with source.open("wb") as source_file:
                writer.write(source_file)

            result = self._run(
                input_root, output_root, _EmptyOcrProvider(), _RecordingReplacementProvider()
            )

            self.assertEqual(1, result.processed_files)
            self.assertEqual(1, result.replaced_native_text_items)
            self.assertIn(b"\\043\\043\\043\\043\\043", (output_root / "document.pdf").read_bytes())

    # Verifies FR-2026-08-23-01.
    def test_basic_layout_pdf_content_uses_the_safe_replacement_font(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"; output_root = root / "output"; input_root.mkdir()
            source = input_root / "document.pdf"
            writer = PdfWriter(); page = writer.add_blank_page(100, 100)
            contents = DecodedStreamObject(); contents.set_data(b"BT /F1 12 Tf 10 10 Td [(Hel) 0 (lo)] TJ ET")
            page.replace_contents(ContentStream(contents, writer))
            with source.open("wb") as source_file: writer.write(source_file)

            provider = _RecordingReplacementProvider()
            self._run(
                input_root, output_root, _EmptyOcrProvider(), provider,
                document_text_layout="preserve-basic-layout",
            )
            self.assertEqual(["Hello"], [
                request.text for request in provider.requests if not request.is_filename
            ])

            output = PdfReader(output_root / "document.pdf")
            stream = ContentStream(output.pages[0].get_contents(), output)
            self.assertTrue(any(operator == b"Tf" and operands[0] == "/PipelineNoto" for operands, operator in stream.operations))
            self.assertTrue(any(
                operator == b"Tr" and operands and int(operands[0]) == 3
                for operands, operator in stream.operations
            ))
            self.assertNotIn("Hello", output.pages[0].extract_text())
            fonts = cast(DictionaryObject, cast(DictionaryObject, output.pages[0]["/Resources"])["/Font"])
            type_zero = cast(DictionaryObject, fonts["/PipelineNoto"].get_object())
            descendant = cast(DictionaryObject, cast(ArrayObject, type_zero["/DescendantFonts"])[0].get_object())
            self.assertTrue(descendant.get("/W"))

    # Verifies FR-2026-08-02-06 and FR-2026-08-24-02.
    def test_basic_layout_pdf_character_mask_preserves_distinct_whitespace_semantics(self) -> None:
        """Assign separate CIDs when the embedded font reuses a whitespace glyph."""
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"; output_root = root / "output"; input_root.mkdir()
            source = input_root / "document.pdf"
            writer = PdfWriter(); page = writer.add_blank_page(100, 100)
            contents = DecodedStreamObject()
            contents.set_data(
                b"BT /F1 12 Tf 10 10 Td <FEFF0041002000A00042> Tj ET"
            )
            page.replace_contents(ContentStream(contents, writer))
            with source.open("wb") as source_file:
                writer.write(source_file)

            result = self._run(
                input_root,
                output_root,
                _EmptyOcrProvider(),
                CharacterMaskProvider(),
                document_text_layout="preserve-basic-layout",
                diagnostics_enabled=True,
            )

            self.assertEqual(1, result.replaced_native_text_items)
            self.assertEqual([], result.diagnostic_sidecars)
            output = PdfReader(output_root / "document.pdf")
            self.assertIn("# \u00A0#", output.pages[0].extract_text())
            fonts = cast(DictionaryObject, cast(DictionaryObject, output.pages[0]["/Resources"])["/Font"])
            type_zero = cast(DictionaryObject, fonts["/PipelineNoto"].get_object())
            to_unicode = cast(DecodedStreamObject, type_zero["/ToUnicode"].get_object())
            mapping = to_unicode.get_data()
            self.assertIn(b"<0002> <0020>", mapping)
            self.assertIn(b"<0003> <00A0>", mapping)

    # Verifies FR-2026-08-27-06.
    def test_basic_layout_pdf_reports_font_serialization_not_region_reconstruction(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"; output_root = root / "output"; input_root.mkdir()
            source = input_root / "document.pdf"
            writer = PdfWriter(); page = writer.add_blank_page(100, 100)
            contents = DecodedStreamObject()
            contents.set_data(b"BT /F1 12 Tf 10 10 Td (source) Tj ET")
            page.replace_contents(ContentStream(contents, writer))
            with source.open("wb") as source_file:
                writer.write(source_file)

            with patch(
                "pipeline.folder_replacement.pdf._pdf_static_glyph_bytes",
                side_effect=_PdfReplacementSerializationError(
                    "pdf_replacement_font_encoding_invalid",
                    "Synthetic portable-font serialization failure.",
                    "######",
                ),
            ):
                result = self._run(
                    input_root,
                    output_root,
                    _EmptyOcrProvider(),
                    _RecordingReplacementProvider(),
                    document_text_layout="preserve-basic-layout",
                    diagnostics_enabled=True,
                )

            self.assertEqual(0, result.replaced_native_text_items)
            self.assertEqual(1, len(result.diagnostic_sidecars))
            report = json.loads(result.diagnostic_sidecars[0].read_text(encoding="utf-8"))
            self.assertEqual(
                [{
                    "kind": "unsupported",
                    "reason_code": "pdf_replacement_font_encoding_invalid",
                    "container_kind": "pdf_visual_text",
                    "page": 1,
                    "replacement_text": "######",
                }],
                [{key: entry[key] for key in (
                    "kind", "reason_code", "container_kind", "page", "replacement_text"
                )} for entry in report["entries"]],
            )

    # Verifies FR-2026-08-27-06.
    def test_basic_layout_pdf_records_unsupported_portable_text_and_keeps_the_region(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"; output_root = root / "output"; input_root.mkdir()
            source = input_root / "document.pdf"
            writer = PdfWriter(); page = writer.add_blank_page(100, 100)
            contents = DecodedStreamObject()
            contents.set_data(b"BT /F1 12 Tf 10 10 Td (Hello) Tj ET")
            page.replace_contents(ContentStream(contents, writer))
            with source.open("wb") as source_file:
                writer.write(source_file)

            with patch(
                "pipeline.folder_replacement.pdf._pdf_fitted_region_operations",
                side_effect=PortableTextUnsupportedError(
                    "portable_font_coverage_unsupported",
                    "\U0001f9ea",
                    ("Noto Sans", "Noto Sans Symbols 2"),
                    replacement_text="\U0001f9ea",
                ),
            ):
                result = self._run(
                    input_root,
                    output_root,
                    _EmptyOcrProvider(),
                    _RecordingReplacementProvider(replacement_text="\U0001f9ea"),
                    document_text_layout="preserve-basic-layout",
                    diagnostics_enabled=True,
                )

            self.assertEqual(1, result.processed_files)
            self.assertEqual(0, result.failed_files)
            self.assertIn(b"(Hello)", (output_root / "document.pdf").read_bytes())
            sidecar = json.loads(
                (output_root / "document.pdf.diagnostics.json").read_text(encoding="utf-8")
            )
            self.assertEqual("pdf_visual_text", sidecar["entries"][0]["container_kind"])
            self.assertEqual(
                ["Noto Sans", "Noto Sans Symbols 2"],
                sidecar["entries"][0]["candidate_faces"],
            )
            self.assertEqual(["U+1F9EA"], sidecar["entries"][0]["code_points"])
            self.assertEqual("Hello", sidecar["entries"][0]["source_text"])
            self.assertEqual("\U0001f9ea", sidecar["entries"][0]["replacement_text"])
            self.assertEqual(
                "pdf_page_user_space",
                sidecar["entries"][0]["region_location"]["coordinate_space"],
            )
            self.assertEqual(
                ["source_text", "replacement_text"],
                list(sidecar["entries"][0])[:2],
            )

    # Verifies FR-2026-08-23-01.
    def test_basic_layout_pdf_splits_widely_spaced_tj_fragments_into_visual_regions(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"; output_root = root / "output"; input_root.mkdir()
            source = input_root / "document.pdf"
            writer = PdfWriter(); page = writer.add_blank_page(240, 100)
            contents = DecodedStreamObject()
            contents.set_data(
                b"BT /F1 12 Tf 10 60 Td [(one) -5000 (two) -5000 (three) -5000 (four)] TJ ET"
            )
            page.replace_contents(ContentStream(contents, writer))
            with source.open("wb") as source_file:
                writer.write(source_file)

            provider = _RecordingReplacementProvider()
            result = self._run(
                input_root, output_root, _EmptyOcrProvider(), provider,
                document_text_layout="preserve-basic-layout",
            )

            self.assertEqual(4, result.replaced_native_text_items)
            self.assertEqual(["one", "two", "three", "four"], [
                request.text for request in provider.requests if not request.is_filename
            ])
            output = PdfReader(output_root / "document.pdf")
            stream = ContentStream(output.pages[0].get_contents(), output)
            font_name = ""
            text_matrix: list[float] | None = None
            generated_x_positions: list[float] = []
            for operands, operator in stream.operations:
                if operator == b"Tf":
                    font_name = str(operands[0])
                elif operator == b"Tm":
                    text_matrix = [float(value) for value in operands]
                elif operator == b"Tj" and font_name == "/PipelineNoto" and text_matrix is not None:
                    generated_x_positions.append(text_matrix[4])
            self.assertEqual(4, len(generated_x_positions))
            self.assertEqual(generated_x_positions, sorted(generated_x_positions))
            self.assertGreater(min(
                right - left for left, right in zip(generated_x_positions, generated_x_positions[1:])
            ), 40.0)

    # Verifies FR-2026-08-23-01.
    def test_basic_layout_pdf_splits_close_tj_fragments_at_a_table_cell_gap(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"; output_root = root / "output"; input_root.mkdir()
            source = input_root / "document.pdf"
            writer = PdfWriter(); page = writer.add_blank_page(120, 100)
            contents = DecodedStreamObject()
            contents.set_data(b"BT /F1 12 Tf 10 60 Td [(value) -600 (unit)] TJ ET")
            page.replace_contents(ContentStream(contents, writer))
            with source.open("wb") as source_file:
                writer.write(source_file)

            provider = _RecordingReplacementProvider()
            result = self._run(
                input_root, output_root, _EmptyOcrProvider(), provider,
                document_text_layout="preserve-basic-layout",
            )

            self.assertEqual(2, result.replaced_native_text_items)
            self.assertEqual(["value", "unit"], [
                request.text for request in provider.requests if not request.is_filename
            ])

    # Verifies FR-2026-08-24-02.
    def test_basic_layout_pdf_identity_replaces_through_the_fitted_output_path(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"; output_root = root / "output"; input_root.mkdir()
            source = input_root / "document.pdf"
            writer = PdfWriter(); page = writer.add_blank_page(100, 100)
            contents = DecodedStreamObject()
            contents.set_data(b"BT /F1 12 Tf 10 10 Td (Hello) Tj ET")
            page.replace_contents(ContentStream(contents, writer))
            with source.open("wb") as output_file: writer.write(output_file)

            provider = _RecordingReplacementProvider(replacement_text="Hello")
            result = self._run(
                input_root, output_root, _EmptyOcrProvider(), provider,
                document_text_layout="preserve-basic-layout",
            )

            self.assertEqual(1, result.replaced_native_text_items)
            self.assertEqual(["Hello"], [request.text for request in provider.requests if not request.is_filename])
            output = PdfReader(output_root / "document.pdf")
            stream = ContentStream(output.pages[0].get_contents(), output)
            self.assertTrue(any(
                operator == b"Tf" and operands[0] == "/PipelineNoto"
                for operands, operator in stream.operations
            ))
            self.assertFalse(any(
                operator == b"Tj" and operands and str(operands[0]) == "Hello"
                for operands, operator in stream.operations
            ))
            self.assertEqual("Hello", output.pages[0].extract_text().strip())
            pipeline_font_index = next(
                index for index, (operands, operator) in enumerate(stream.operations)
                if operator == b"Tf" and operands[0] == "/PipelineNoto"
            )
            self.assertTrue(any(
                operator == b"Tj" and operands and str(operands[0])
                for operands, operator in stream.operations[pipeline_font_index:]
            ))

    # Verifies FR-2026-08-24-02.
    def test_basic_layout_pdf_replaces_the_extractable_source_text(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"; output_root = root / "output"; input_root.mkdir()
            source = input_root / "document.pdf"
            writer = PdfWriter(); page = writer.add_blank_page(150, 100)
            contents = DecodedStreamObject()
            contents.set_data(b"BT /F1 12 Tf 10 10 Td (source) Tj ET")
            page.replace_contents(ContentStream(contents, writer))
            with source.open("wb") as output_file:
                writer.write(output_file)

            self._run(
                input_root, output_root, _EmptyOcrProvider(),
                _RecordingReplacementProvider(replacement_text="replacement"),
                document_text_layout="preserve-basic-layout",
            )

            output = PdfReader(output_root / "document.pdf")
            extracted = output.pages[0].extract_text()
            self.assertIn("replacement", extracted)
            self.assertNotIn("source", extracted)
            fonts = cast(DictionaryObject, cast(DictionaryObject, output.pages[0]["/Resources"])["/Font"])
            type_zero = cast(DictionaryObject, fonts["/PipelineNoto"].get_object())
            self.assertIsNotNone(type_zero.get("/ToUnicode"))

    # Verifies FR-2026-08-24-02 and FR-2026-08-28-04.
    def test_basic_layout_pdf_replaces_a_text_only_actual_text_scope(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"; output_root = root / "output"; input_root.mkdir()
            source = input_root / "document.pdf"
            writer = PdfWriter(); page = writer.add_blank_page(150, 100)
            contents = DecodedStreamObject()
            contents.set_data(
                b"BT /F1 12 Tf 10 10 Td /Span << /ActualText (alternate) >> BDC "
                b"(protected) Tj EMC 0 -20 Td (replaceable) Tj ET"
            )
            page.replace_contents(ContentStream(contents, writer))
            with source.open("wb") as output_file:
                writer.write(output_file)

            provider = _RecordingReplacementProvider(replacement_text="replacement")
            result = self._run(
                input_root, output_root, _EmptyOcrProvider(), provider,
                document_text_layout="preserve-basic-layout",
            )

            self.assertEqual(2, result.replaced_native_text_items)
            self.assertEqual(["protected", "replaceable"], [
                request.text for request in provider.requests if not request.is_filename
            ])
            output = PdfReader(output_root / "document.pdf")
            self.assertNotIn("protected", output.pages[0].extract_text())
            self.assertNotIn("alternate", output.pages[0].extract_text())
            self.assertIn("replacement", output.pages[0].extract_text())
            stream = ContentStream(output.pages[0].get_contents(), output)
            self.assertTrue(any(
                operator == b"BDC" and len(operands) >= 2
                and isinstance(operands[1], DictionaryObject)
                and "/ActualText" not in operands[1]
                for operands, operator in stream.operations
            ))

    # Verifies FR-2026-08-28-04.
    def test_basic_layout_pdf_rewrites_only_the_replaced_actual_text_invocation(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"; output_root = root / "output"; input_root.mkdir()
            source = input_root / "document.pdf"
            writer = PdfWriter(); page = writer.add_blank_page(150, 100)
            page[NameObject("/Resources")] = DictionaryObject({
                NameObject("/Properties"): DictionaryObject({
                    NameObject("/Shared"): DictionaryObject({
                        NameObject("/ActualText"): TextStringObject("alternate"),
                    }),
                }),
            })
            contents = DecodedStreamObject()
            contents.set_data(
                b"BT /F1 12 Tf 10 10 Td /Span /Shared BDC (replaceable) Tj EMC "
                b"0 -20 Td /Span /Shared BDC /Artifact BMC (retained) Tj EMC EMC ET"
            )
            page.replace_contents(ContentStream(contents, writer))
            with source.open("wb") as output_file:
                writer.write(output_file)

            result = self._run(
                input_root, output_root, _EmptyOcrProvider(),
                _RecordingReplacementProvider(replacement_text="replacement"),
                document_text_layout="preserve-basic-layout",
                diagnostics_enabled=True,
            )

            self.assertEqual(1, result.replaced_native_text_items)
            output = PdfReader(output_root / "document.pdf")
            stream = ContentStream(output.pages[0].get_contents(), output)
            shared_scopes = [
                operands for operands, operator in stream.operations
                if operator == b"BDC" and len(operands) >= 2
            ]
            self.assertIsInstance(shared_scopes[0][1], DictionaryObject)
            self.assertNotIn("/ActualText", shared_scopes[0][1])
            self.assertEqual("/Shared", str(shared_scopes[1][1]))
            resources = cast(DictionaryObject, output.pages[0].get("/Resources"))
            properties = cast(DictionaryObject, resources.get("/Properties"))
            shared_property = cast(DictionaryObject, properties.get("/Shared"))
            self.assertEqual("alternate", str(shared_property.get("/ActualText")))
            diagnostic = json.loads((output_root / "document.pdf.diagnostics.json").read_text())
            self.assertTrue(any(
                entry.get("reason_code") == "pdf_text_marked_content_actual_text"
                for entry in diagnostic["entries"]
            ))

    # Verifies FR-2026-08-23-02 and FR-2026-08-24-02.
    def test_basic_layout_pdf_keeps_each_fitted_region_in_its_source_paint_state(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"; output_root = root / "output"; input_root.mkdir()
            source = input_root / "document.pdf"
            writer = PdfWriter(); page = writer.add_blank_page(200, 100)
            page[NameObject("/Resources")] = DictionaryObject({
                NameObject("/ExtGState"): DictionaryObject({
                    NameObject("/Full"): DictionaryObject({NameObject("/ca"): NumberObject(1)}),
                    NameObject("/Half"): DictionaryObject({NameObject("/ca"): FloatObject(0.5)}),
                }),
            })
            contents = DecodedStreamObject()
            contents.set_data(
                b"BT /F1 12 Tf 10 70 Td 1 0 0 rg /Full gs (first) Tj "
                b"0 0 1 rg /Half gs 0 -20 Td (second) Tj "
                b"1 0 0 rg /Full gs 0 -20 Td (third) Tj 14 TL T* 1 Tr (outline) Tj ET"
            )
            page.replace_contents(ContentStream(contents, writer))
            with source.open("wb") as output_file: writer.write(output_file)

            provider = _RecordingReplacementProvider(replacement_text="mask")
            result = self._run(
                input_root, output_root, _EmptyOcrProvider(), provider,
                document_text_layout="preserve-basic-layout",
            )

            self.assertEqual(3, result.replaced_native_text_items)
            self.assertEqual(["first", "second", "third"], [
                request.text for request in provider.requests if not request.is_filename
            ])
            output = PdfReader(output_root / "document.pdf")
            stream = ContentStream(output.pages[0].get_contents(), output)
            self.assertEqual(1, sum(operator == b"BT" for _operands, operator in stream.operations))
            self.assertEqual(1, sum(operator == b"ET" for _operands, operator in stream.operations))
            self.assertFalse(any(
                operator == b"Tj" and operands and str(operands[0]) in {"first", "second", "third"}
                for operands, operator in stream.operations
            ))
            outline_index = next(
                index for index, (operands, operator) in enumerate(stream.operations)
                if operator == b"Tj" and operands and str(operands[0]) == "outline"
            )
            self.assertEqual(b"Tr", stream.operations[outline_index - 1][1])
            self.assertEqual(1, int(stream.operations[outline_index - 1][0][0]))

            fill_colour: tuple[float, float, float] | None = None
            opacity: str | None = None
            fitted_paint_states: list[tuple[tuple[float, float, float] | None, str | None]] = []
            for operands, operator in stream.operations:
                if operator == b"rg":
                    fill_colour = (float(operands[0]), float(operands[1]), float(operands[2]))
                elif operator == b"gs":
                    opacity = str(operands[0])
                elif operator == b"Tf" and operands[0] == "/PipelineNoto":
                    fitted_paint_states.append((fill_colour, opacity))
            self.assertEqual([
                ((1.0, 0.0, 0.0), "/Full"),
                ((0.0, 0.0, 1.0), "/Half"),
                ((1.0, 0.0, 0.0), "/Full"),
            ], fitted_paint_states)

    # Verifies FR-2026-08-23-03.
    def test_basic_layout_pdf_replaces_fill_and_stroke_with_fill_only_output(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"; output_root = root / "output"; input_root.mkdir()
            source = input_root / "document.pdf"
            writer = PdfWriter(); page = writer.add_blank_page(200, 100)
            contents = DecodedStreamObject()
            contents.set_data(
                b"2 w 1 0 0 rg 0 0 1 RG BT /F1 12 Tf 14 TL 10 70 Td "
                b"2 Tr (filled and stroked) Tj 0 Tr T* (fill only) Tj ET"
            )
            page.replace_contents(ContentStream(contents, writer))
            with source.open("wb") as output_file: writer.write(output_file)

            provider = _RecordingReplacementProvider(replacement_text="mask")
            result = self._run(
                input_root, output_root, _EmptyOcrProvider(), provider,
                document_text_layout="preserve-basic-layout",
            )

            self.assertEqual(2, result.replaced_native_text_items)
            self.assertEqual(["filled and stroked", "fill only"], [
                request.text for request in provider.requests if not request.is_filename
            ])
            output = PdfReader(output_root / "document.pdf")
            stream = ContentStream(output.pages[0].get_contents(), output)
            self.assertEqual(1, sum(operator == b"BT" for _operands, operator in stream.operations))
            self.assertEqual(1, sum(operator == b"ET" for _operands, operator in stream.operations))

            text_rendering_mode = 0
            generated_modes: list[int] = []
            for operands, operator in stream.operations:
                if operator == b"Tr":
                    text_rendering_mode = int(operands[0])
                elif operator == b"Tf" and operands[0] == "/PipelineNoto":
                    generated_modes.append(text_rendering_mode)
            self.assertEqual([0, 0], generated_modes)
            self.assertTrue(any(
                operator == b"w" and float(operands[0]) == 2.0
                for operands, operator in stream.operations
            ))

    # Verifies FR-2026-08-23-01 and FR-2026-08-24-02.
    def test_basic_layout_pdf_uses_text_matrix_and_ctm_for_fitted_geometry(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"; output_root = root / "output"; input_root.mkdir()
            source = input_root / "document.pdf"
            writer = PdfWriter(); page = writer.add_blank_page(100, 100)
            contents = DecodedStreamObject()
            contents.set_data(
                b"q 2 0 0 2 0 0 cm BT /F1 10 Tf .5 0 0 .5 10 20 Tm (test) Tj ET Q"
            )
            page.replace_contents(ContentStream(contents, writer))
            with source.open("wb") as output_file: writer.write(output_file)

            result = self._run(
                input_root, output_root, _EmptyOcrProvider(), _RecordingReplacementProvider(),
                document_text_layout="preserve-basic-layout",
            )

            self.assertEqual(1, result.replaced_native_text_items)
            output = PdfReader(output_root / "document.pdf")
            stream = ContentStream(output.pages[0].get_contents(), output)
            pipeline_font_index = next(
                index for index, (operands, operator) in enumerate(stream.operations)
                if operator == b"Tf" and operands[0] == "/PipelineNoto"
            )
            generated_matrix = next(
                operands for operands, operator in stream.operations[pipeline_font_index:]
                if operator == b"Tm"
            )
            self.assertEqual([0.5, 0.0, 0.0, 0.5], [round(float(value), 4) for value in generated_matrix[:4]])

    # Verifies FR-2026-08-23-01.
    def test_basic_layout_pdf_preserves_an_axis_aligned_nonuniform_transform(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"; output_root = root / "output"; input_root.mkdir()
            source = input_root / "document.pdf"
            writer = PdfWriter(); page = writer.add_blank_page(100, 100)
            contents = DecodedStreamObject()
            contents.set_data(
                b"q 1.2 0 0 .8 0 0 cm BT /F1 10 Tf 1 0 0 1 10 70 Tm (cell) Tj ET Q"
            )
            page.replace_contents(ContentStream(contents, writer))
            with source.open("wb") as output_file: writer.write(output_file)

            provider = _RecordingReplacementProvider(replacement_text="mask")
            result = self._run(
                input_root, output_root, _EmptyOcrProvider(), provider,
                document_text_layout="preserve-basic-layout",
            )

            self.assertEqual(1, result.replaced_native_text_items)
            self.assertEqual(["cell"], [request.text for request in provider.requests if not request.is_filename])
            output = PdfReader(output_root / "document.pdf")
            stream = ContentStream(output.pages[0].get_contents(), output)
            self.assertFalse(any(
                operator == b"Tj" and operands and str(operands[0]) == "cell"
                for operands, operator in stream.operations
            ))
            pipeline_font_index = next(
                index for index, (operands, operator) in enumerate(stream.operations)
                if operator == b"Tf" and operands[0] == "/PipelineNoto"
            )
            generated_matrix = next(
                operands for operands, operator in stream.operations[pipeline_font_index:]
                if operator == b"Tm"
            )
            # The fitted font size absorbs the source's vertical 0.8 scale;
            # the matrix retains the resulting 1.5 horizontal-to-vertical ratio.
            self.assertEqual([1.25, 0.0, 0.0, 1.25], [round(float(value), 4) for value in generated_matrix[:4]])

    # Verifies FR-2026-08-23-01.
    def test_basic_layout_pdf_groups_compatible_text_objects_into_one_visual_block(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"; output_root = root / "output"; input_root.mkdir()
            source = input_root / "document.pdf"
            writer = PdfWriter(); page = writer.add_blank_page(200, 100)
            contents = DecodedStreamObject()
            contents.set_data(
                b"BT /F1 12 Tf 10 70 Td (first line) Tj ET BT /F1 12 Tf 10 56 Td (second line) Tj ET"
            )
            page.replace_contents(ContentStream(contents, writer))
            with source.open("wb") as output_file: writer.write(output_file)

            provider = _RecordingReplacementProvider(replacement_text="a fitted replacement")
            result = self._run(
                input_root, output_root, _EmptyOcrProvider(), provider,
                document_text_layout="preserve-basic-layout",
            )

            self.assertEqual(1, result.replaced_native_text_items)
            self.assertEqual(["first line\nsecond line"], [
                request.text for request in provider.requests if not request.is_filename
            ])

    # Verifies FR-2026-08-23-01.
    def test_basic_layout_pdf_keeps_an_undecodable_operation_but_replaces_a_reset_neighbour(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"; output_root = root / "output"; input_root.mkdir()
            source = input_root / "document.pdf"
            writer = PdfWriter(); page = writer.add_blank_page(100, 100)
            to_unicode = DecodedStreamObject()
            to_unicode.set_data(b"1 beginbfchar\n<0001> <0041>\nendbfchar\n")
            font = writer._add_object(DictionaryObject({
                NameObject("/Type"): NameObject("/Font"), NameObject("/Subtype"): NameObject("/Type0"),
                NameObject("/BaseFont"): NameObject("/SyntheticIdentity"), NameObject("/Encoding"): NameObject("/Identity-H"),
                NameObject("/ToUnicode"): writer._add_object(to_unicode),
            }))
            page[NameObject("/Resources")] = DictionaryObject({
                NameObject("/Font"): DictionaryObject({NameObject("/F1"): font})
            })
            contents = DecodedStreamObject()
            contents.set_data(b"BT /F1 12 Tf 10 10 Td <0002> Tj 1 0 0 1 10 30 Tm <0001> Tj ET")
            page.replace_contents(ContentStream(contents, writer))
            with source.open("wb") as output_file: writer.write(output_file)

            provider = _RecordingReplacementProvider()
            result = self._run(
                input_root, output_root, _EmptyOcrProvider(), provider,
                document_text_layout="preserve-basic-layout",
            )

            self.assertEqual(1, result.replaced_native_text_items)
            self.assertEqual(["A"], [request.text for request in provider.requests if not request.is_filename])
            output = PdfReader(output_root / "document.pdf")
            stream = ContentStream(output.pages[0].get_contents(), output)
            self.assertTrue(any(
                operator == b"Tj" and operands and str(operands[0]) == "\x02"
                for operands, operator in stream.operations
            ))

    # Verifies FR-2026-08-23-05.
    def test_basic_layout_pdf_advances_past_undecodable_identity_text(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"; output_root = root / "output"; input_root.mkdir()
            source = input_root / "document.pdf"
            writer = PdfWriter(); page = writer.add_blank_page(100, 100)
            to_unicode = DecodedStreamObject()
            to_unicode.set_data(b"1 beginbfchar\n<0001> <0041>\nendbfchar\n")
            descendant = writer._add_object(DictionaryObject({NameObject("/DW"): NumberObject(1000)}))
            descendants = writer._add_object(ArrayObject([descendant]))
            font = writer._add_object(DictionaryObject({
                NameObject("/Type"): NameObject("/Font"), NameObject("/Subtype"): NameObject("/Type0"),
                NameObject("/BaseFont"): NameObject("/SyntheticIdentity"), NameObject("/Encoding"): NameObject("/Identity-H"),
                NameObject("/ToUnicode"): writer._add_object(to_unicode),
                NameObject("/DescendantFonts"): descendants,
            }))
            page[NameObject("/Resources")] = DictionaryObject({
                NameObject("/Font"): DictionaryObject({NameObject("/F1"): font})
            })
            contents = DecodedStreamObject()
            contents.set_data(b"BT /F1 10 Tf 10 30 Td <0002> Tj <0001> Tj ET")
            page.replace_contents(ContentStream(contents, writer))
            with source.open("wb") as output_file: writer.write(output_file)

            provider = _RecordingReplacementProvider()
            result = self._run(
                input_root, output_root, _EmptyOcrProvider(), provider,
                document_text_layout="preserve-basic-layout",
            )

            self.assertEqual(1, result.replaced_native_text_items)
            self.assertEqual(["A"], [request.text for request in provider.requests if not request.is_filename])
            output = PdfReader(output_root / "document.pdf")
            stream = ContentStream(output.pages[0].get_contents(), output)
            unknown_index = next(
                index for index, (operands, operator) in enumerate(stream.operations)
                if operator == b"Tj" and operands and str(operands[0]) == "\x02"
            )
            self.assertNotEqual(b"Tr", stream.operations[unknown_index - 1][1])
            pipeline_font_index = next(
                index for index, (operands, operator) in enumerate(stream.operations)
                if operator == b"Tf" and operands[0] == "/PipelineNoto"
            )
            generated_matrix = next(
                operands for operands, operator in stream.operations[pipeline_font_index:]
                if operator == b"Tm"
            )
            self.assertEqual(20.0, float(generated_matrix[4]))

    # Verifies FR-2026-08-23-05.
    def test_basic_layout_pdf_keeps_barrier_when_undecodable_encoding_is_unknown(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"; output_root = root / "output"; input_root.mkdir()
            source = input_root / "document.pdf"
            writer = PdfWriter(); page = writer.add_blank_page(100, 100)
            to_unicode = DecodedStreamObject()
            to_unicode.set_data(b"1 beginbfchar\n<0001> <0041>\nendbfchar\n")
            encoding = writer._add_object(DecodedStreamObject())
            descendant = writer._add_object(DictionaryObject({NameObject("/DW"): NumberObject(1000)}))
            font = writer._add_object(DictionaryObject({
                NameObject("/Type"): NameObject("/Font"), NameObject("/Subtype"): NameObject("/Type0"),
                NameObject("/BaseFont"): NameObject("/SyntheticCustom"), NameObject("/Encoding"): encoding,
                NameObject("/ToUnicode"): writer._add_object(to_unicode),
                NameObject("/DescendantFonts"): ArrayObject([descendant]),
            }))
            page[NameObject("/Resources")] = DictionaryObject({
                NameObject("/Font"): DictionaryObject({NameObject("/F1"): font})
            })
            contents = DecodedStreamObject()
            contents.set_data(b"BT /F1 10 Tf 10 30 Td <0002> Tj <0001> Tj ET")
            page.replace_contents(ContentStream(contents, writer))
            with source.open("wb") as output_file: writer.write(output_file)

            provider = _RecordingReplacementProvider()
            result = self._run(
                input_root, output_root, _EmptyOcrProvider(), provider,
                document_text_layout="preserve-basic-layout",
            )

            self.assertEqual(0, result.replaced_native_text_items)
            self.assertFalse([request for request in provider.requests if not request.is_filename])

    # Verifies FR-2026-08-23-01.
    def test_basic_layout_pdf_recovers_at_line_position_resets_after_unknown_advance(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"; output_root = root / "output"; input_root.mkdir()
            source = input_root / "document.pdf"
            writer = PdfWriter(); page = writer.add_blank_page(100, 120)
            to_unicode = DecodedStreamObject()
            to_unicode.set_data(b"1 beginbfchar\n<0001> <0041>\nendbfchar\n")
            descendant = writer._add_object(DictionaryObject({NameObject("/DW"): NumberObject(1000)}))
            font = writer._add_object(DictionaryObject({
                NameObject("/Type"): NameObject("/Font"), NameObject("/Subtype"): NameObject("/Type0"),
                NameObject("/BaseFont"): NameObject("/SyntheticIdentity"), NameObject("/Encoding"): NameObject("/Identity-H"),
                NameObject("/ToUnicode"): writer._add_object(to_unicode),
                NameObject("/DescendantFonts"): ArrayObject([descendant]),
            }))
            page[NameObject("/Resources")] = DictionaryObject({
                NameObject("/Font"): DictionaryObject({NameObject("/F1"): font})
            })
            contents = DecodedStreamObject()
            contents.set_data(
                b"BT /F1 10 Tf 12 TL 10 105 Td <02> Tj 0 -15 Td <0001> Tj ET "
                b"BT /F1 10 Tf 12 TL 10 85 Td <02> Tj 0 -15 TD <0001> Tj ET "
                b"BT /F1 10 Tf 12 TL 10 65 Td <02> Tj T* <0001> Tj ET "
                b"BT /F1 10 Tf 12 TL 10 45 Td <02> Tj <0001> ' ET "
                b"BT /F1 10 Tf 12 TL 10 25 Td <02> Tj 0 0 <0001> \" ET"
            )
            page.replace_contents(ContentStream(contents, writer))
            with source.open("wb") as output_file:
                writer.write(output_file)

            provider = _RecordingReplacementProvider()
            result = self._run(
                input_root, output_root, _EmptyOcrProvider(), provider,
                document_text_layout="preserve-basic-layout",
            )

            self.assertEqual(5, result.replaced_native_text_items)
            self.assertEqual(
                ["A"] * 5,
                [request.text for request in provider.requests if not request.is_filename],
            )

    # Verifies FR-2026-08-23-06.
    def test_basic_layout_pdf_recovers_unicode_from_an_embedded_identity_font(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"; output_root = root / "output"; input_root.mkdir()
            source = input_root / "document.pdf"
            font_data = static_noto_bytes("sans-serif", False)
            source_face = skia.Typeface.MakeFromData(font_data)
            self.assertIsNotNone(source_face)
            glyph = skia.Font(source_face).textToGlyphs("A")[0]
            self.assertLess(glyph, 256)

            writer = PdfWriter(); page = writer.add_blank_page(100, 100)
            font_stream = DecodedStreamObject(); font_stream.set_data(font_data)
            descriptor = writer._add_object(DictionaryObject({
                NameObject("/FontFile2"): writer._add_object(font_stream),
            }))
            descendant = writer._add_object(DictionaryObject({
                NameObject("/DW"): NumberObject(1000),
                NameObject("/CIDToGIDMap"): NameObject("/Identity"),
                NameObject("/FontDescriptor"): descriptor,
            }))
            font = writer._add_object(DictionaryObject({
                NameObject("/Type"): NameObject("/Font"), NameObject("/Subtype"): NameObject("/Type0"),
                NameObject("/BaseFont"): NameObject("/SyntheticEmbeddedIdentity"),
                NameObject("/Encoding"): NameObject("/Identity-H"),
                NameObject("/DescendantFonts"): ArrayObject([descendant]),
            }))
            page[NameObject("/Resources")] = DictionaryObject({
                NameObject("/Font"): DictionaryObject({NameObject("/F1"): font})
            })
            source_font = cast(DictionaryObject, font.get_object())
            self.assertEqual(
                "A", _pdf_decode_composite_bytes(glyph.to_bytes(2, "big"), source_font)
            )
            self.assertEqual("A", _pdf_decode_composite_bytes(bytes([glyph]), source_font))
            self.assertIsNone(
                _pdf_decode_composite_bytes(bytes([glyph, glyph, glyph]), source_font)
            )
            non_identity_font = DictionaryObject(source_font)
            non_identity_font[NameObject("/Encoding")] = NameObject("/Identity-V")
            self.assertIsNone(_pdf_decode_composite_bytes(bytes([glyph]), non_identity_font))
            contents = DecodedStreamObject()
            contents.set_data(f"BT /F1 10 Tf 10 30 Td <{glyph:02X}> Tj ET".encode("ascii"))
            page.replace_contents(ContentStream(contents, writer))
            with source.open("wb") as output_file:
                writer.write(output_file)

            provider = _RecordingReplacementProvider()
            result = self._run(
                input_root, output_root, _EmptyOcrProvider(), provider,
                document_text_layout="preserve-basic-layout",
            )

            self.assertEqual(1, result.replaced_native_text_items)
            self.assertEqual(
                ["A"],
                [request.text for request in provider.requests if not request.is_filename],
            )

    # Verifies FR-2026-08-24-01.
    def test_type0_pdf_uses_direct_tounicode_parsing_when_the_helper_returns_whitespace(self) -> None:
        to_unicode = DecodedStreamObject()
        to_unicode.set_data(
            b"1 begincodespacerange\n<0000> <FFFF>\nendcodespacerange\n"
            b"1 beginbfrange\n<0001> <0001> <0041>\nendbfrange\n"
        )
        font = DictionaryObject({
            NameObject("/Type"): NameObject("/Font"),
            NameObject("/Subtype"): NameObject("/Type0"),
            NameObject("/Encoding"): NameObject("/Identity-H"),
            NameObject("/ToUnicode"): to_unicode,
        })
        with patch(
            "pipeline.folder_replacement.pdf.get_encoding",
            return_value=("utf-8", {"\x00\x01": " "}),
        ):
            self.assertEqual("A", _pdf_decode_composite_bytes(b"\x00\x01", font))

    # Verifies FR-2026-08-23-01.
    def test_basic_layout_pdf_replaces_tj_text_with_a_whitespace_fragment(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"; output_root = root / "output"; input_root.mkdir()
            source = input_root / "document.pdf"
            writer = PdfWriter(); page = writer.add_blank_page(100, 100)
            to_unicode = DecodedStreamObject()
            to_unicode.set_data(
                b"1 begincodespacerange\n<0000> <FFFF>\nendcodespacerange\n"
                b"2 beginbfchar\n<0001> <0041>\n<0002> <0020>\nendbfchar\n"
            )
            descendant = writer._add_object(DictionaryObject({NameObject("/DW"): NumberObject(1000)}))
            font = writer._add_object(DictionaryObject({
                NameObject("/Type"): NameObject("/Font"), NameObject("/Subtype"): NameObject("/Type0"),
                NameObject("/BaseFont"): NameObject("/SyntheticIdentity"), NameObject("/Encoding"): NameObject("/Identity-H"),
                NameObject("/ToUnicode"): writer._add_object(to_unicode),
                NameObject("/DescendantFonts"): ArrayObject([descendant]),
            }))
            page[NameObject("/Resources")] = DictionaryObject({
                NameObject("/Font"): DictionaryObject({NameObject("/F1"): font})
            })
            contents = DecodedStreamObject()
            contents.set_data(b"BT /F1 10 Tf 10 30 Td [<0001> 0 <0002> 0 <0001>] TJ ET")
            page.replace_contents(ContentStream(contents, writer))
            with source.open("wb") as output_file:
                writer.write(output_file)

            provider = _RecordingReplacementProvider()
            result = self._run(
                input_root, output_root, _EmptyOcrProvider(), provider,
                document_text_layout="preserve-basic-layout",
            )

            self.assertEqual(1, result.replaced_native_text_items)
            self.assertEqual(["A A"], [
                request.text for request in provider.requests if not request.is_filename
            ])

    # Verifies FR-2026-08-23-01.
    def test_basic_layout_pdf_reflows_one_visual_paragraph_and_replaces_its_source_operations(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"; output_root = root / "output"; input_root.mkdir()
            source = input_root / "document.pdf"
            writer = PdfWriter(); page = writer.add_blank_page(200, 100)
            contents = DecodedStreamObject()
            contents.set_data(
                b"BT /F1 12 Tf 14 TL 10 70 Td (first line) Tj T* (second line) Tj ET"
            )
            page.replace_contents(ContentStream(contents, writer))
            with source.open("wb") as output_file: writer.write(output_file)

            provider = _RecordingReplacementProvider(
                replacement_text="A replacement sentence that needs several wrapped lines."
            )
            result = self._run(
                input_root, output_root, _EmptyOcrProvider(), provider,
                document_text_layout="preserve-basic-layout",
            )

            self.assertEqual(1, result.replaced_native_text_items)
            requests = [request.text for request in provider.requests if not request.is_filename]
            self.assertEqual(["first line\nsecond line"], requests)
            output_reader = PdfReader(output_root / "document.pdf")
            stream = ContentStream(output_reader.pages[0].get_contents(), output_reader)
            self.assertGreaterEqual(sum(operator == b"Tj" for _operands, operator in stream.operations), 3)
            self.assertTrue(any(
                operator == b"Tf" and operands[0] == "/PipelineNoto"
                for operands, operator in stream.operations
            ))

    # Verifies FR-2026-08-23-01.
    def test_basic_layout_pdf_keeps_distant_same_line_labels_as_independent_regions(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"; output_root = root / "output"; input_root.mkdir()
            source = input_root / "document.pdf"
            writer = PdfWriter(); page = writer.add_blank_page(200, 100)
            contents = DecodedStreamObject()
            contents.set_data(b"BT /F1 12 Tf 10 70 Td (left) Tj 100 0 Td (right) Tj ET")
            page.replace_contents(ContentStream(contents, writer))
            with source.open("wb") as output_file: writer.write(output_file)

            provider = _RecordingReplacementProvider(replacement_text="replacement")
            result = self._run(
                input_root, output_root, _EmptyOcrProvider(), provider,
                document_text_layout="preserve-basic-layout",
            )

            self.assertEqual(2, result.replaced_native_text_items)
            requests = [request.text for request in provider.requests if not request.is_filename]
            self.assertEqual(["left", "right"], requests)

    # Verifies FR-2026-08-23-01.
    def test_basic_layout_pdf_retains_clipping_text_and_skips_the_provider(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"; output_root = root / "output"; input_root.mkdir()
            source = input_root / "document.pdf"
            writer = PdfWriter(); page = writer.add_blank_page(200, 100)
            contents = DecodedStreamObject()
            contents.set_data(b"BT /F1 12 Tf 7 Tr 10 70 Td (clip only) Tj ET")
            page.replace_contents(ContentStream(contents, writer))
            with source.open("wb") as output: writer.write(output)

            provider = _RecordingReplacementProvider()
            result = self._run(
                input_root, output_root, _EmptyOcrProvider(), provider,
                document_text_layout="preserve-basic-layout",
            )

            self.assertEqual(0, result.replaced_native_text_items)
            self.assertFalse([request for request in provider.requests if not request.is_filename])
            self.assertIn(b"clip only", (output_root / "document.pdf").read_bytes())
            self.assertFalse((output_root / "document.pdf.diagnostics.json").exists())

            debug_output_root = root / "debug-output"
            debug_provider = _RecordingReplacementProvider()
            debug_result = self._run(
                input_root,
                debug_output_root,
                _EmptyOcrProvider(),
                debug_provider,
                document_text_layout="preserve-basic-layout",
                diagnostics_enabled=True,
            )

            self.assertEqual(0, debug_result.replaced_native_text_items)
            self.assertFalse([request for request in debug_provider.requests if not request.is_filename])
            sidecar = json.loads(
                (debug_output_root / "document.pdf.diagnostics.json").read_text(encoding="utf-8")
            )
            entry = sidecar["entries"][0]
            self.assertEqual("retained", entry["kind"])
            self.assertEqual("pdf_text_rendering_mode_ineligible", entry["reason_code"])
            self.assertEqual("clip only", entry["source_text"])
            self.assertEqual(1, entry["page"])
            self.assertEqual("Tj", entry["operator"])

    # Verifies FR-2026-08-23-01.
    def test_basic_layout_pdf_preserves_a_right_angle_text_orientation(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"; output_root = root / "output"; input_root.mkdir()
            source = input_root / "document.pdf"
            writer = PdfWriter(); page = writer.add_blank_page(200, 100)
            contents = DecodedStreamObject()
            contents.set_data(b"BT /F1 12 Tf 0 1 -1 0 100 10 Tm (vertical) Tj ET")
            page.replace_contents(ContentStream(contents, writer))
            with source.open("wb") as output: writer.write(output)

            self._run(
                input_root, output_root, _EmptyOcrProvider(), _RecordingReplacementProvider(),
                document_text_layout="preserve-basic-layout",
            )

            output_reader = PdfReader(output_root / "document.pdf")
            stream = ContentStream(output_reader.pages[0].get_contents(), output_reader)
            self.assertTrue(any(
                operator == b"Tm" and [float(value) for value in operands[:4]] == [0.0, 1.0, -1.0, 0.0]
                for operands, operator in stream.operations
            ))

    # Verifies FR-2026-08-23-01.
    def test_basic_layout_pdf_replaces_text_inside_a_reusable_form_xobject(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"; output_root = root / "output"; input_root.mkdir()
            source = input_root / "document.pdf"
            writer = PdfWriter(); page = writer.add_blank_page(200, 100)
            form = DecodedStreamObject()
            form.set_data(b"BT /F1 12 Tf 10 10 Td (form text) Tj ET")
            form.update({
                NameObject("/Type"): NameObject("/XObject"), NameObject("/Subtype"): NameObject("/Form"),
                NameObject("/BBox"): ArrayObject([NumberObject(0), NumberObject(0), NumberObject(100), NumberObject(30)]),
            })
            form_reference = writer._add_object(form)
            page[NameObject("/Resources")] = DictionaryObject({
                NameObject("/XObject"): DictionaryObject({NameObject("/Form1"): form_reference})
            })
            contents = DecodedStreamObject(); contents.set_data(b"q /Form1 Do Q")
            page.replace_contents(ContentStream(contents, writer))
            with source.open("wb") as output_file: writer.write(output_file)

            provider = _RecordingReplacementProvider(replacement_text="fitted replacement")
            result = self._run(
                input_root, output_root, _EmptyOcrProvider(), provider,
                document_text_layout="preserve-basic-layout",
            )

            self.assertEqual(1, result.replaced_native_text_items)
            output_reader = PdfReader(output_root / "document.pdf")
            xobjects = cast(DictionaryObject, cast(DictionaryObject, output_reader.pages[0]["/Resources"])["/XObject"])
            updated_form = cast(DecodedStreamObject, xobjects["/Form1"].get_object())
            stream = ContentStream(updated_form, output_reader)
            self.assertTrue(any(
                operator == b"Tf" and operands[0] == "/PipelineNoto"
                for operands, operator in stream.operations
            ))

    # Verifies FR-2026-08-03-03.
    def test_replaces_type0_pdf_text_with_a_tounicode_map(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"; output_root = root / "output"; input_root.mkdir()
            source = input_root / "document.pdf"
            writer = PdfWriter(); page = writer.add_blank_page(100, 100)
            to_unicode = DecodedStreamObject()
            to_unicode.set_data(
                b"1 beginbfchar\n<0001> <0041>\nendbfchar\n"
            )
            font = writer._add_object(DictionaryObject({
                NameObject("/Type"): NameObject("/Font"),
                NameObject("/Subtype"): NameObject("/Type0"),
                NameObject("/BaseFont"): NameObject("/SyntheticIdentity"),
                NameObject("/Encoding"): NameObject("/Identity-H"),
                NameObject("/ToUnicode"): writer._add_object(to_unicode),
            }))
            page[NameObject("/Resources")] = DictionaryObject({
                NameObject("/Font"): DictionaryObject({NameObject("/F1"): font})
            })
            contents = DecodedStreamObject()
            contents.set_data(b"BT /F1 12 Tf 10 10 Td <0001> Tj ET")
            page.replace_contents(ContentStream(contents, writer))
            with source.open("wb") as output_file: writer.write(output_file)

            result = self._run(
                input_root,
                output_root,
                _EmptyOcrProvider(),
                _RecordingReplacementProvider(),
                document_text_layout="preserve-basic-layout",
            )

            self.assertEqual(1, result.replaced_native_text_items)
            basic_reader = PdfReader(output_root / "document.pdf")
            basic_stream = ContentStream(basic_reader.pages[0].get_contents(), basic_reader)
            self.assertTrue(any(
                operator == b"Tf" and operands[0] == "/PipelineNoto"
                for operands, operator in basic_stream.operations
            ))

            identity_basic_output = root / "identity-basic-output"
            self._run(
                input_root,
                identity_basic_output,
                _EmptyOcrProvider(),
                _RecordingReplacementProvider(replacement_text="A"),
                document_text_layout="preserve-basic-layout",
            )
            identity_reader = PdfReader(identity_basic_output / "document.pdf")
            identity_stream = ContentStream(identity_reader.pages[0].get_contents(), identity_reader)
            self.assertFalse(any(
                operator == b"Tf" and operands[0] == "/PipelineFallback"
                for operands, operator in identity_stream.operations
            ))
            self.assertIn(b"<0001>", (identity_basic_output / "document.pdf").read_bytes())

            source_font_output = root / "source-font-output"
            self._run(
                input_root,
                source_font_output,
                _EmptyOcrProvider(),
                _RecordingReplacementProvider(replacement_text="A"),
                document_text_layout="preserve-basic-layout-source-font",
            )
            source_font_reader = PdfReader(source_font_output / "document.pdf")
            source_font_stream = ContentStream(source_font_reader.pages[0].get_contents(), source_font_reader)
            self.assertFalse(any(
                operator == b"Tf" and operands[0] == "/PipelineFallback"
                for operands, operator in source_font_stream.operations
            ))
            self.assertIn(b"<0001>", (source_font_output / "document.pdf").read_bytes())

            fallback_output = root / "source-font-fallback-output"
            self._run(
                input_root,
                fallback_output,
                _EmptyOcrProvider(),
                _RecordingReplacementProvider(),
                document_text_layout="preserve-basic-layout-source-font",
            )
            fallback_reader = PdfReader(fallback_output / "document.pdf")
            fallback_stream = ContentStream(fallback_reader.pages[0].get_contents(), fallback_reader)
            self.assertTrue(any(
                operator == b"Tf" and operands[0] == "/PipelineNoto"
                for operands, operator in fallback_stream.operations
            ))
            fallback_fonts = cast(
                DictionaryObject, cast(DictionaryObject, fallback_reader.pages[0]["/Resources"])["/Font"]
            )
            fallback_font = cast(DictionaryObject, fallback_fonts["/PipelineNoto"].get_object())
            fallback_mapping = cast(
                DecodedStreamObject, fallback_font["/ToUnicode"].get_object()
            ).get_data()
            self.assertIn(b"<0001> <0023>", fallback_mapping)

    # Verifies FR-2026-08-23-04.
    def test_type0_pdf_widths_use_w_overrides_and_dw_fallback(self) -> None:
        to_unicode = DecodedStreamObject()
        to_unicode.set_data(
            b"4 beginbfchar\n<0001> <0041>\n<0002> <0042>\n<0003> <0043>\n<0004> <003F>\nendbfchar\n"
        )
        descendant = DictionaryObject({
            NameObject("/DW"): NumberObject(1000),
            NameObject("/W"): ArrayObject([
                NumberObject(1), ArrayObject([NumberObject(250), NumberObject(750)]),
                NumberObject(3), NumberObject(3), NumberObject(500),
            ]),
        })
        font = DictionaryObject({
            NameObject("/Subtype"): NameObject("/Type0"),
            NameObject("/Encoding"): NameObject("/Identity-H"),
            NameObject("/ToUnicode"): to_unicode,
            NameObject("/DescendantFonts"): ArrayObject([descendant]),
        })
        value = ByteStringObject(b"\x00\x01\x00\x02\x00\x03\x00\x04")
        self.assertEqual("ABC", _pdf_decode_composite_bytes(value[:6], font))
        advance = _pdf_text_advance(
            value,
            "ABC?",
            (NameObject("/F1"), NumberObject(10)),
            {"/F1": font},
            10.0,
            0.0,
            0.0,
            1.0,
        )
        self.assertEqual(25.0, advance)

    # Verifies FR-2026-08-23-04.
    def test_type0_pdf_widths_support_one_byte_codes_and_explicit_nonidentity_cids(self) -> None:
        to_unicode = DecodedStreamObject()
        to_unicode.set_data(b"2 beginbfchar\n<01> <0041>\n<02> <0042>\nendbfchar\n")
        encoding = DecodedStreamObject()
        encoding.set_data(b"2 begincidchar\n<01> 7\n<02> 9\nendcidchar\n")
        descendant = DictionaryObject({
            NameObject("/DW"): NumberObject(1000),
            NameObject("/W"): ArrayObject([
                NumberObject(7), ArrayObject([NumberObject(300)]),
                NumberObject(9), NumberObject(9), NumberObject(700),
            ]),
        })
        font = DictionaryObject({
            NameObject("/Subtype"): NameObject("/Type0"),
            NameObject("/Encoding"): encoding,
            NameObject("/ToUnicode"): to_unicode,
            NameObject("/DescendantFonts"): ArrayObject([descendant]),
        })
        value = ByteStringObject(b"\x01\x02")
        self.assertEqual("AB", _pdf_decode_composite_bytes(value, font))
        self.assertEqual(
            10.0,
            _pdf_text_advance(
                value,
                "AB",
                (NameObject("/F1"), NumberObject(10)),
                {"/F1": font},
                10.0,
                0.0,
                0.0,
                1.0,
            ),
        )

    # Verifies FR-2026-08-23-04.
    def test_type0_pdf_widths_fall_back_to_dw_when_cid_mapping_is_incomplete(self) -> None:
        to_unicode = DecodedStreamObject()
        to_unicode.set_data(b"2 beginbfchar\n<01> <0041>\n<02> <0042>\nendbfchar\n")
        encoding = DecodedStreamObject()
        encoding.set_data(b"1 begincidchar\n<01> 7\nendcidchar\n")
        descendant = DictionaryObject({
            NameObject("/DW"): NumberObject(1000),
            NameObject("/W"): ArrayObject([NumberObject(7), ArrayObject([NumberObject(300)])]),
        })
        font = DictionaryObject({
            NameObject("/Subtype"): NameObject("/Type0"),
            NameObject("/Encoding"): encoding,
            NameObject("/ToUnicode"): to_unicode,
            NameObject("/DescendantFonts"): ArrayObject([descendant]),
        })
        self.assertEqual(
            20.0,
            _pdf_text_advance(
                ByteStringObject(b"\x01\x02"),
                "AB",
                (NameObject("/F1"), NumberObject(10)),
                {"/F1": font},
                10.0,
                0.0,
                0.0,
                1.0,
            ),
        )

    # Verifies FR-2026-08-23-04.
    def test_basic_layout_pdf_type0_widths_control_the_fitted_replacement_size(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"; input_root.mkdir()
            narrow_source = input_root / "narrow.pdf"
            wide_source = input_root / "wide.pdf"

            def write_source(path: Path, widths: ArrayObject | None) -> None:
                writer = PdfWriter(); page = writer.add_blank_page(100, 100)
                to_unicode = DecodedStreamObject()
                to_unicode.set_data(b"2 beginbfchar\n<0001> <0041>\n<0002> <0042>\nendbfchar\n")
                descendant = DictionaryObject({NameObject("/DW"): NumberObject(1000)})
                if widths is not None:
                    descendant[NameObject("/W")] = widths
                font = writer._add_object(DictionaryObject({
                    NameObject("/Type"): NameObject("/Font"), NameObject("/Subtype"): NameObject("/Type0"),
                    NameObject("/BaseFont"): NameObject("/SyntheticIdentity"), NameObject("/Encoding"): NameObject("/Identity-H"),
                    NameObject("/ToUnicode"): writer._add_object(to_unicode),
                    NameObject("/DescendantFonts"): ArrayObject([writer._add_object(descendant)]),
                }))
                page[NameObject("/Resources")] = DictionaryObject({
                    NameObject("/Font"): DictionaryObject({NameObject("/F1"): font})
                })
                contents = DecodedStreamObject()
                contents.set_data(b"BT /F1 10 Tf 10 20 Td <00010002> Tj ET")
                page.replace_contents(ContentStream(contents, writer))
                with path.open("wb") as output_file: writer.write(output_file)

            write_source(narrow_source, ArrayObject([NumberObject(1), ArrayObject([NumberObject(250), NumberObject(750)])]))
            write_source(wide_source, None)

            output_root = root / "output"
            result = self._run(
                input_root,
                output_root,
                _EmptyOcrProvider(),
                _RecordingReplacementProvider(replacement_text="########"),
                document_text_layout="preserve-basic-layout",
            )

            self.assertEqual(2, result.replaced_native_text_items)

            def fitted_size(path: Path) -> float:
                output = PdfReader(path)
                stream = ContentStream(output.pages[0].get_contents(), output)
                return next(
                    float(operands[1]) for operands, operator in stream.operations
                    if operator == b"Tf" and operands[0] == "/PipelineNoto"
                )

            self.assertLess(fitted_size(output_root / "narrow.pdf"), fitted_size(output_root / "wide.pdf"))

    # Verifies FR-2026-08-03-03.
    def test_replaces_one_byte_type0_tj_fragments(self) -> None:
        """Use ToUnicode code widths, rather than the declared Identity encoding."""
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"; output_root = root / "output"; input_root.mkdir()
            source = input_root / "document.pdf"
            writer = PdfWriter(); page = writer.add_blank_page(100, 100)
            to_unicode = DecodedStreamObject()
            to_unicode.set_data(b"1 beginbfchar\n<01> <0041>\nendbfchar\n")
            font = writer._add_object(DictionaryObject({
                NameObject("/Type"): NameObject("/Font"),
                NameObject("/Subtype"): NameObject("/Type0"),
                NameObject("/BaseFont"): NameObject("/SyntheticIdentity"),
                NameObject("/Encoding"): NameObject("/Identity-H"),
                NameObject("/ToUnicode"): writer._add_object(to_unicode),
            }))
            page[NameObject("/Resources")] = DictionaryObject({
                NameObject("/Font"): DictionaryObject({NameObject("/F1"): font})
            })
            contents = DecodedStreamObject()
            contents.set_data(b"BT /F1 12 Tf 10 10 Td [<01> 20 <01>] TJ ET")
            page.replace_contents(ContentStream(contents, writer))
            with source.open("wb") as output: writer.write(output)

            result = self._run(
                input_root,
                output_root,
                _EmptyOcrProvider(),
                _RecordingReplacementProvider(),
                document_text_layout="preserve-basic-layout-source-font",
            )

            self.assertEqual(1, result.replaced_native_text_items)
            output_reader = PdfReader(output_root / "document.pdf")
            stream = ContentStream(output_reader.pages[0].get_contents(), output_reader)
            self.assertTrue(any(
                operator == b"Tf" and operands[0] == "/PipelineNoto"
                for operands, operator in stream.operations
            ))
            fonts = cast(DictionaryObject, cast(DictionaryObject, output_reader.pages[0]["/Resources"])["/Font"])
            type_zero = cast(DictionaryObject, fonts["/PipelineNoto"].get_object())
            mapping = cast(DecodedStreamObject, type_zero["/ToUnicode"].get_object()).get_data()
            self.assertIn(b"<0001> <0023>", mapping)

    # Verifies FR-2026-08-03-03.
    def test_source_font_pdf_falls_back_when_tounicode_lacks_replacement_glyph(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"; output_root = root / "output"; input_root.mkdir()
            source = input_root / "document.pdf"
            writer = PdfWriter(); page = writer.add_blank_page(100, 100)
            to_unicode = DecodedStreamObject()
            to_unicode.set_data(b"1 beginbfchar\n<41> <0041>\nendbfchar\n")
            font = writer._add_object(DictionaryObject({
                NameObject("/Type"): NameObject("/Font"),
                NameObject("/Subtype"): NameObject("/TrueType"),
                NameObject("/BaseFont"): NameObject("/SyntheticSubset"),
                NameObject("/Encoding"): NameObject("/WinAnsiEncoding"),
                NameObject("/ToUnicode"): writer._add_object(to_unicode),
            }))
            page[NameObject("/Resources")] = DictionaryObject({
                NameObject("/Font"): DictionaryObject({NameObject("/F1"): font})
            })
            contents = DecodedStreamObject()
            contents.set_data(b"BT /F1 12 Tf 10 10 Td (A) Tj ET")
            page.replace_contents(ContentStream(contents, writer))
            with source.open("wb") as output_file: writer.write(output_file)

            self._run(
                input_root,
                output_root,
                _EmptyOcrProvider(),
                _RecordingReplacementProvider(),
                document_text_layout="preserve-basic-layout-source-font",
            )

            output_reader = PdfReader(output_root / "document.pdf")
            stream = ContentStream(output_reader.pages[0].get_contents(), output_reader)
            self.assertTrue(any(
                operator == b"Tf" and operands[0] == "/PipelineNoto"
                for operands, operator in stream.operations
            ))
            fonts = cast(DictionaryObject, cast(DictionaryObject, output_reader.pages[0]["/Resources"])["/Font"])
            type_zero = cast(DictionaryObject, fonts["/PipelineNoto"].get_object())
            mapping = cast(DecodedStreamObject, type_zero["/ToUnicode"].get_object()).get_data()
            self.assertIn(b"<0001> <0023>", mapping)

    # Verifies FR-2026-08-04-13.
    def test_replaces_pptx_speaker_notes_in_every_document_text_layout_mode(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"
            output_root = root / "output"
            input_root.mkdir()
            source = input_root / "notes.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = "Slide title"
            presentation.save(str(source))
            self._add_speaker_note_part(source)

            for layout_mode in (
                "preserve-source-formatting",
                "preserve-basic-layout",
                "preserve-basic-layout-source-font",
            ):
                with self.subTest(document_text_layout=layout_mode):
                    destination_root = output_root / layout_mode
                    result = self._run(
                        input_root,
                        destination_root,
                        _EmptyOcrProvider(),
                        _RecordingReplacementProvider(replacement_text="Translated speaker note"),
                        document_text_layout=layout_mode,
                    )

                    self.assertEqual(1, result.processed_files)
                    self.assertGreaterEqual(result.replaced_native_text_items, 1)
                    with ZipFile(destination_root / "notes.pptx") as archive:
                        note_xml = archive.read("ppt/notesSlides/notesSlide1.xml")
                        slide_relationships = archive.read("ppt/slides/_rels/slide1.xml.rels")
                    self.assertIn(b"Translated speaker note", note_xml)
                    self.assertNotIn(b"Speaker note source", note_xml)
                    self.assertIn(b'uri="keep-note-extension"', note_xml)
                    self.assertIn(b"../notesSlides/notesSlide1.xml", slide_relationships)
                    Presentation(str(destination_root / "notes.pptx"))

    @staticmethod
    def _write_png(path: Path) -> None:
        Image.new("RGB", (30, 20), "white").save(path, "PNG")

    @staticmethod
    def _write_complete_docx(path: Path) -> None:
        """Write a complete synthetic DOCX package with a bounded text box."""
        with ZipFile(path, "w", ZIP_DEFLATED) as archive:
            archive.writestr(
                "[Content_Types].xml",
                b'''<?xml version="1.0" encoding="UTF-8"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  <Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>
  <Override PartName="/word/fontTable.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.fontTable+xml"/>
</Types>''',
            )
            archive.writestr(
                "_rels/.rels",
                b'''<?xml version="1.0" encoding="UTF-8"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rIdOfficeDocument" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>
</Relationships>''',
            )
            archive.writestr(
                "word/_rels/document.xml.rels",
                b'''<?xml version="1.0" encoding="UTF-8"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rIdExistingFontTable" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/fontTable" Target="fontTable.xml"/>
  <Relationship Id="rIdExistingSettings" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/settings" Target="settings.xml"/>
</Relationships>''',
            )
            archive.writestr(
                "word/settings.xml",
                b'''<?xml version="1.0" encoding="UTF-8"?>
<w:settings xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
  <w:writeProtection/><w:zoom/><w:bordersDoNotSurroundHeader/>
</w:settings>''',
            )
            archive.writestr(
                "word/fontTable.xml",
                b'''<?xml version="1.0" encoding="UTF-8"?>
<w:fonts xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
  <w:font w:name="Noto Sans JP"><w:altName w:val="Synthetic Noto"/></w:font>
</w:fonts>''',
            )
            archive.writestr(
                "word/document.xml",
                b'''<?xml version="1.0" encoding="UTF-8"?>
<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"
 xmlns:wp="http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing"
 xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
 xmlns:wps="http://schemas.microsoft.com/office/word/2010/wordprocessingShape">
  <w:body>
    <w:p><w:r><w:t>Flow text</w:t></w:r></w:p>
    <w:p><w:r><w:drawing><wp:inline>
      <wp:extent cx="914400" cy="457200"/><wp:docPr id="1" name="Synthetic text box"/>
      <a:graphic><a:graphicData uri="http://schemas.microsoft.com/office/word/2010/wordprocessingShape">
        <wps:wsp><wps:txbx><w:txbxContent>
          <w:p><w:r><w:rPr><w:b/><w:sz w:val="48"/></w:rPr><w:t>Bold text</w:t></w:r></w:p>
          <w:p><w:r><w:rPr><w:i/><w:sz w:val="48"/></w:rPr><w:t>Italic text</w:t></w:r></w:p>
          <w:p><w:r><w:rPr><w:sz w:val="48"/><w:u w:val="single"/></w:rPr><w:t>Underlined text</w:t></w:r></w:p>
        </w:txbxContent></wps:txbx></wps:wsp>
      </a:graphicData></a:graphic>
    </wp:inline></w:drawing></w:r></w:p>
    <w:sectPr><w:pgSz w:w="12240" w:h="15840"/><w:pgMar w:top="1440" w:right="1440" w:bottom="1440" w:left="1440"/></w:sectPr>
  </w:body>
</w:document>''',
            )

    @staticmethod
    def _add_reachable_smartart_data_part(path: Path) -> None:
        """Add synthetic canonical SmartArt labels linked from the first slide."""
        content_types_namespace = "http://schemas.openxmlformats.org/package/2006/content-types"
        relationships_namespace = "http://schemas.openxmlformats.org/package/2006/relationships"
        with ZipFile(path) as source_archive:
            entries = source_archive.infolist()
            payloads = {
                entry.filename: source_archive.read(entry.filename) for entry in entries
            }
        relationships = ElementTree.fromstring(payloads["ppt/slides/_rels/slide1.xml.rels"])
        ElementTree.SubElement(
            relationships,
            f"{{{relationships_namespace}}}Relationship",
            {
                "Id": "rIdSmartArtData",
                "Type": (
                    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/"
                    "diagramData"
                ),
                "Target": "../diagrams/data1.xml",
            },
        )
        payloads["ppt/slides/_rels/slide1.xml.rels"] = ElementTree.tostring(
            relationships, encoding="utf-8", xml_declaration=True
        )
        content_types = ElementTree.fromstring(payloads["[Content_Types].xml"])
        ElementTree.SubElement(
            content_types,
            f"{{{content_types_namespace}}}Override",
            {
                "PartName": "/ppt/diagrams/data1.xml",
                "ContentType": "application/vnd.ms-office.drawingml.diagramData+xml",
            },
        )
        payloads["[Content_Types].xml"] = ElementTree.tostring(
            content_types, encoding="utf-8", xml_declaration=True
        )
        payloads["ppt/diagrams/data1.xml"] = b"""<?xml version=\"1.0\" encoding=\"UTF-8\" standalone=\"yes\"?>
<dgm:dataModel xmlns:dgm=\"http://schemas.openxmlformats.org/drawingml/2006/diagram\" xmlns:a=\"http://schemas.openxmlformats.org/drawingml/2006/main\">
  <dgm:ptLst>
    <dgm:pt modelId=\"node-1\"><dgm:t><a:bodyPr/><a:lstStyle/><a:p><a:r><a:t>SmartArt first</a:t></a:r></a:p></dgm:t></dgm:pt>
    <dgm:pt modelId=\"node-2\"><dgm:t><a:bodyPr/><a:lstStyle/><a:p><a:r><a:t>SmartArt second</a:t></a:r></a:p></dgm:t></dgm:pt>
  </dgm:ptLst>
  <dgm:cxnLst/>
  <dgm:bg/>
  <dgm:whole/>
</dgm:dataModel>"""
        temporary_path = path.with_name(f".{path.name}.smartart-fixture.tmp")
        with ZipFile(temporary_path, "w", ZIP_DEFLATED) as destination_archive:
            for entry in entries:
                destination_archive.writestr(entry, payloads[entry.filename])
            destination_archive.writestr("ppt/diagrams/data1.xml", payloads["ppt/diagrams/data1.xml"])
        temporary_path.replace(path)

    @staticmethod
    def _add_speaker_note_part(path: Path) -> None:
        """Add one reachable synthetic speaker-note part to a PPTX fixture."""
        content_types_namespace = "http://schemas.openxmlformats.org/package/2006/content-types"
        relationships_namespace = "http://schemas.openxmlformats.org/package/2006/relationships"
        with ZipFile(path) as source_archive:
            entries = source_archive.infolist()
            payloads = {
                entry.filename: source_archive.read(entry.filename) for entry in entries
            }
        relationships = ElementTree.fromstring(payloads["ppt/slides/_rels/slide1.xml.rels"])
        ElementTree.SubElement(
            relationships,
            f"{{{relationships_namespace}}}Relationship",
            {
                "Id": "rIdSpeakerNotes",
                "Type": (
                    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/"
                    "notesSlide"
                ),
                "Target": "../notesSlides/notesSlide1.xml",
            },
        )
        payloads["ppt/slides/_rels/slide1.xml.rels"] = ElementTree.tostring(
            relationships, encoding="utf-8", xml_declaration=True
        )
        content_types = ElementTree.fromstring(payloads["[Content_Types].xml"])
        ElementTree.SubElement(
            content_types,
            f"{{{content_types_namespace}}}Override",
            {
                "PartName": "/ppt/notesSlides/notesSlide1.xml",
                "ContentType": "application/vnd.openxmlformats-officedocument.presentationml.notesSlide+xml",
            },
        )
        payloads["[Content_Types].xml"] = ElementTree.tostring(
            content_types, encoding="utf-8", xml_declaration=True
        )
        note_name = "ppt/notesSlides/notesSlide1.xml"
        payloads[note_name] = b'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:notes xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:cSld><p:spTree>
    <p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr><p:grpSpPr/>
    <p:sp><p:nvSpPr><p:cNvPr id="2" name="Notes Placeholder"/><p:cNvSpPr txBox="1"/><p:nvPr/></p:nvSpPr><p:spPr/>
      <p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:r><a:t>Speaker note source</a:t></a:r></a:p></p:txBody>
    </p:sp>
  </p:spTree></p:cSld>
  <p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr><p:extLst><p:ext uri="keep-note-extension"/></p:extLst>
</p:notes>'''
        temporary_path = path.with_name(f".{path.name}.notes-fixture.tmp")
        with ZipFile(temporary_path, "w", ZIP_DEFLATED) as destination_archive:
            for entry in entries:
                destination_archive.writestr(entry, payloads[entry.filename])
            destination_archive.writestr(note_name, payloads[note_name])
        temporary_path.replace(path)

    def _assert_drawingml_paragraph_property_order(self, data: bytes) -> None:
        """Verify the schema order that PowerPoint requires for ``a:pPr`` children."""
        namespace = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
        order = (
            "lnSpc",
            "spcBef",
            "spcAft",
            "buClrTx",
            "buClr",
            "buSzTx",
            "buSzPct",
            "buSzPts",
            "buFontTx",
            "buFont",
            "buNone",
            "buAutoNum",
            "buChar",
            "buBlip",
            "tabLst",
            "defRPr",
            "extLst",
        )
        ranks = {name: index for index, name in enumerate(order)}
        for properties in ElementTree.fromstring(data).iter(f"{namespace}pPr"):
            child_names = [child.tag.removeprefix(namespace) for child in properties]
            self.assertEqual(
                sorted(child_names, key=lambda name: ranks.get(name, len(ranks))),
                child_names,
            )

    def _assert_word_run_property_order(self, root: ElementTree.Element) -> None:
        """Verify the WordprocessingML ``CT_RPr`` child order."""
        namespace = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"
        order = (
            "rStyle", "rFonts", "b", "bCs", "i", "iCs", "caps", "smallCaps",
            "strike", "dstrike", "outline", "shadow", "emboss", "imprint", "noProof",
            "snapToGrid", "vanish", "webHidden", "color", "spacing", "w", "kern",
            "position", "sz", "szCs", "highlight", "u", "effect", "bdr", "shd",
            "fitText", "vertAlign", "rtl", "cs", "em", "lang", "eastAsianLayout",
            "specVanish", "oMath",
        )
        ranks = {name: index for index, name in enumerate(order)}
        for properties in root.iter(f"{namespace}rPr"):
            child_names = [child.tag.removeprefix(namespace) for child in properties]
            self.assertEqual(
                sorted(child_names, key=lambda name: ranks.get(name, len(ranks))),
                child_names,
            )

    @staticmethod
    def _sfnt_font_metadata(data: bytes) -> dict[str, dict[str, str]]:
        table_count = int.from_bytes(data[4:6], "big")
        tables = {
            data[12 + index * 16:16 + index * 16].decode("ascii"):
            int.from_bytes(data[20 + index * 16:24 + index * 16], "big")
            for index in range(table_count)
        }
        os2, post = tables["OS/2"], tables["post"]
        family_class = int.from_bytes(data[os2 + 30:os2 + 32], "big") >> 8
        family = {**{value: "roman" for value in range(1, 8)}, 8: "swiss", 9: "decorative", 10: "script"}.get(family_class, "auto")
        code_pages = int.from_bytes(data[os2 + 78:os2 + 82], "big")
        charset = next(
            (value for bit, value in ((17, "80"), (18, "81"), (19, "82"), (20, "86"), (21, "88"), (0, "00")) if code_pages & (1 << bit)),
            "00",
        )
        return {
            "panose1": {"val": data[os2 + 32:os2 + 42].hex().upper()},
            "charset": {"val": charset},
            "family": {"val": family},
            "pitch": {"val": "fixed" if int.from_bytes(data[post + 12:post + 16], "big") else "variable"},
            "sig": {
                "usb0": data[os2 + 42:os2 + 46].hex().upper(),
                "usb1": data[os2 + 46:os2 + 50].hex().upper(),
                "usb2": data[os2 + 50:os2 + 54].hex().upper(),
                "usb3": data[os2 + 54:os2 + 58].hex().upper(),
                "csb0": data[os2 + 78:os2 + 82].hex().upper(),
                "csb1": data[os2 + 82:os2 + 86].hex().upper(),
            },
        }

    def _assert_valid_drawingml_font_sizes(self, data: bytes) -> None:
        """Verify explicit run sizes stay within the OOXML DrawingML range."""
        namespace = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
        run_tags = {f"{namespace}rPr", f"{namespace}endParaRPr", f"{namespace}defRPr"}
        for properties in ElementTree.fromstring(data).iter():
            if properties.tag not in run_tags or "sz" not in properties.attrib:
                continue
            self.assertGreaterEqual(int(properties.attrib["sz"]), 100)
            self.assertLessEqual(int(properties.attrib["sz"]), 400_000)

    def _run(
        self,
        input_root: Path,
        output_root: Path,
        ocr_provider: _EmptyOcrProvider,
        replacement_provider: TextReplacementProvider,
        *,
        show_progress: bool = False,
        progress_factory: ProgressFactory | None = None,
        document_text_layout: str = "preserve-source-formatting",
        include_patterns: tuple[str, ...] = (),
        diagnostics_enabled: bool = False,
    ) -> FolderReplacementResult:
        typeface = skia.Typeface.MakeFromFile(str(FONT_PATH))
        if typeface is None:
            self.fail("Could not load test typeface.")
        return replace_input_folder(
            input_root,
            output_root,
            ocr_provider=ocr_provider,
            text_replacement_provider=replacement_provider,
            source_language="en",
            target_language="en",
            typeface=typeface,
            document_text_layout=document_text_layout,
            include_patterns=include_patterns,
            diagnostics_enabled=diagnostics_enabled,
            show_progress=show_progress,
            progress_factory=progress_factory,
        )


if __name__ == "__main__":
    unittest.main()
