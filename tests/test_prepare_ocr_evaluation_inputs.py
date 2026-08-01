#!/usr/bin/env python3
"""Regression tests for OCR evaluation input preparation."""

from __future__ import annotations

import importlib.util
import shutil
import sys
import tempfile
import unittest
from pathlib import Path
from zipfile import ZIP_DEFLATED, ZipFile

from docx import Document
from openpyxl import Workbook
from openpyxl.drawing.image import Image as SpreadsheetImage
from openpyxl.worksheet.worksheet import Worksheet
from PIL import Image
from pptx import Presentation


SCRIPT_PATH = Path(__file__).parents[1] / "scripts" / "prepare_ocr_evaluation_inputs.py"
SPECIFICATION = importlib.util.spec_from_file_location("ocr_preparation", SCRIPT_PATH)
assert SPECIFICATION is not None and SPECIFICATION.loader is not None
ocr_preparation = importlib.util.module_from_spec(SPECIFICATION)
sys.modules[SPECIFICATION.name] = ocr_preparation
SPECIFICATION.loader.exec_module(ocr_preparation)


class PrepareOcrEvaluationInputsTests(unittest.TestCase):
    # Verifies FR-2026-08-01-01.
    def test_copies_and_extracts_raster_images_then_skips_unchanged_documents(self) -> None:
        with tempfile.TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            source_root = root / "sample-data"
            output_root = root / "outputs" / "evaluations" / "ocr" / "input"
            language_root = source_root / "collection" / "ja"
            documents_root = language_root / "documents"
            documents_root.mkdir(parents=True)
            image_path = language_root / "source.png"
            Image.new("RGB", (8, 8), "red").save(image_path)

            self._create_docx(documents_root / "document.docx", image_path)
            self._create_pptx(documents_root / "presentation.pptx", image_path)
            self._create_xlsx(documents_root / "spreadsheet.xlsx", image_path)
            self._create_pdf(documents_root / "document.pdf", image_path)
            self._add_unreachable_pptx_media(documents_root / "presentation.pptx", image_path)

            unlabelled_image = source_root / "unlabelled.png"
            unlabelled_image.parent.mkdir(parents=True, exist_ok=True)
            Image.new("RGB", (8, 8), "blue").save(unlabelled_image)

            first_run = ocr_preparation.prepare_evaluation_inputs(source_root, output_root)

            self.assertEqual(1, first_run.copied_bitmaps)
            self.assertEqual(4, first_run.processed_documents)
            self.assertEqual(0, first_run.skipped_documents)
            self.assertEqual(4, first_run.extracted_images)
            self.assertEqual(0, first_run.removed_directories)
            self.assertEqual(
                image_path.read_bytes(),
                (output_root / "collection" / "ja" / "source.png").read_bytes(),
            )
            self.assertFalse((output_root / "unlabelled.png").exists())

            for filename in ("document.docx", "presentation.pptx", "spreadsheet.xlsx", "document.pdf"):
                document_output = output_root / "collection" / "ja" / "documents" / filename
                self.assertTrue((document_output / ".source.sha256").is_file())
                self.assertEqual(1, len(list(document_output.glob("image-*"))))
            self.assertFalse(
                any(
                    (output_root / "collection" / "ja" / "documents" / "presentation.pptx").glob(
                        "image-0002.*"
                    )
                )
            )

            second_run = ocr_preparation.prepare_evaluation_inputs(source_root, output_root)
            self.assertEqual(0, second_run.processed_documents)
            self.assertEqual(4, second_run.skipped_documents)

    # Verifies FR-2026-08-01-01.
    def test_removes_only_stale_directories_from_a_language_subtree(self) -> None:
        with tempfile.TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            source_root = root / "sample-data"
            language_root = source_root / "en"
            live_source_directory = language_root / "live"
            live_source_directory.mkdir(parents=True)
            Image.new("RGB", (8, 8), "green").save(live_source_directory / "source.png")

            output_root = root / "outputs" / "evaluations" / "ocr" / "input"
            stale_directory = output_root / "en" / "stale"
            stale_directory.mkdir(parents=True)
            (stale_directory / "extra.txt").write_text("remove", encoding="ascii")
            retained_directory = output_root / "en" / "live"
            retained_directory.mkdir(parents=True)
            (retained_directory / "extra.txt").write_text("retain", encoding="ascii")
            stale_language_directory = output_root / "fr"
            stale_language_directory.mkdir(parents=True)
            (stale_language_directory / "extra.txt").write_text("remove", encoding="ascii")
            stale_document_directory = output_root / "old-document.pptx"
            stale_document_directory.mkdir(parents=True)
            (stale_document_directory / "extra.txt").write_text("remove", encoding="ascii")

            result = ocr_preparation.prepare_evaluation_inputs(source_root, output_root)

            self.assertEqual(1, result.copied_bitmaps)
            self.assertEqual(3, result.removed_directories)
            self.assertFalse(stale_directory.exists())
            self.assertFalse(stale_language_directory.exists())
            self.assertFalse(stale_document_directory.exists())
            self.assertTrue((retained_directory / "extra.txt").is_file())

    @staticmethod
    def _create_docx(path: Path, image_path: Path) -> None:
        document = Document()
        document.add_picture(str(image_path))
        document.save(str(path))

    @staticmethod
    def _create_pptx(path: Path, image_path: Path) -> None:
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        slide.shapes.add_picture(str(image_path), 0, 0)
        presentation.save(str(path))

    @staticmethod
    def _create_xlsx(path: Path, image_path: Path) -> None:
        workbook = Workbook()
        worksheet = workbook.active
        if not isinstance(worksheet, Worksheet):
            raise RuntimeError("The workbook did not create a worksheet.")
        worksheet.add_image(SpreadsheetImage(str(image_path)), "A1")
        workbook.save(path)

    @staticmethod
    def _create_pdf(path: Path, image_path: Path) -> None:
        with Image.open(image_path) as image:
            image.save(path, "PDF")

    @staticmethod
    def _add_unreachable_pptx_media(path: Path, image_path: Path) -> None:
        rewritten_path = path.with_suffix(".rewritten")
        with ZipFile(path) as source_archive, ZipFile(
            rewritten_path, "w", ZIP_DEFLATED
        ) as destination_archive:
            for entry in source_archive.infolist():
                destination_archive.writestr(entry, source_archive.read(entry.filename))
            destination_archive.writestr("ppt/media/unreachable.png", image_path.read_bytes())
        shutil.move(rewritten_path, path)


if __name__ == "__main__":
    unittest.main()
