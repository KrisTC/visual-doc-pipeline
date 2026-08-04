"""PPTX-specific native text replacement."""

from __future__ import annotations

from collections.abc import Callable, Iterable, Iterator
from dataclasses import dataclass, replace
import os
from pathlib import Path
import posixpath
from typing import Protocol, cast
import xml.etree.ElementTree as ElementTree
from zipfile import ZIP_DEFLATED, ZipFile

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
from pptx.shapes.base import BaseShape
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.group import GroupShape
from pptx.table import Table, _Cell
from pptx.text.text import TextFrame, _Paragraph, _Run
from pptx.util import Length
# skia-python does not publish PEP 561 stubs; this is the native rendering boundary.
import skia  # type: ignore[import-not-found]

from pipeline.bounded_text_layout import (
    BoundedTextBox,
    BoundedTextParagraph,
    BoundedTextRun,
    noto_typefaces,
    replace_and_fit_text_box,
    source_occupied_text_box,
)
from pipeline.ocr import OcrProvider
from pipeline.folder_replacement.office_xml import replace_drawing_diagram_xml_text
from pipeline.text_replacement import TextReplacementProvider, TextReplacementRequest


_DRAWING_NAMESPACE = "http://schemas.openxmlformats.org/drawingml/2006/main"
_PARAGRAPH_PROPERTY_CHILD_ORDER = (
    "lnSpc",
    "spcBef",
    "spcAft",
    "buClrTx",
    "buClr",
    "buSzTx",
    "buSzPct",
    "buSzPts",
    "buFontTx",
    "buFont",
    "buNone",
    "buAutoNum",
    "buChar",
    "buBlip",
    "tabLst",
    "defRPr",
    "extLst",
)
_RUN_PROPERTY_CHILD_ORDER = (
    "noFill",
    "solidFill",
    "gradFill",
    "blipFill",
    "pattFill",
    "grpFill",
    "effectLst",
    "effectDag",
    "highlight",
    "uLnTx",
    "uLn",
    "uFillTx",
    "uFill",
    "latin",
    "ea",
    "cs",
    "sym",
    "hlinkClick",
    "hlinkMouseOver",
    "rtl",
    "extLst",
)
_OOXML_MINIMUM_FONT_SIZE_CENTIPOINTS = 100
_OOXML_MAXIMUM_FONT_SIZE_CENTIPOINTS = 400_000
_RELATIONSHIPS_NAMESPACE = "http://schemas.openxmlformats.org/package/2006/relationships"


def replace_pptx_file(
    source: Path,
    destination: Path,
    ocr: OcrProvider,
    replacement: TextReplacementProvider,
    source_language: str,
    target_language: str,
    typeface: skia.Typeface,
    completed: Callable[[str], None],
    document_text_layout: str = "preserve-source-formatting",
) -> tuple[int, int, int]:
    """Replace PPTX content, optionally fitting bounded slide text frames."""
    from pipeline.folder_replacement.processor import _replace_office_file

    smartart_parts, smartart_data_parts = _reachable_smartart_parts(source)
    if document_text_layout == "preserve-source-formatting":
        native_items, image_regions, retained_vectors = _replace_office_file(
            source,
            destination,
            ocr,
            replacement,
            source_language,
            target_language,
            typeface,
            completed,
            skip_native_xml_part=smartart_parts.__contains__,
        )
        native_items += _replace_smartart_data_parts(
            destination,
            smartart_data_parts,
            replacement,
            source_language,
            target_language,
        )
        return native_items, image_regions, retained_vectors
    if document_text_layout not in {
        "preserve-basic-layout",
        "preserve-basic-layout-source-font",
    }:
        raise ValueError(f"Unsupported document text layout mode: {document_text_layout!r}")
    preserve_source_font_family = document_text_layout == "preserve-basic-layout-source-font"

    # Preserve the established embedded bitmap and vector paths before python-pptx
    # rewrites supported slide text frames.
    native_items, image_regions, retained_vectors = _replace_office_file(
        source,
        destination,
        ocr,
        replacement,
        source_language,
        target_language,
        typeface,
        completed,
        replace_native_xml=False,
    )
    native_items += _replace_smartart_data_parts(
        destination,
        smartart_data_parts,
        replacement,
        source_language,
        target_language,
    )
    presentation = Presentation(str(destination))
    layout_typefaces = noto_typefaces()
    for slide in presentation.slides:
        native_items += _replace_slide_text_frames(
            slide.shapes,
            slide.slide_layout,
            replacement,
            source_language,
            target_language,
            layout_typefaces,
            preserve_source_font_family,
        )
    presentation.save(str(destination))
    completed("native text layout")
    return native_items, image_regions, retained_vectors


def _reachable_smartart_parts(source: Path) -> tuple[frozenset[str], frozenset[str]]:
    """Find reachable SmartArt package parts and their canonical data parts."""
    with ZipFile(source) as archive:
        archive_parts = frozenset(archive.namelist())
        pending: list[str | None] = [None]
        visited: set[str | None] = set()
        reachable: set[str] = set()
        data_parts: set[str] = set()
        while pending:
            source_part = pending.pop()
            if source_part in visited:
                continue
            visited.add(source_part)
            relationships_part = _relationships_part_name(source_part)
            if relationships_part not in archive_parts:
                continue
            try:
                relationships = ElementTree.fromstring(archive.read(relationships_part))
            except ElementTree.ParseError:
                continue
            for relationship in relationships:
                if relationship.tag != f"{{{_RELATIONSHIPS_NAMESPACE}}}Relationship":
                    continue
                if relationship.get("TargetMode") == "External":
                    continue
                target = relationship.get("Target")
                if target is None:
                    continue
                target_part = _relationship_target_part_name(source_part, target)
                if target_part is None or target_part not in archive_parts:
                    continue
                reachable.add(target_part)
                pending.append(target_part)
                if (relationship.get("Type") or "").endswith("/diagramData"):
                    data_parts.add(target_part)

    smartart_parts = {
        part for part in reachable if part.startswith("ppt/diagrams/")
    } | data_parts
    return frozenset(smartart_parts), frozenset(data_parts)


def _relationships_part_name(source_part: str | None) -> str:
    if source_part is None:
        return "_rels/.rels"
    parent, basename = posixpath.split(source_part)
    return posixpath.join(parent, "_rels", f"{basename}.rels")


def _relationship_target_part_name(source_part: str | None, target: str) -> str | None:
    base = "" if source_part is None else posixpath.dirname(source_part)
    candidate = target.lstrip("/") if target.startswith("/") else posixpath.join(base, target)
    normalized = posixpath.normpath(candidate)
    if normalized in {"", ".", ".."} or normalized.startswith("../"):
        return None
    return normalized


def _replace_smartart_data_parts(
    presentation_path: Path,
    data_parts: frozenset[str],
    replacement: TextReplacementProvider,
    source_language: str,
    target_language: str,
) -> int:
    """Replace only canonical SmartArt labels, not generated diagram drawings."""
    if not data_parts:
        return 0
    temporary_path = presentation_path.with_name(f".{presentation_path.name}.smartart.tmp")
    replaced_items = 0
    try:
        with (
            ZipFile(presentation_path) as source_archive,
            ZipFile(temporary_path, "w", ZIP_DEFLATED) as destination_archive,
        ):
            for entry in source_archive.infolist():
                data = source_archive.read(entry.filename)
                if entry.filename in data_parts:
                    data, replaced = replace_drawing_diagram_xml_text(
                        data,
                        replacement,
                        source_language,
                        target_language,
                    )
                    replaced_items += replaced
                destination_archive.writestr(entry, data)
        os.replace(temporary_path, presentation_path)
    finally:
        if temporary_path.exists():
            temporary_path.unlink()
    return replaced_items


def _replace_slide_text_frames(
    shapes: Iterable[BaseShape],
    slide_layout: object,
    replacement: TextReplacementProvider,
    source_language: str,
    target_language: str,
    typefaces: dict[str, skia.Typeface],
    preserve_source_font_family: bool,
) -> int:
    replaced = 0
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            replaced += _replace_slide_text_frames(
                cast(GroupShape, shape).shapes,
                slide_layout,
                replacement,
                source_language,
                target_language,
                typefaces,
                preserve_source_font_family,
            )
            continue
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            replaced += _replace_table_cells(
                cast(GraphicFrame, shape).table,
                replacement,
                source_language,
                target_language,
                typefaces,
                preserve_source_font_family,
            )
            continue
        if not shape.has_text_frame:
            continue
        text_shape = cast(_TextShape, shape)
        if not _has_text(text_shape):
            continue
        if _is_wordart(text_shape):
            replaced += _replace_text_frame_source_formatting(
                text_shape.text_frame,
                replacement,
                source_language,
                target_language,
            )
            continue
        text_box = _text_box(text_shape, slide_layout)
        fit_box = (
            source_occupied_text_box(text_box, typefaces)
            if _has_explicit_no_autofit(text_shape.text_frame)
            else text_box
        )
        fitted = replace_and_fit_text_box(
            fit_box,
            replacement,
            source_language,
            target_language,
            typefaces,
            preserve_source_font_family=preserve_source_font_family,
        )
        _write_explicit_text_frame(text_shape.text_frame, fitted.text_box)
        replaced += sum(
            1
            for paragraph in text_box.paragraphs
            if "".join(run.text for run in paragraph.runs).strip()
        )
    return replaced


def _replace_table_cells(
    table: Table,
    replacement: TextReplacementProvider,
    source_language: str,
    target_language: str,
    typefaces: dict[str, skia.Typeface],
    preserve_source_font_family: bool,
) -> int:
    """Replace each merge-origin table cell using its resolved grid rectangle."""
    replaced = 0
    for row_index, row in enumerate(table.rows):
        for column_index in range(len(table.columns)):
            cell = table.cell(row_index, column_index)
            if cell.is_spanned:
                continue
            bounds = _table_cell_bounds(table, cell, row_index, column_index)
            if bounds is None:
                replaced += _replace_text_frame_source_formatting(
                    cell.text_frame, replacement, source_language, target_language
                )
                continue
            width, height = bounds
            text_box = _table_cell_text_box(cell, width, height)
            if not _has_non_whitespace_text(text_box):
                continue
            fitted = replace_and_fit_text_box(
                text_box,
                replacement,
                source_language,
                target_language,
                typefaces,
                preserve_source_font_family=preserve_source_font_family,
            )
            _write_explicit_text_frame(
                cell.text_frame,
                fitted.text_box,
                write_text_frame_geometry=False,
            )
            replaced += sum(
                1
                for paragraph in text_box.paragraphs
                if "".join(run.text for run in paragraph.runs).strip()
            )
    return replaced


def _table_cell_bounds(
    table: Table, cell: _Cell, row_index: int, column_index: int
) -> tuple[int, int] | None:
    """Return the merge-origin cell's finite grid rectangle, when available."""
    span_width = cell.span_width if cell.is_merge_origin else 1
    span_height = cell.span_height if cell.is_merge_origin else 1
    columns = tuple(
        table.columns[index] for index in range(column_index, column_index + span_width)
    )
    rows = tuple(table.rows[index] for index in range(row_index, row_index + span_height))
    if len(columns) != span_width or len(rows) != span_height:
        return None
    width = sum(int(column.width) for column in columns)
    height = sum(int(row.height) for row in rows)
    return (width, height) if width > 0 and height > 0 else None


def _table_cell_text_box(cell: _Cell, width: int, height: int) -> BoundedTextBox:
    """Build a bounded layout model using the table cell's own margins."""
    text_frame = cell.text_frame
    return BoundedTextBox(
        width_emu=width,
        height_emu=height,
        margin_left_emu=int(cell.margin_left),
        margin_top_emu=int(cell.margin_top),
        margin_right_emu=int(cell.margin_right),
        margin_bottom_emu=int(cell.margin_bottom),
        text_direction=text_frame._element.bodyPr.get("vert"),
        paragraphs=tuple(
            _effective_table_cell_paragraph(paragraph) for paragraph in text_frame.paragraphs
        ),
    )


def _effective_table_cell_paragraph(paragraph: _Paragraph) -> BoundedTextParagraph:
    direct = _paragraph_properties(paragraph)
    style_properties = (paragraph._p.pPr,) if paragraph._p.pPr is not None else ()
    defaults = _run_defaults(style_properties)
    bullet_kind, bullet_marker = _effective_bullet(style_properties)
    return replace(
        direct,
        alignment=direct.alignment or _inherited_alignment(style_properties) or "left",
        margin_left_emu=(
            direct.margin_left_emu
            if direct.margin_left_emu is not None
            else _inherited_integer(style_properties, "marL")
        ),
        indent_emu=(
            direct.indent_emu
            if direct.indent_emu is not None
            else _inherited_integer(style_properties, "indent")
        ),
        bullet_kind=bullet_kind,
        bullet_marker=bullet_marker,
        empty_line_font_size_points=direct.empty_line_font_size_points
        or defaults.font_size_points,
        runs=tuple(_effective_run(run, defaults) for run in direct.runs),
    )


def _has_non_whitespace_text(text_box: BoundedTextBox) -> bool:
    return any(run.text.strip() for paragraph in text_box.paragraphs for run in paragraph.runs)


def _replace_text_frame_source_formatting(
    text_frame: TextFrame,
    replacement: TextReplacementProvider,
    source_language: str,
    target_language: str,
) -> int:
    """Replace a cell while retaining its original runs and text-frame settings."""
    replaced = 0
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if not run.text:
                continue
            run.text = replacement.replace(
                TextReplacementRequest(run.text, False, source_language, target_language)
            ).text
            replaced += 1
    return replaced


def _has_text(shape: "_TextShape") -> bool:
    return any(
        run.text.strip()
        for paragraph in shape.text_frame.paragraphs
        for run in paragraph.runs
    )


def _has_explicit_no_autofit(text_frame: TextFrame) -> bool:
    """Whether the source, rather than our output, explicitly disables autofit."""
    return text_frame._element.bodyPr.find(qn("a:noAutofit")) is not None


def _is_wordart(shape: "_TextShape") -> bool:
    """Whether fitted rewriting would discard WordArt text styling.

    The bounded-text writer intentionally creates new runs. It cannot safely
    reproduce DrawingML WordArt transforms or run paint/effect definitions, so
    these shapes retain source formatting while their text is replaced in place.
    """
    text_body = cast(_XmlElement, shape.text_frame._element)
    body_properties = text_body.find(qn("a:bodyPr"))
    if body_properties is not None and body_properties.find(qn("a:prstTxWarp")) is not None:
        return True
    unsupported_run_properties = {
        qn("a:noFill"),
        qn("a:solidFill"),
        qn("a:gradFill"),
        qn("a:blipFill"),
        qn("a:pattFill"),
        qn("a:grpFill"),
        qn("a:effectLst"),
        qn("a:effectDag"),
        qn("a:highlight"),
        qn("a:uLnTx"),
        qn("a:uLn"),
        qn("a:uFillTx"),
        qn("a:uFill"),
    }
    return any(
        child.tag in unsupported_run_properties
        for properties in text_body.iter()
        if properties.tag in {qn("a:rPr"), qn("a:defRPr"), qn("a:endParaRPr")}
        for child in properties
    )


def _text_box(shape: "_TextShape", slide_layout: object) -> BoundedTextBox:
    text_frame = shape.text_frame
    paragraphs = tuple(
        _effective_paragraph(paragraph, shape, slide_layout)
        for paragraph in text_frame.paragraphs
    )
    return BoundedTextBox(
        width_emu=int(shape.width),
        height_emu=int(shape.height),
        margin_left_emu=int(text_frame.margin_left),
        margin_top_emu=int(text_frame.margin_top),
        margin_right_emu=int(text_frame.margin_right),
        margin_bottom_emu=int(text_frame.margin_bottom),
        text_direction=text_frame._element.bodyPr.get("vert"),
        paragraphs=paragraphs,
    )


def _effective_paragraph(
    paragraph: _Paragraph, shape: "_TextShape", slide_layout: object
) -> BoundedTextParagraph:
    direct = _paragraph_properties(paragraph)
    style_properties = _paragraph_style_properties(
        shape, slide_layout, direct.level, paragraph._p.pPr
    )
    defaults = _run_defaults(style_properties)
    bullet_kind, bullet_marker = _effective_bullet(style_properties)
    return replace(
        direct,
        alignment=direct.alignment or _inherited_alignment(style_properties) or "left",
        margin_left_emu=(
            direct.margin_left_emu
            if direct.margin_left_emu is not None
            else _inherited_integer(style_properties, "marL")
        ),
        indent_emu=(
            direct.indent_emu
            if direct.indent_emu is not None
            else _inherited_integer(style_properties, "indent")
        ),
        bullet_kind=bullet_kind,
        bullet_marker=bullet_marker,
        empty_line_font_size_points=direct.empty_line_font_size_points
        or defaults.font_size_points,
        runs=tuple(_effective_run(run, defaults) for run in direct.runs),
    )


def _paragraph_properties(paragraph: _Paragraph) -> BoundedTextParagraph:
    line_spacing = paragraph.line_spacing
    if isinstance(line_spacing, Length):
        line_spacing_value, line_spacing_kind = float(line_spacing.pt), "points"
    elif isinstance(line_spacing, (int, float)):
        line_spacing_value, line_spacing_kind = float(line_spacing), "multiple"
    else:
        line_spacing_value, line_spacing_kind = None, None
    paragraph_properties = paragraph._p.pPr
    bullet_kind, bullet_marker = _paragraph_bullet(paragraph_properties)
    return BoundedTextParagraph(
        alignment=_enum_name(paragraph.alignment),
        space_before_points=_length_points(paragraph.space_before),
        space_after_points=_length_points(paragraph.space_after),
        line_spacing=line_spacing_value,
        line_spacing_kind=line_spacing_kind,
        level=paragraph.level,
        margin_left_emu=_xml_integer(paragraph_properties, "marL"),
        indent_emu=_xml_integer(paragraph_properties, "indent"),
        bullet_kind=bullet_kind,
        bullet_marker=bullet_marker,
        empty_line_font_size_points=_end_paragraph_font_size_points(paragraph),
        runs=tuple(_run_properties(run) for run in paragraph.runs),
    )


def _run_properties(run: _Run) -> BoundedTextRun:
    return BoundedTextRun(
        text=run.text,
        font_family=run.font.name,
        font_classification=_font_classification(run.font.name),
        font_size_points=_length_points(run.font.size),
        bold=run.font.bold,
        italic=run.font.italic,
        underline=_enum_name(run.font.underline),
        baseline=_xml_integer(run.font._element, "baseline"),
    )


@dataclass(frozen=True, slots=True)
class _RunDefaults:
    font_family: str | None = None
    font_size_points: float | None = None
    bold: bool | None = None
    italic: bool | None = None
    underline: str | None = None
    baseline: int | None = None


def _paragraph_style_properties(
    shape: "_TextShape", slide_layout: object, level: int, direct: object | None
) -> tuple[object, ...]:
    properties: list[object] = []
    layout = cast("_SlideLayout", slide_layout)
    master_level = _list_level_properties(_master_text_style(shape, layout), level)
    if master_level is not None:
        properties.append(master_level)
    layout_level = _list_level_properties(_layout_list_style(shape, layout), level)
    if layout_level is not None:
        properties.append(layout_level)
    text_frame_level = _list_level_properties(shape.text_frame._element.find(qn("a:lstStyle")), level)
    if text_frame_level is not None:
        properties.append(text_frame_level)
    if direct is not None:
        properties.append(direct)
    return tuple(properties)


def _master_text_style(shape: "_TextShape", layout: "_SlideLayout") -> object | None:
    if not shape.is_placeholder:
        return None
    placeholder_type = str(shape.placeholder_format.type)
    style_name = (
        "titleStyle"
        if "TITLE" in placeholder_type
        else "bodyStyle"
        if "BODY" in placeholder_type
        else "otherStyle"
    )
    styles = cast("_XmlElement", layout.slide_master._element).find(qn("p:txStyles"))
    return None if styles is None else styles.find(qn(f"p:{style_name}"))


def _layout_list_style(shape: "_TextShape", layout: "_SlideLayout") -> object | None:
    if not shape.is_placeholder:
        return None
    placeholder_index = shape.placeholder_format.idx
    for layout_shape in layout.shapes:
        if layout_shape.is_placeholder:
            candidate = cast(_TextShape, layout_shape)
            if candidate.placeholder_format.idx == placeholder_index:
                return cast(object, candidate.text_frame._element.find(qn("a:lstStyle")))
    return None


def _list_level_properties(list_style: object | None, level: int) -> object | None:
    if list_style is None:
        return None
    return cast("_XmlElement", list_style).find(qn(f"a:lvl{level + 1}pPr"))


def _run_defaults(properties: tuple[object, ...]) -> _RunDefaults:
    defaults = _RunDefaults()
    for paragraph_properties in properties:
        run_properties = cast("_XmlElement", paragraph_properties).find(qn("a:defRPr"))
        if run_properties is None:
            continue
        size = run_properties.get("sz")
        defaults = _RunDefaults(
            font_family=_font_from_properties(run_properties) or defaults.font_family,
            font_size_points=float(size) / 100.0 if size is not None else defaults.font_size_points,
            bold=_xml_boolean(run_properties.get("b"))
            if run_properties.get("b") is not None
            else defaults.bold,
            italic=_xml_boolean(run_properties.get("i"))
            if run_properties.get("i") is not None
            else defaults.italic,
            underline=run_properties.get("u") or defaults.underline,
            baseline=_xml_integer(run_properties, "baseline") or defaults.baseline,
        )
    return defaults


def _effective_run(run: BoundedTextRun, defaults: _RunDefaults) -> BoundedTextRun:
    family = run.font_family or defaults.font_family
    return replace(
        run,
        font_family=family,
        font_classification=_font_classification(family),
        font_size_points=run.font_size_points or defaults.font_size_points,
        bold=run.bold if run.bold is not None else defaults.bold,
        italic=run.italic if run.italic is not None else defaults.italic,
        underline=run.underline if run.underline is not None else defaults.underline,
        baseline=run.baseline if run.baseline is not None else defaults.baseline,
    )


def _effective_bullet(properties: tuple[object, ...]) -> tuple[str | None, str | None]:
    kind: str | None = None
    marker: str | None = None
    for paragraph_properties in properties:
        candidate_kind, candidate_marker = _paragraph_bullet(paragraph_properties)
        if candidate_kind is not None:
            kind, marker = candidate_kind, candidate_marker
    return kind, marker


def _paragraph_bullet(properties: object | None) -> tuple[str | None, str | None]:
    if properties is None:
        return None, None
    element = cast("_XmlElement", properties)
    if element.find(qn("a:buNone")) is not None:
        return "none", None
    character = element.find(qn("a:buChar"))
    if character is not None:
        return "character", character.get("char")
    if element.find(qn("a:buAutoNum")) is not None:
        return "automatic-number", None
    if element.find(qn("a:buBlip")) is not None:
        return "picture", None
    return None, None


def _write_explicit_text_frame(
    text_frame: TextFrame,
    text_box: BoundedTextBox,
    *,
    write_text_frame_geometry: bool = True,
) -> None:
    body_properties = text_frame._element.bodyPr
    if write_text_frame_geometry:
        body_properties.set("lIns", str(int(text_frame.margin_left)))
        body_properties.set("tIns", str(int(text_frame.margin_top)))
        body_properties.set("rIns", str(int(text_frame.margin_right)))
        body_properties.set("bIns", str(int(text_frame.margin_bottom)))
        body_properties.set("wrap", "square" if text_frame.word_wrap is not False else "none")
    body_properties.autofit = MSO_AUTO_SIZE.NONE
    for paragraph, explicit in zip(text_frame.paragraphs, text_box.paragraphs, strict=True):
        _write_paragraph(paragraph, explicit)


def _write_paragraph(paragraph: _Paragraph, explicit: BoundedTextParagraph) -> None:
    paragraph_properties = paragraph._p.get_or_add_pPr()
    paragraph_properties.set("lvl", str(explicit.level))
    paragraph_properties.set("algn", _drawing_alignment(explicit.alignment))
    _set_optional_integer(paragraph_properties, "marL", explicit.margin_left_emu)
    _set_optional_integer(paragraph_properties, "indent", explicit.indent_emu)
    _write_spacing(paragraph_properties, "a:spcBef", explicit.space_before_points, "points")
    _write_spacing(paragraph_properties, "a:spcAft", explicit.space_after_points, "points")
    _write_spacing(
        paragraph_properties,
        "a:lnSpc",
        explicit.line_spacing,
        explicit.line_spacing_kind,
    )
    _write_bullet(paragraph_properties, explicit)
    for child in tuple(paragraph._p):
        if child.tag != qn("a:pPr"):
            paragraph._p.remove(child)
    for run in explicit.runs:
        destination_run = paragraph.add_run()
        destination_run.text = run.text
        _write_run_properties(destination_run._r.get_or_add_rPr(), run)
    end_properties = paragraph._p.get_or_add_endParaRPr()
    empty_style = explicit.runs[0] if explicit.runs else BoundedTextRun(
        "", "Noto Sans JP", "sans-serif", explicit.empty_line_font_size_points, False, False, "none", 0
    )
    _write_run_properties(end_properties, empty_style)
    _reorder_children(paragraph_properties, _PARAGRAPH_PROPERTY_CHILD_ORDER)


def _write_spacing(
    paragraph_properties: object,
    container_tag: str,
    value: float | None,
    kind: str | None,
) -> None:
    element = cast("_XmlElement", paragraph_properties)
    existing = element.find(qn(container_tag))
    if existing is not None:
        element.remove(existing)
    if value is None:
        return
    container = _new_element(container_tag)
    child = _new_element("a:spcPts" if kind == "points" else "a:spcPct")
    child.set("val", str(round(value * (100.0 if kind == "points" else 100_000.0))))
    container.append(child)
    element.append(container)


def _write_bullet(properties: object, paragraph: BoundedTextParagraph) -> None:
    element = cast("_XmlElement", properties)
    for tag in ("a:buNone", "a:buChar"):
        existing = element.find(qn(tag))
        if existing is not None:
            element.remove(existing)
    if paragraph.bullet_kind == "none":
        element.append(_new_element("a:buNone"))
    elif paragraph.bullet_kind == "character":
        bullet = _new_element("a:buChar")
        bullet.set("char", paragraph.bullet_marker or "•")
        element.append(bullet)


def _write_run_properties(properties: object, run: BoundedTextRun) -> None:
    element = cast("_XmlElement", properties)
    element.set("sz", str(_ooxml_font_size_centipoints(run.font_size_points or 18.0)))
    element.set("b", "1" if run.bold else "0")
    element.set("i", "1" if run.italic else "0")
    element.set("u", _drawing_underline(run.underline))
    element.set("baseline", str(run.baseline or 0))
    for tag in ("a:latin", "a:ea"):
        child = element.find(qn(tag))
        if child is None:
            child = _new_element(tag)
            element.append(child)
        child.set("typeface", run.font_family or "Noto Sans JP")
    _reorder_children(element, _RUN_PROPERTY_CHILD_ORDER)


def _ooxml_font_size_centipoints(size_points: float) -> int:
    """Return a valid DrawingML ``sz`` value for an explicitly fitted run."""
    return min(
        _OOXML_MAXIMUM_FONT_SIZE_CENTIPOINTS,
        max(_OOXML_MINIMUM_FONT_SIZE_CENTIPOINTS, round(size_points * 100.0)),
    )


def _new_element(tag: str) -> "_XmlElement":
    """Create a DrawingML element through python-pptx's lxml element factory."""
    return cast(
        _XmlElement,
        parse_xml(
            f'<a:{tag.partition(":")[2]} xmlns:a="{_DRAWING_NAMESPACE}"/>'
        ),
    )


def _reorder_children(element: object, order: tuple[str, ...]) -> None:
    """Restore the schema-required order after writing DrawingML properties."""
    xml_element = cast(_XmlElement, element)
    ranks = {name: index for index, name in enumerate(order)}
    children = list(xml_element)
    children.sort(key=lambda child: ranks.get(child.tag.rsplit("}", 1)[-1], len(ranks)))
    for child in children:
        xml_element.remove(child)
    for child in children:
        xml_element.append(child)


def _drawing_alignment(value: str | None) -> str:
    return {"left": "l", "center": "ctr", "right": "r", "justify": "just"}.get(value or "left", "l")


def _drawing_underline(value: str | None) -> str:
    return {"single": "sng", "double": "dbl", "none": "none", "false": "none"}.get(
        value or "none", "none"
    )


def _set_optional_integer(element: object, name: str, value: int | None) -> None:
    xml_element = cast("_XmlElement", element)
    if value is None:
        xml_element.attrib.pop(name, None)
    else:
        xml_element.set(name, str(value))


def _inherited_integer(properties: tuple[object, ...], name: str) -> int | None:
    value: int | None = None
    for paragraph_properties in properties:
        candidate = _xml_integer(paragraph_properties, name)
        if candidate is not None:
            value = candidate
    return value


def _inherited_alignment(properties: tuple[object, ...]) -> str | None:
    value: str | None = None
    for paragraph_properties in properties:
        raw = cast("_XmlElement", paragraph_properties).get("algn")
        if raw is not None:
            value = {"l": "left", "ctr": "center", "r": "right", "just": "justify"}.get(raw, raw)
    return value


def _font_from_properties(properties: object) -> str | None:
    element = cast("_XmlElement", properties)
    for tag in ("a:latin", "a:ea"):
        candidate = element.find(qn(tag))
        if candidate is not None and candidate.get("typeface"):
            return candidate.get("typeface")
    return None


def _xml_integer(element: object | None, name: str) -> int | None:
    if element is None:
        return None
    value = cast("_XmlElement", element).get(name)
    return int(value) if value is not None else None


def _xml_boolean(value: str | None) -> bool | None:
    return None if value is None else value not in {"0", "false", "False"}


def _end_paragraph_font_size_points(paragraph: _Paragraph) -> float | None:
    size = _xml_integer(paragraph._p.endParaRPr, "sz")
    return None if size is None else size / 100.0


def _length_points(value: Length | None) -> float | None:
    return None if value is None else float(value.pt)


def _enum_name(value: object) -> str | None:
    if value is None:
        return None
    return str(value).split(" ", 1)[0].lower().replace("_", "-")


def _font_classification(family: str | None) -> str:
    normalized = (family or "").lower()
    if any(marker in normalized for marker in ("mono", "courier", "console", "code")):
        return "fixed-width"
    if any(marker in normalized for marker in ("serif", "roman", "times", "georgia")):
        return "serif"
    return "sans-serif"


class _XmlElement(Protocol, Iterable["_XmlElement"]):
    attrib: dict[str, str]
    tag: str

    def append(self, element: object) -> None: ...
    def remove(self, element: object) -> None: ...
    def find(self, path: str) -> "_XmlElement | None": ...
    def get(self, key: str) -> str | None: ...
    def set(self, key: str, value: str) -> None: ...
    def iter(self) -> Iterator["_XmlElement"]: ...
    def __iter__(self) -> Iterator["_XmlElement"]:
        raise NotImplementedError


class _TextShape(Protocol):
    width: int
    height: int
    is_placeholder: bool
    text_frame: TextFrame
    placeholder_format: "_PlaceholderFormat"


class _PlaceholderFormat(Protocol):
    idx: int
    type: object


class _SlideLayout:
    shapes: Iterable[BaseShape]
    slide_master: "_SlideMaster"


class _SlideMaster:
    _element: object
