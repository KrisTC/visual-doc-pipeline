#!/usr/bin/env python3
"""Generate local PowerPoint native-text layout evaluations without OCR."""

from __future__ import annotations

import argparse
from collections.abc import Iterable
from dataclasses import asdict, dataclass, replace
import html
from io import BytesIO
import json
from math import ceil
import os
from pathlib import Path
import re
import sys
from typing import cast
import unicodedata


PROJECT_ROOT = Path(__file__).resolve().parents[1]
if str(PROJECT_ROOT) not in sys.path:
    sys.path.insert(0, str(PROJECT_ROOT))


from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.exc import PackageNotFoundError
from pptx.oxml.ns import qn
from pptx.shapes.autoshape import Shape
from pptx.shapes.base import BaseShape
from pptx.shapes.group import GroupShape
from pptx.text.text import TextFrame, _Paragraph, _Run
from pptx.util import Length
from PIL import Image
# skia-python does not publish PEP 561 stubs; this is the native rendering boundary.
import skia  # type: ignore[import-not-found]

from pipeline.ocr.languages import discover_language_directories
from pipeline.bounded_text_layout import (
    BoundedTextBox,
    BoundedTextParagraph,
    BoundedTextRun,
    fit_explicit_noto_text_box,
    source_occupied_text_box,
)
from pipeline.text_replacement import (
    TextReplacementProviderFactory,
    TextReplacementRequest,
)
from pipeline.text_replacement.provider import TextReplacementProvider


DEFAULT_INPUT_ROOT = Path("sample-data")
DEFAULT_OUTPUT_ROOT = Path("outputs/evaluations/text-replacement")
FONT_DIRECTORY = PROJECT_ROOT / "tests" / "assets" / "fonts"
SANS_FONT_PATH = FONT_DIRECTORY / "NotoSansJP[wght].ttf"
SERIF_FONT_PATH = FONT_DIRECTORY / "NotoSerifJP[wght].ttf"
MONO_FONT_PATH = FONT_DIRECTORY / "NotoSansMono[wdth,wght].ttf"
EMU_PER_PIXEL = 9_525
PIXELS_PER_POINT = 4.0 / 3.0
DEFAULT_FONT_SIZE_POINTS = 18.0
DEFAULT_TARGET_LANGUAGE = "en"
_DRAWING_NAMESPACE = "http://schemas.openxmlformats.org/drawingml/2006/main"
_TOKEN_PATTERN = re.compile(r"[\n\v]|\S+\s*|\s+")
_DERIVED_FIT_GUIDE_COLOR = skia.ColorSetRGB(217, 119, 6)


@dataclass(frozen=True, slots=True)
class TextRunProperties:
    """Directly specified typography for one PowerPoint text run."""

    text: str
    font_family: str | None
    font_classification: str
    font_size_points: float | None
    bold: bool | None
    italic: bool | None
    underline: str | None
    baseline: int | None


@dataclass(frozen=True, slots=True)
class ParagraphProperties:
    """Directly specified layout settings and runs for one paragraph."""

    alignment: str | None
    space_before_points: float | None
    space_after_points: float | None
    line_spacing: float | None
    line_spacing_kind: str | None
    level: int
    margin_left_emu: int | None
    indent_emu: int | None
    bullet_kind: str | None
    bullet_marker: str | None
    empty_line_font_size_points: float | None
    runs: tuple[TextRunProperties, ...]


@dataclass(frozen=True, slots=True)
class TextBoxProperties:
    """The source layout properties required by the native-text evaluator."""

    source: str
    left_emu: int
    top_emu: int
    width_emu: int
    height_emu: int
    shape_rotation_degrees: float
    text_direction: str | None
    effective_text_rotation_degrees: float
    margin_left_emu: int
    margin_top_emu: int
    margin_right_emu: int
    margin_bottom_emu: int
    word_wrap: bool | None
    vertical_alignment: str | None
    autofit_mode: str
    explicit_no_autofit: bool
    autofit_font_scale: int | None
    autofit_line_spacing_reduction: int | None
    paragraphs: tuple[ParagraphProperties, ...]


@dataclass(frozen=True, slots=True)
class _DrawStyle:
    classification: str
    size_pixels: float
    bold: bool
    italic: bool
    underline: str | None
    baseline: int


@dataclass(frozen=True, slots=True)
class _DrawSegment:
    text: str
    style: _DrawStyle


@dataclass(frozen=True, slots=True)
class _DrawLine:
    segments: tuple[_DrawSegment, ...]
    width: float
    height: float
    paragraph: ParagraphProperties
    is_first_line: bool


@dataclass(frozen=True, slots=True)
class _LayoutFit:
    """The selected pre-render layout and its bounds-fit result."""

    lines: tuple[_DrawLine, ...]
    font_scale: float
    fit_status: str


@dataclass(frozen=True, slots=True)
class _ReplacementFitting:
    """The production-equivalent fit result and its evaluator-only guide."""

    layout_fit: _LayoutFit
    fitting_box: BoundedTextBox
    derived_from_source: bool

    @property
    def canvas_height_emu(self) -> int:
        return max(self.fitting_box.height_emu, self.source_height_emu)

    source_height_emu: int


@dataclass(frozen=True, slots=True)
class _TextBoxEvaluation:
    """Direct source properties and the inherited formatting used for evaluation."""

    source_properties: TextBoxProperties
    effective_properties: TextBoxProperties


@dataclass(frozen=True, slots=True)
class _TextBoxArtifact:
    """Local evaluation artifacts produced for one source text box."""

    rendering_path: Path
    properties_path: Path
    explicit_properties_path: Path
    replacement_artifacts: tuple["_ProviderTextBoxArtifact", ...]


@dataclass(frozen=True, slots=True)
class _ProviderTextBoxArtifact:
    """One provider-specific native-text replacement rendering and definition."""

    provider_name: str
    rendering_path: Path
    explicit_properties_path: Path


@dataclass(slots=True)
class TextReplacementEvaluationRunResult:
    """Counts produced by one native-text layout evaluation run."""

    processed_presentations: int = 0
    skipped_presentations: int = 0
    written_pages: int = 0
    rendered_text_boxes: int = 0


def evaluate_text_replacement_examples(
    input_root: Path,
    output_root: Path,
    target_language: str = DEFAULT_TARGET_LANGUAGE,
) -> TextReplacementEvaluationRunResult:
    """Evaluate each PowerPoint text box below ``input_root`` without OCR."""
    if not input_root.is_dir():
        raise ValueError(f"Input root does not exist or is not a directory: {input_root}")
    typefaces = _load_typefaces()
    factory = TextReplacementProviderFactory.discover_default_plugins()
    provider_names = factory.local_evaluation_provider_names
    result = TextReplacementEvaluationRunResult()
    for source_path, source_language in _presentation_paths(input_root):
        if source_path.name.startswith("~$"):
            continue
        result.processed_presentations += 1
        try:
            text_box_evaluations = tuple(_presentation_text_boxes(source_path))
            output_path = (output_root / source_path.relative_to(input_root)).with_suffix(
                ".html"
            )
            providers = {
                provider_name: factory.create(provider_name) for provider_name in provider_names
            }
            artifacts = _render_text_boxes(
                text_box_evaluations,
                typefaces,
                output_path,
                providers,
                source_language,
                target_language,
            )
            _write_html_page(
                output_path,
                source_path.relative_to(input_root),
                tuple(item.source_properties for item in text_box_evaluations),
                artifacts,
                provider_names,
                source_language,
                target_language,
            )
        except (OSError, PackageNotFoundError, RuntimeError, ValueError) as error:
            result.skipped_presentations += 1
            print(f"Skipping {source_path}: {error}.")
            continue
        result.written_pages += 1
        result.rendered_text_boxes += len(text_box_evaluations)
    return result


def _presentation_paths(input_root: Path) -> Iterable[tuple[Path, str]]:
    """Yield PPTX files beneath the OCR evaluator's established language directories."""
    for language_directory in discover_language_directories(input_root):
        for source_path in sorted(language_directory.rglob("*.pptx")):
            yield source_path, language_directory.name


def _load_typefaces() -> dict[str, skia.Typeface]:
    """Load the committed Noto serif, sans, and fixed-width typefaces."""
    return {
        "sans-serif": _load_typeface(SANS_FONT_PATH),
        "serif": _load_typeface(SERIF_FONT_PATH),
        "fixed-width": _load_typeface(MONO_FONT_PATH),
    }


def _load_typeface(path: Path) -> skia.Typeface:
    if not path.is_file():
        raise RuntimeError(f"The native-text evaluator requires {path}.")
    typeface = skia.Typeface.MakeFromFile(str(path))
    if typeface is None:
        raise RuntimeError(f"Skia could not load evaluator typeface {path}.")
    return typeface


def _presentation_text_boxes(source_path: Path) -> Iterable[_TextBoxEvaluation]:
    """Yield eligible text boxes, including text boxes in grouped shapes."""
    presentation = Presentation(str(source_path))
    for slide_number, slide in enumerate(presentation.slides, 1):
        yield from _slide_text_boxes(slide.shapes, slide_number, (), slide.slide_layout)


def _slide_text_boxes(
    shapes: Iterable[BaseShape],
    slide_number: int,
    group_path: tuple[int, ...],
    slide_layout: object,
) -> Iterable[_TextBoxEvaluation]:
    for shape_number, shape in enumerate(shapes, 1):
        path = group_path + (shape_number,)
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            group_shape = cast(GroupShape, shape)
            yield from _slide_text_boxes(group_shape.shapes, slide_number, path, slide_layout)
            continue
        if not shape.has_text_frame:
            continue
        text_shape = cast(Shape, shape)
        if not _has_non_whitespace_run_text(text_shape.text_frame):
            continue
        source_properties = _text_box_properties(
            text_shape, f"slide {slide_number}, shape {'/'.join(map(str, path))}"
        )
        yield _TextBoxEvaluation(
            source_properties,
            _effective_text_box_properties(text_shape, source_properties, slide_layout),
        )


def _has_non_whitespace_run_text(text_frame: TextFrame) -> bool:
    """Exclude empty PowerPoint placeholders but preserve blank content within a box."""
    return any(run.text.strip() for paragraph in text_frame.paragraphs for run in paragraph.runs)


def _text_box_properties(shape: Shape, source: str) -> TextBoxProperties:
    text_frame = shape.text_frame
    body_properties = text_frame._element.bodyPr
    text_direction = body_properties.get("vert")
    shape_rotation = float(shape.rotation)
    return TextBoxProperties(
        source=source,
        left_emu=int(shape.left),
        top_emu=int(shape.top),
        width_emu=int(shape.width),
        height_emu=int(shape.height),
        shape_rotation_degrees=shape_rotation,
        text_direction=text_direction,
        effective_text_rotation_degrees=_effective_text_rotation(shape_rotation, text_direction),
        margin_left_emu=int(text_frame.margin_left),
        margin_top_emu=int(text_frame.margin_top),
        margin_right_emu=int(text_frame.margin_right),
        margin_bottom_emu=int(text_frame.margin_bottom),
        word_wrap=text_frame.word_wrap,
        vertical_alignment=_enum_name(text_frame.vertical_anchor),
        autofit_mode=_autofit_mode(text_frame),
        explicit_no_autofit=body_properties.find(qn("a:noAutofit")) is not None,
        autofit_font_scale=_autofit_integer(body_properties, "fontScale"),
        autofit_line_spacing_reduction=_autofit_integer(body_properties, "lnSpcReduction"),
        paragraphs=tuple(_paragraph_properties(paragraph) for paragraph in text_frame.paragraphs),
    )


def _effective_text_box_properties(
    shape: Shape, source_properties: TextBoxProperties, slide_layout: object
) -> TextBoxProperties:
    """Resolve list-style defaults needed to make editable layout self-contained."""
    paragraphs = tuple(
        _effective_paragraph_properties(paragraph, source_paragraph, shape, slide_layout)
        for paragraph, source_paragraph in zip(
            shape.text_frame.paragraphs, source_properties.paragraphs, strict=True
        )
    )
    return replace(source_properties, paragraphs=paragraphs)


def _effective_paragraph_properties(
    paragraph: _Paragraph,
    source_properties: ParagraphProperties,
    shape: Shape,
    slide_layout: object,
) -> ParagraphProperties:
    style_properties = _paragraph_style_properties(
        shape, slide_layout, source_properties.level, paragraph._p.pPr
    )
    bullet_kind, bullet_marker = _effective_bullet(style_properties)
    defaults = _run_defaults(style_properties)
    runs = tuple(
        _effective_run_properties(run, defaults) for run in source_properties.runs
    )
    return replace(
        source_properties,
        alignment=source_properties.alignment
        or _inherited_alignment(style_properties),
        margin_left_emu=source_properties.margin_left_emu
        if source_properties.margin_left_emu is not None
        else _inherited_integer(style_properties, "marL"),
        indent_emu=source_properties.indent_emu
        if source_properties.indent_emu is not None
        else _inherited_integer(style_properties, "indent"),
        bullet_kind=bullet_kind,
        bullet_marker=bullet_marker,
        empty_line_font_size_points=source_properties.empty_line_font_size_points
        or defaults.font_size_points,
        runs=runs,
    )


def _paragraph_style_properties(
    shape: Shape,
    slide_layout: object,
    level: int,
    direct_properties: object | None,
) -> tuple[object, ...]:
    """Return master, layout, text-frame, then direct list properties in precedence order."""
    properties: list[object] = []
    layout = cast("_SlideLayout", slide_layout)
    master_text_style = _master_text_style(shape, layout)
    master_level = _list_level_properties(master_text_style, level)
    if master_level is not None:
        properties.append(master_level)
    layout_level = _list_level_properties(_layout_list_style(shape, layout), level)
    if layout_level is not None:
        properties.append(layout_level)
    text_frame_level = _list_level_properties(
        shape.text_frame._element.find(qn("a:lstStyle")), level
    )
    if text_frame_level is not None:
        properties.append(text_frame_level)
    if direct_properties is not None:
        properties.append(direct_properties)
    return tuple(properties)


def _master_text_style(shape: Shape, layout: "_SlideLayout") -> object | None:
    if not shape.is_placeholder:
        return None
    placeholder_type = str(shape.placeholder_format.type)
    if "TITLE" in placeholder_type:
        style_name = "titleStyle"
    elif "BODY" in placeholder_type:
        style_name = "bodyStyle"
    else:
        style_name = "otherStyle"
    text_styles = cast("_XmlSearchElement", layout.slide_master._element).find(qn("p:txStyles"))
    return None if text_styles is None else text_styles.find(qn(f"p:{style_name}"))


def _layout_list_style(shape: Shape, layout: "_SlideLayout") -> object | None:
    if not shape.is_placeholder:
        return None
    placeholder_index = shape.placeholder_format.idx
    for base_layout_shape in layout.shapes:
        layout_shape = cast(Shape, base_layout_shape)
        if layout_shape.is_placeholder and layout_shape.placeholder_format.idx == placeholder_index:
            return cast(
                object, layout_shape.text_frame._element.find(qn("a:lstStyle"))
            )
    return None


def _list_level_properties(list_style: object | None, level: int) -> object | None:
    if list_style is None:
        return None
    return cast("_XmlSearchElement", list_style).find(qn(f"a:lvl{level + 1}pPr"))


def _effective_bullet(style_properties: tuple[object, ...]) -> tuple[str | None, str | None]:
    bullet_kind: str | None = None
    bullet_marker: str | None = None
    for properties in style_properties:
        candidate_kind, candidate_marker = _paragraph_bullet(properties)
        if candidate_kind is not None:
            bullet_kind, bullet_marker = candidate_kind, candidate_marker
    return bullet_kind, bullet_marker


def _inherited_integer(style_properties: tuple[object, ...], attribute_name: str) -> int | None:
    value: int | None = None
    for properties in style_properties:
        candidate = _xml_integer(properties, attribute_name)
        if candidate is not None:
            value = candidate
    return value


def _inherited_alignment(style_properties: tuple[object, ...]) -> str | None:
    value: str | None = None
    for properties in style_properties:
        raw_value = cast("_XmlSearchElement", properties).get("algn")
        if raw_value is not None:
            value = {"l": "left", "ctr": "center", "r": "right", "just": "justify"}.get(
                raw_value, raw_value
            )
    return value


@dataclass(frozen=True, slots=True)
class _RunDefaults:
    font_family: str | None = None
    font_size_points: float | None = None
    bold: bool | None = None
    italic: bool | None = None
    underline: str | None = None
    baseline: int | None = None


def _run_defaults(style_properties: tuple[object, ...]) -> _RunDefaults:
    defaults = _RunDefaults()
    for properties in style_properties:
        default_run_properties = cast("_XmlSearchElement", properties).find(qn("a:defRPr"))
        if default_run_properties is not None:
            defaults = _merge_run_defaults(defaults, default_run_properties)
    return defaults


def _merge_run_defaults(defaults: _RunDefaults, properties: object) -> _RunDefaults:
    element = cast("_XmlSearchElement", properties)
    font_size = element.get("sz")
    bold = _xml_boolean(element.get("b"))
    italic = _xml_boolean(element.get("i"))
    baseline = element.get("baseline")
    underline = element.get("u")
    return _RunDefaults(
        font_family=_font_family_from_properties(properties) or defaults.font_family,
        font_size_points=(float(font_size) / 100.0) if font_size is not None else defaults.font_size_points,
        bold=bold if bold is not None else defaults.bold,
        italic=italic if italic is not None else defaults.italic,
        underline=underline if underline is not None else defaults.underline,
        baseline=int(baseline) if baseline is not None else defaults.baseline,
    )


def _font_family_from_properties(properties: object) -> str | None:
    element = cast("_XmlSearchElement", properties)
    for tag_name in ("latin", "ea"):
        family_element = element.find(qn(f"a:{tag_name}"))
        if family_element is not None:
            family = family_element.get("typeface")
            if family:
                return family
    return None


def _xml_boolean(value: str | None) -> bool | None:
    if value is None:
        return None
    return value not in {"0", "false", "False"}


def _effective_run_properties(
    source_properties: TextRunProperties, defaults: _RunDefaults
) -> TextRunProperties:
    font_family = source_properties.font_family or defaults.font_family
    return replace(
        source_properties,
        font_family=font_family,
        font_classification=_font_classification(font_family),
        font_size_points=source_properties.font_size_points or defaults.font_size_points,
        bold=source_properties.bold if source_properties.bold is not None else defaults.bold,
        italic=source_properties.italic if source_properties.italic is not None else defaults.italic,
        underline=source_properties.underline
        if source_properties.underline is not None
        else defaults.underline,
        baseline=source_properties.baseline
        if source_properties.baseline is not None
        else defaults.baseline,
    )


def _effective_text_rotation(shape_rotation: float, text_direction: str | None) -> float:
    direction_rotation = (
        {"vert": 90.0, "vert270": 270.0}.get(text_direction, 0.0)
        if text_direction is not None
        else 0.0
    )
    return (shape_rotation + direction_rotation) % 360.0


def _autofit_mode(text_frame: TextFrame) -> str:
    body_properties = text_frame._element.bodyPr
    if body_properties.find(f"{{{_DRAWING_NAMESPACE}}}normAutofit") is not None:
        return "text-to-fit-shape"
    if body_properties.find(f"{{{_DRAWING_NAMESPACE}}}spAutoFit") is not None:
        return "shape-to-fit-text"
    if body_properties.find(f"{{{_DRAWING_NAMESPACE}}}noAutofit") is not None:
        return "none"
    if text_frame.auto_size is None:
        return "inherited-or-unspecified"
    if text_frame.auto_size == MSO_AUTO_SIZE.NONE:
        return "none"
    return _enum_name(text_frame.auto_size) or "unknown"


def _autofit_integer(body_properties: object, attribute_name: str) -> int | None:
    """Read optional ``a:normAutofit`` integers from the XML boundary."""
    # python-pptx does not expose normAutofit scale attributes through its public API.
    element = cast("_XmlSearchElement", body_properties)
    autofit = element.find(f"{{{_DRAWING_NAMESPACE}}}normAutofit")
    if autofit is None:
        return None
    value = autofit.get(attribute_name)
    return int(value) if value is not None else None


class _XmlSearchElement:
    """Narrow protocol for the private python-pptx XML element boundary."""

    def find(self, path: str) -> _XmlSearchElement | None: ...

    def get(self, key: str) -> str | None: ...


class _SlideLayout:
    """Narrow protocol for the python-pptx layout inheritance boundary."""

    shapes: Iterable[BaseShape]
    slide_master: "_SlideMaster"


class _SlideMaster:
    """Narrow protocol for the master XML inheritance boundary."""

    _element: object


def _paragraph_properties(paragraph: _Paragraph) -> ParagraphProperties:
    line_spacing = paragraph.line_spacing
    if isinstance(line_spacing, Length):
        line_spacing_value = float(line_spacing.pt)
        line_spacing_kind = "points"
    elif isinstance(line_spacing, (int, float)):
        line_spacing_value = float(line_spacing)
        line_spacing_kind = "multiple"
    else:
        line_spacing_value = None
        line_spacing_kind = None
    paragraph_properties = paragraph._p.pPr
    bullet_kind, bullet_marker = _paragraph_bullet(paragraph_properties)
    return ParagraphProperties(
        alignment=_enum_name(paragraph.alignment),
        space_before_points=_length_points(paragraph.space_before),
        space_after_points=_length_points(paragraph.space_after),
        line_spacing=line_spacing_value,
        line_spacing_kind=line_spacing_kind,
        level=paragraph.level,
        margin_left_emu=_xml_integer(paragraph_properties, "marL"),
        indent_emu=_xml_integer(paragraph_properties, "indent"),
        bullet_kind=bullet_kind,
        bullet_marker=bullet_marker,
        empty_line_font_size_points=_end_paragraph_font_size_points(paragraph),
        runs=tuple(_run_properties(run) for run in paragraph.runs),
    )


def _run_properties(run: _Run) -> TextRunProperties:
    font = run.font
    font_name = font.name
    return TextRunProperties(
        text=run.text,
        font_family=font_name,
        font_classification=_font_classification(font_name),
        font_size_points=_length_points(font.size),
        bold=font.bold,
        italic=font.italic,
        underline=_enum_name(font.underline),
        baseline=_font_baseline(font._element),
    )


def _font_baseline(font_properties: object) -> int | None:
    """Read an optional direct DrawingML baseline value from a run."""
    # python-pptx has no public baseline API for DrawingML character properties.
    baseline = cast("_XmlSearchElement", font_properties).get("baseline")
    return int(baseline) if baseline is not None else None


def _paragraph_bullet(paragraph_properties: object | None) -> tuple[str | None, str | None]:
    """Return one directly specified DrawingML bullet marker, when available."""
    if paragraph_properties is None:
        return None, None
    element = cast("_XmlSearchElement", paragraph_properties)
    if element.find(f"{{{_DRAWING_NAMESPACE}}}buNone") is not None:
        return "none", None
    character_bullet = element.find(f"{{{_DRAWING_NAMESPACE}}}buChar")
    if character_bullet is not None:
        return "character", character_bullet.get("char")
    if element.find(f"{{{_DRAWING_NAMESPACE}}}buAutoNum") is not None:
        return "automatic-number", None
    if element.find(f"{{{_DRAWING_NAMESPACE}}}buBlip") is not None:
        return "picture", None
    return None, None


def _xml_integer(element: object | None, attribute_name: str) -> int | None:
    if element is None:
        return None
    value = cast("_XmlSearchElement", element).get(attribute_name)
    return int(value) if value is not None else None


def _end_paragraph_font_size_points(paragraph: _Paragraph) -> float | None:
    """Read an empty paragraph's direct end-paragraph run font size."""
    raw_size = _xml_integer(paragraph._p.endParaRPr, "sz")
    return None if raw_size is None else raw_size / 100.0


def _length_points(length: Length | None) -> float | None:
    return None if length is None else float(length.pt)


def _enum_name(value: object) -> str | None:
    if value is None:
        return None
    return str(value).split(" ", 1)[0].lower().replace("_", "-")


def _font_classification(font_name: str | None) -> str:
    if font_name is None:
        return "sans-serif"
    normalized = font_name.lower()
    if any(marker in normalized for marker in ("mono", "courier", "console", "code")):
        return "fixed-width"
    if any(marker in normalized for marker in ("serif", "roman", "times", "georgia")):
        return "serif"
    return "sans-serif"


def _render_text_boxes(
    text_box_evaluations: tuple[_TextBoxEvaluation, ...],
    typefaces: dict[str, skia.Typeface],
    output_path: Path,
    providers: dict[str, TextReplacementProvider],
    source_language: str,
    target_language: str,
) -> tuple[_TextBoxArtifact, ...]:
    artifact_directory = output_path.with_name(f"{output_path.stem}.text-layout-artifacts")
    artifacts: list[_TextBoxArtifact] = []
    for text_box_index, text_box_evaluation in enumerate(text_box_evaluations, 1):
        source_properties = text_box_evaluation.source_properties
        artifact_stem = artifact_directory / f"text-box-{text_box_index:04d}"
        rendering_path = artifact_stem.with_suffix(".png")
        properties_path = artifact_stem.with_suffix(".json")
        explicit_properties_path = artifact_stem.with_suffix(".explicit.json")
        rendering_path.parent.mkdir(parents=True, exist_ok=True)
        source_preview = _source_preview_fitting(text_box_evaluation.effective_properties, typefaces)
        rendering, layout_fit = _render_text_box(
            text_box_evaluation.effective_properties,
            typefaces,
            layout_fit=source_preview.layout_fit,
            canvas_height_emu=source_preview.canvas_height_emu,
            source_shape_guide=(
                _bounded_text_box(text_box_evaluation.effective_properties)
                if source_preview.derived_from_source
                else None
            ),
            derived_fit_box=(
                source_preview.fitting_box if source_preview.derived_from_source else None
            ),
        )
        rendering.save(rendering_path, format="PNG")
        properties = asdict(source_properties)
        properties["rendering"] = {
            "font_scale": layout_fit.font_scale,
            "fit_status": layout_fit.fit_status,
        }
        properties_path.write_text(
            json.dumps(properties, ensure_ascii=False, indent=2) + "\n", encoding="utf-8"
        )
        explicit_properties = asdict(
            _explicit_text_box_properties(
                text_box_evaluation.effective_properties, layout_fit, typefaces
            )
        )
        explicit_properties["rendering"] = {
            "font_scale": layout_fit.font_scale,
            "fit_status": layout_fit.fit_status,
        }
        explicit_properties_path.write_text(
            json.dumps(explicit_properties, ensure_ascii=False, indent=2) + "\n",
            encoding="utf-8",
        )
        replacement_artifacts: list[_ProviderTextBoxArtifact] = []
        for provider_index, (provider_name, provider) in enumerate(providers.items(), 1):
            replacement_text_box = _replace_text_box_paragraphs(
                text_box_evaluation.effective_properties,
                provider,
                source_language,
                target_language,
            )
            replacement_fitting = _replacement_fitting(
                text_box_evaluation.effective_properties, replacement_text_box, typefaces
            )
            replacement_rendering, replacement_layout_fit = _render_text_box(
                replacement_text_box,
                typefaces,
                layout_fit=replacement_fitting.layout_fit,
                derived_fit_box=(
                    replacement_fitting.fitting_box
                    if replacement_fitting.derived_from_source
                    else None
                ),
                canvas_height_emu=replacement_fitting.canvas_height_emu,
                source_shape_guide=(
                    _bounded_text_box(text_box_evaluation.effective_properties)
                    if replacement_fitting.derived_from_source
                    else None
                ),
            )
            provider_filename = f"{artifact_stem.name}.provider-{provider_index:04d}"
            replacement_rendering_path = artifact_stem.parent / f"{provider_filename}.png"
            replacement_explicit_path = artifact_stem.parent / f"{provider_filename}.explicit.json"
            replacement_rendering.save(replacement_rendering_path, format="PNG")
            replacement_explicit_properties = asdict(
                _explicit_text_box_properties(
                    replacement_text_box, replacement_layout_fit, typefaces
                )
            )
            replacement_explicit_properties["rendering"] = {
                "font_scale": replacement_layout_fit.font_scale,
                "fit_status": replacement_layout_fit.fit_status,
            }
            replacement_explicit_properties["replacement"] = {
                "provider": provider_name,
                "source_language": source_language,
                "target_language": target_language,
                "unit": "paragraph",
                "style_strategy": "dominant-source-run",
            }
            replacement_explicit_properties["fitting"] = {
                "derived_from_source": replacement_fitting.derived_from_source,
                "rectangle": _fitting_rectangle(replacement_fitting.fitting_box),
            }
            replacement_explicit_path.write_text(
                json.dumps(replacement_explicit_properties, ensure_ascii=False, indent=2) + "\n",
                encoding="utf-8",
            )
            replacement_artifacts.append(
                _ProviderTextBoxArtifact(
                    provider_name, replacement_rendering_path, replacement_explicit_path
                )
            )
        artifacts.append(
            _TextBoxArtifact(
                rendering_path,
                properties_path,
                explicit_properties_path,
                tuple(replacement_artifacts),
            )
        )
    return tuple(artifacts)


def _replacement_fitting(
    source_text_box: TextBoxProperties,
    replacement_text_box: TextBoxProperties,
    typefaces: dict[str, skia.Typeface],
) -> _ReplacementFitting:
    """Use the PPTX adapter's bound selection before rendering a replacement."""
    source_bounded = _bounded_text_box(source_text_box)
    derived_from_source = source_text_box.explicit_no_autofit
    fitting_box = (
        source_occupied_text_box(source_bounded, typefaces)
        if derived_from_source
        else source_bounded
    )
    replacement_bounded = replace(
        _bounded_text_box(replacement_text_box),
        width_emu=fitting_box.width_emu,
        height_emu=fitting_box.height_emu,
    )
    fitted = fit_explicit_noto_text_box(replacement_bounded, typefaces)
    full_width, _ = _content_dimensions(replacement_text_box)
    lines = _layout_lines(
        replacement_text_box.paragraphs, full_width, typefaces, fitted.font_scale
    )
    return _ReplacementFitting(
        _LayoutFit(lines, fitted.font_scale, fitted.fit_status),
        fitting_box,
        derived_from_source,
        source_text_box.height_emu,
    )


def _source_preview_fitting(
    text_box: TextBoxProperties, typefaces: dict[str, skia.Typeface]
) -> _ReplacementFitting:
    """Render no-autofit source text at its source scale and natural height."""
    source_bounded = _bounded_text_box(text_box)
    if not text_box.explicit_no_autofit:
        width, height = _content_dimensions(text_box)
        return _ReplacementFitting(
            _fit_layout(text_box.paragraphs, width, height, typefaces),
            source_bounded,
            False,
            text_box.height_emu,
        )
    fitting_box = source_occupied_text_box(source_bounded, typefaces)
    full_width, _ = _content_dimensions(text_box)
    return _ReplacementFitting(
        _LayoutFit(
            _layout_lines(text_box.paragraphs, full_width, typefaces),
            1.0,
            "source-no-autofit",
        ),
        fitting_box,
        True,
        text_box.height_emu,
    )


def _bounded_text_box(text_box: TextBoxProperties) -> BoundedTextBox:
    return BoundedTextBox(
        width_emu=text_box.width_emu,
        height_emu=text_box.height_emu,
        margin_left_emu=text_box.margin_left_emu,
        margin_top_emu=text_box.margin_top_emu,
        margin_right_emu=text_box.margin_right_emu,
        margin_bottom_emu=text_box.margin_bottom_emu,
        text_direction=text_box.text_direction,
        paragraphs=tuple(
            BoundedTextParagraph(
                alignment=paragraph.alignment,
                space_before_points=paragraph.space_before_points,
                space_after_points=paragraph.space_after_points,
                line_spacing=paragraph.line_spacing,
                line_spacing_kind=paragraph.line_spacing_kind,
                level=paragraph.level,
                margin_left_emu=paragraph.margin_left_emu,
                indent_emu=paragraph.indent_emu,
                bullet_kind=paragraph.bullet_kind,
                bullet_marker=paragraph.bullet_marker,
                empty_line_font_size_points=paragraph.empty_line_font_size_points,
                runs=tuple(
                    BoundedTextRun(
                        text=run.text,
                        font_family=run.font_family,
                        font_classification=run.font_classification,
                        font_size_points=run.font_size_points,
                        bold=run.bold,
                        italic=run.italic,
                        underline=run.underline,
                        baseline=run.baseline,
                    )
                    for run in paragraph.runs
                ),
            )
            for paragraph in text_box.paragraphs
        ),
    )


def _fitting_rectangle(text_box: BoundedTextBox) -> dict[str, int]:
    return {
        "left_emu": 0,
        "top_emu": 0,
        "width_emu": text_box.width_emu,
        "height_emu": text_box.height_emu,
    }


def _replace_text_box_paragraphs(
    text_box: TextBoxProperties,
    provider: TextReplacementProvider,
    source_language: str,
    target_language: str,
) -> TextBoxProperties:
    """Replace each populated paragraph using its most common source run style."""
    paragraphs = tuple(
        _replace_paragraph(paragraph, provider, source_language, target_language)
        for paragraph in text_box.paragraphs
    )
    return replace(text_box, paragraphs=paragraphs)


def _replace_paragraph(
    paragraph: ParagraphProperties,
    provider: TextReplacementProvider,
    source_language: str,
    target_language: str,
) -> ParagraphProperties:
    source_text = "".join(run.text for run in paragraph.runs)
    if not source_text.strip():
        return paragraph
    replacement = provider.replace(
        TextReplacementRequest(source_text, False, source_language, target_language)
    )
    dominant_run = _dominant_run(paragraph.runs)
    return replace(paragraph, runs=(replace(dominant_run, text=replacement.text),))


def _dominant_run(runs: tuple[TextRunProperties, ...]) -> TextRunProperties:
    """Return the first run with the greatest non-whitespace source character count."""
    return max(runs, key=lambda run: sum(not character.isspace() for character in run.text))


def _explicit_text_box_properties(
    text_box: TextBoxProperties,
    layout_fit: _LayoutFit,
    typefaces: dict[str, skia.Typeface],
) -> TextBoxProperties:
    """Return evaluator-writable properties with selected fonts and fitted sizes."""
    paragraphs = tuple(
        replace(
            paragraph,
            runs=tuple(
                replace(
                    run,
                    font_family=typefaces[run.font_classification].getFamilyName(),
                    font_size_points=(run.font_size_points or DEFAULT_FONT_SIZE_POINTS)
                    * layout_fit.font_scale,
                )
                for run in paragraph.runs
            ),
            empty_line_font_size_points=(
                paragraph.empty_line_font_size_points or DEFAULT_FONT_SIZE_POINTS
            )
            * layout_fit.font_scale,
        )
        for paragraph in text_box.paragraphs
    )
    return replace(
        text_box,
        autofit_mode="none",
        autofit_font_scale=None,
        autofit_line_spacing_reduction=None,
        paragraphs=paragraphs,
    )


def _render_text_box(
    text_box: TextBoxProperties,
    typefaces: dict[str, skia.Typeface],
    *,
    layout_fit: _LayoutFit | None = None,
    derived_fit_box: BoundedTextBox | None = None,
    canvas_height_emu: int | None = None,
    source_shape_guide: BoundedTextBox | None = None,
) -> tuple[Image.Image, _LayoutFit]:
    width = max(1, ceil(text_box.width_emu / EMU_PER_PIXEL))
    height = max(1, ceil((canvas_height_emu or text_box.height_emu) / EMU_PER_PIXEL))
    surface = skia.Surface(width, height)
    if surface is None:
        raise RuntimeError(f"Skia could not create a {width}x{height} native-text surface.")
    canvas = surface.getCanvas()
    canvas.clear(skia.ColorWHITE)
    content_left = text_box.margin_left_emu / EMU_PER_PIXEL
    content_top = text_box.margin_top_emu / EMU_PER_PIXEL
    content_right = width - text_box.margin_right_emu / EMU_PER_PIXEL
    content_bottom = height - text_box.margin_bottom_emu / EMU_PER_PIXEL
    if layout_fit is None:
        layout_fit = _fit_layout(
            text_box.paragraphs,
            max(0.0, content_right - content_left),
            max(0.0, content_bottom - content_top),
            typefaces,
        )
    canvas.save()
    canvas.clipRect(skia.Rect.MakeLTRB(content_left, content_top, content_right, content_bottom))
    canvas.rotate(text_box.effective_text_rotation_degrees, width / 2.0, height / 2.0)
    lines = layout_fit.lines
    total_height = _layout_height(lines)
    baseline = _vertical_start(text_box.vertical_alignment, content_top, content_bottom, total_height)
    for line in lines:
        baseline += _space_before_pixels(line.paragraph)
        line_baseline = baseline + line.height * 0.8
        _draw_line(canvas, line, content_left, content_right, line_baseline, typefaces)
        baseline += _line_advance(line)
        baseline += _space_after_pixels(line.paragraph)
    canvas.restore()
    if source_shape_guide is not None:
        _draw_source_shape_guide(canvas, source_shape_guide)
    if derived_fit_box is not None:
        _draw_derived_fit_guide(canvas, derived_fit_box)
    image = surface.makeImageSnapshot()
    data = image.encodeToData()
    if data is None:
        raise RuntimeError("Skia could not encode a native-text layout bitmap.")
    return Image.open(BytesIO(bytes(data))).copy(), layout_fit


def _content_dimensions(text_box: TextBoxProperties) -> tuple[float, float]:
    return (
        max(
            0.0,
            (text_box.width_emu - text_box.margin_left_emu - text_box.margin_right_emu)
            / EMU_PER_PIXEL,
        ),
        max(
            0.0,
            (text_box.height_emu - text_box.margin_top_emu - text_box.margin_bottom_emu)
            / EMU_PER_PIXEL,
        ),
    )


def _draw_derived_fit_guide(canvas: skia.Canvas, fit_box: BoundedTextBox) -> None:
    """Draw the local fit input without changing the replacement preview itself."""
    paint = skia.Paint(
        Color=_DERIVED_FIT_GUIDE_COLOR,
        Style=skia.Paint.kStroke_Style,
        StrokeWidth=1.0,
        AntiAlias=False,
    )
    paint.setPathEffect(skia.DashPathEffect.Make([4.0, 3.0], 0.0))
    canvas.drawRect(
        skia.Rect.MakeLTRB(
            0.5,
            0.5,
            max(0.5, fit_box.width_emu / EMU_PER_PIXEL - 0.5),
            max(0.5, fit_box.height_emu / EMU_PER_PIXEL - 0.5),
        ),
        paint,
    )


def _draw_source_shape_guide(canvas: skia.Canvas, source_box: BoundedTextBox) -> None:
    """Show the original shape boundary when an evaluation canvas is taller."""
    paint = skia.Paint(
        Color=skia.ColorRED,
        Style=skia.Paint.kStroke_Style,
        StrokeWidth=1.0,
        AntiAlias=False,
    )
    canvas.drawRect(
        skia.Rect.MakeLTRB(
            0.5,
            0.5,
            max(0.5, source_box.width_emu / EMU_PER_PIXEL - 0.5),
            max(0.5, source_box.height_emu / EMU_PER_PIXEL - 0.5),
        ),
        paint,
    )


def _fit_layout(
    paragraphs: tuple[ParagraphProperties, ...],
    width: float,
    height: float,
    typefaces: dict[str, skia.Typeface],
) -> _LayoutFit:
    """Find the largest uniform scale that fits while testing down to one pixel."""
    full_size_lines = _layout_lines(paragraphs, width, typefaces)
    if _layout_fits(full_size_lines, width, height):
        return _LayoutFit(full_size_lines, 1.0, "fit")

    maximum_font_size = _maximum_font_size_pixels(paragraphs)
    minimum_scale = 1.0 / maximum_font_size
    minimum_lines = _layout_lines(paragraphs, width, typefaces, minimum_scale)
    if not _layout_fits(minimum_lines, width, height):
        return _LayoutFit(minimum_lines, minimum_scale, "overflow")

    fitting_scale = minimum_scale
    fitting_lines = minimum_lines
    non_fitting_scale = 1.0
    for _ in range(16):
        candidate_scale = (fitting_scale + non_fitting_scale) / 2.0
        candidate_lines = _layout_lines(paragraphs, width, typefaces, candidate_scale)
        if _layout_fits(candidate_lines, width, height):
            fitting_scale = candidate_scale
            fitting_lines = candidate_lines
        else:
            non_fitting_scale = candidate_scale
    return _LayoutFit(fitting_lines, fitting_scale, "fit")


def _maximum_font_size_pixels(paragraphs: tuple[ParagraphProperties, ...]) -> float:
    sizes = [
        (run.font_size_points or DEFAULT_FONT_SIZE_POINTS) * PIXELS_PER_POINT
        for paragraph in paragraphs
        for run in paragraph.runs
    ]
    sizes.extend(
        (paragraph.empty_line_font_size_points or DEFAULT_FONT_SIZE_POINTS)
        * PIXELS_PER_POINT
        for paragraph in paragraphs
        if not paragraph.runs
    )
    return max(sizes, default=DEFAULT_FONT_SIZE_POINTS * PIXELS_PER_POINT)


def _layout_fits(lines: tuple[_DrawLine, ...], width: float, height: float) -> bool:
    if _layout_height(lines) > height:
        return False
    return all(
        line.width <= max(0.0, width - _paragraph_margin_pixels(line.paragraph))
        for line in lines
    )


def _layout_lines(
    paragraphs: tuple[ParagraphProperties, ...],
    width: float,
    typefaces: dict[str, skia.Typeface],
    font_scale: float = 1.0,
) -> tuple[_DrawLine, ...]:
    lines: list[_DrawLine] = []
    for paragraph in paragraphs:
        current_segments: list[_DrawSegment] = []
        current_width = 0.0
        paragraph_width = max(0.0, width - _paragraph_margin_pixels(paragraph))
        is_first_line = True
        for run in paragraph.runs:
            style = _draw_style(run, font_scale)
            for token in _layout_tokens(run.text):
                if token in {"\n", "\v"}:
                    lines.append(
                        _draw_line_info(
                            current_segments, current_width, paragraph, is_first_line, font_scale
                        )
                    )
                    current_segments = []
                    current_width = 0.0
                    is_first_line = False
                    continue
                for layout_token in _emergency_wrap_tokens(
                    token, style, paragraph_width, typefaces
                ):
                    token_width = _measure(layout_token, style, typefaces)
                    if (
                        current_segments
                        and paragraph_width > 0
                        and current_width + token_width > paragraph_width
                    ):
                        lines.append(
                            _draw_line_info(
                                current_segments, current_width, paragraph, is_first_line, font_scale
                            )
                        )
                        current_segments = []
                        current_width = 0.0
                        is_first_line = False
                    current_segments.append(_DrawSegment(layout_token, style))
                    current_width += token_width
        lines.append(
            _draw_line_info(current_segments, current_width, paragraph, is_first_line, font_scale)
        )
    return tuple(lines)


def _layout_tokens(text: str) -> Iterable[str]:
    """Yield word tokens, allowing East Asian text to wrap between characters."""
    for token in _TOKEN_PATTERN.findall(text):
        if token in {"\n", "\v"} or not _contains_wide_character(token):
            yield token
        else:
            yield from token


def _contains_wide_character(text: str) -> bool:
    return any(unicodedata.east_asian_width(character) in {"W", "F"} for character in text)


def _emergency_wrap_tokens(
    token: str,
    style: _DrawStyle,
    available_width: float,
    typefaces: dict[str, skia.Typeface],
) -> tuple[str, ...]:
    """Split an over-wide unbreakable token only when ordinary wrapping cannot fit it."""
    if available_width <= 0 or _measure(token, style, typefaces) <= available_width:
        return (token,)
    return tuple(token)


def _draw_style(run: TextRunProperties, font_scale: float = 1.0) -> _DrawStyle:
    return _DrawStyle(
        classification=run.font_classification,
        size_pixels=(run.font_size_points or DEFAULT_FONT_SIZE_POINTS)
        * PIXELS_PER_POINT
        * font_scale,
        bold=run.bold is True,
        italic=run.italic is True,
        underline=None if run.underline in (None, "none", "false") else run.underline,
        baseline=run.baseline or 0,
    )


def _draw_line_info(
    segments: list[_DrawSegment],
    width: float,
    paragraph: ParagraphProperties,
    is_first_line: bool,
    font_scale: float,
) -> _DrawLine:
    default_height = (
        paragraph.empty_line_font_size_points or DEFAULT_FONT_SIZE_POINTS
    ) * PIXELS_PER_POINT * 1.2 * font_scale
    height = max((segment.style.size_pixels * 1.2 for segment in segments), default=default_height)
    return _DrawLine(tuple(segments), width, height, paragraph, is_first_line)


def _measure(text: str, style: _DrawStyle, typefaces: dict[str, skia.Typeface]) -> float:
    return float(_font(style, typefaces).measureText(text))


def _font(style: _DrawStyle, typefaces: dict[str, skia.Typeface]) -> skia.Font:
    font = skia.Font(typefaces[style.classification], style.size_pixels)
    font.setEmbolden(style.bold)
    if style.italic:
        font.setSkewX(-0.2)
    return font


def _layout_height(lines: tuple[_DrawLine, ...]) -> float:
    return sum(
        _space_before_pixels(line.paragraph)
        + _line_advance(line)
        + _space_after_pixels(line.paragraph)
        for line in lines
    )


def _line_advance(line: _DrawLine) -> float:
    if line.paragraph.line_spacing_kind == "multiple" and line.paragraph.line_spacing is not None:
        return line.height * line.paragraph.line_spacing
    if line.paragraph.line_spacing_kind == "points" and line.paragraph.line_spacing is not None:
        return max(line.height, line.paragraph.line_spacing * PIXELS_PER_POINT)
    return line.height


def _space_before_pixels(paragraph: ParagraphProperties) -> float:
    return (paragraph.space_before_points or 0.0) * PIXELS_PER_POINT


def _space_after_pixels(paragraph: ParagraphProperties) -> float:
    return (paragraph.space_after_points or 0.0) * PIXELS_PER_POINT


def _paragraph_margin_pixels(paragraph: ParagraphProperties) -> float:
    return (paragraph.margin_left_emu or 0) / EMU_PER_PIXEL


def _vertical_start(alignment: str | None, top: float, bottom: float, total_height: float) -> float:
    if alignment == "middle":
        return top + max(0.0, (bottom - top - total_height) / 2.0)
    if alignment == "bottom":
        return max(top, bottom - total_height)
    return top


def _draw_line(
    canvas: skia.Canvas,
    line: _DrawLine,
    content_left: float,
    content_right: float,
    baseline: float,
    typefaces: dict[str, skia.Typeface],
) -> None:
    paragraph_left = content_left + _paragraph_margin_pixels(line.paragraph)
    x = _horizontal_start(line.paragraph.alignment, paragraph_left, content_right, line.width)
    paint = skia.Paint(Color=skia.ColorBLACK, AntiAlias=True)
    if line.is_first_line and line.segments and line.paragraph.bullet_kind == "character":
        bullet_style = line.segments[0].style if line.segments else _DrawStyle(
            "sans-serif", DEFAULT_FONT_SIZE_POINTS * PIXELS_PER_POINT, False, False, None, 0
        )
        bullet_x = paragraph_left + (line.paragraph.indent_emu or 0) / EMU_PER_PIXEL
        canvas.drawString(line.paragraph.bullet_marker or "", bullet_x, baseline, _font(bullet_style, typefaces), paint)
    for segment in line.segments:
        font = _font(segment.style, typefaces)
        baseline_offset = -segment.style.size_pixels * segment.style.baseline / 100_000.0
        y = baseline + baseline_offset
        canvas.drawString(segment.text, x, y, font, paint)
        segment_width = font.measureText(segment.text)
        if segment.style.underline not in (None, "none"):
            _draw_underline(
                canvas, x, y, segment_width, font, paint, segment.style.underline
            )
        x += segment_width


def _draw_underline(
    canvas: skia.Canvas,
    x: float,
    baseline: float,
    width: float,
    font: skia.Font,
    paint: skia.Paint,
    underline_kind: str,
) -> None:
    """Draw one source underline using the selected typeface metrics."""
    metrics = font.getMetrics()
    underline_y = baseline + float(metrics.fUnderlinePosition)
    underline_paint = skia.Paint(paint)
    underline_paint.setStyle(skia.Paint.kStroke_Style)
    underline_paint.setStrokeWidth(max(1.0, float(metrics.fUnderlineThickness)))
    canvas.drawLine(x, underline_y, x + width, underline_y, underline_paint)
    if underline_kind in {"double-line", "wavy-double-line"}:
        second_y = underline_y + max(2.0, float(metrics.fUnderlineThickness) * 2.0)
        canvas.drawLine(x, second_y, x + width, second_y, underline_paint)


def _horizontal_start(alignment: str | None, left: float, right: float, line_width: float) -> float:
    if alignment == "center":
        return left + max(0.0, (right - left - line_width) / 2.0)
    if alignment == "right":
        return max(left, right - line_width)
    return left


def _write_html_page(
    output_path: Path,
    source_path: Path,
    text_boxes: tuple[TextBoxProperties, ...],
    artifacts: tuple[_TextBoxArtifact, ...],
    provider_names: tuple[str, ...],
    source_language: str,
    target_language: str,
) -> None:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    rows = "\n".join(
        _html_row(index, text_box, artifact, output_path.parent)
        for index, (text_box, artifact) in enumerate(zip(text_boxes, artifacts, strict=True), 1)
    )
    if not rows:
        rows = f'<tr><td colspan="{2 + len(provider_names)}">No eligible text boxes.</td></tr>'
    title = html.escape(source_path.as_posix())
    provider_headers = "".join(f"<th>{html.escape(name)}</th>" for name in provider_names)
    output_path.write_text(
        f"""<!doctype html>
<html lang="en">
<head>
  <meta charset="utf-8">
  <meta name="viewport" content="width=device-width, initial-scale=1">
  <title>Native text-layout evaluation: {title}</title>
  <style>
    body {{ background: #f6f7f9; color: #1f2937; font-family: system-ui, sans-serif; margin: 0; overflow-x: auto; }}
    main {{ padding: 2rem; width: max-content; }}
    table {{ background: white; border-collapse: collapse; width: max-content; }}
    th, td {{ border: 1px solid #cbd5e1; padding: .55rem; text-align: left; vertical-align: top; }}
    th {{ background: #e9eef5; white-space: nowrap; }}
    .rendering img {{ background: white; border: 1px solid #dc2626; display: block; height: auto; }}
  </style>
</head>
<body>
  <main>
    <h1>Native text-layout evaluation</h1>
    <p>PowerPoint source: <code>{title}</code> <span>{html.escape(source_language)}→{html.escape(target_language)}</span></p>
    <table>
      <thead><tr><th>Text box</th><th>Original</th>{provider_headers}</tr></thead>
      <tbody>{rows}</tbody>
    </table>
  </main>
</body>
</html>
""",
        encoding="utf-8",
    )


def _html_row(
    index: int, text_box: TextBoxProperties, artifact: _TextBoxArtifact, page_directory: Path
) -> str:
    image_url = Path(os.path.relpath(artifact.rendering_path, page_directory)).as_posix()
    properties_url = Path(os.path.relpath(artifact.properties_path, page_directory)).as_posix()
    explicit_properties_url = Path(
        os.path.relpath(artifact.explicit_properties_path, page_directory)
    ).as_posix()
    replacement_cells = "".join(
        _provider_html_cell(provider_artifact, page_directory)
        for provider_artifact in artifact.replacement_artifacts
    )
    return (
        f"<tr><td>{index}<br><code>{html.escape(text_box.source)}</code>"
        f'<br><a href="{html.escape(properties_url, quote=True)}" target="_blank" '
        f'rel="noopener">Properties JSON</a>'
        f'<br><a href="{html.escape(explicit_properties_url, quote=True)}" target="_blank" '
        f'rel="noopener">Explicit properties JSON</a></td>'
        f'<td class="rendering"><img alt="Skia layout for text box {index}" '
        f'src="{html.escape(image_url, quote=True)}"></td>{replacement_cells}</tr>'
    )


def _provider_html_cell(artifact: _ProviderTextBoxArtifact, page_directory: Path) -> str:
    image_url = Path(os.path.relpath(artifact.rendering_path, page_directory)).as_posix()
    properties_url = Path(
        os.path.relpath(artifact.explicit_properties_path, page_directory)
    ).as_posix()
    return (
        f'<td class="rendering"><img alt="Skia layout replacement by '
        f'{html.escape(artifact.provider_name, quote=True)}" '
        f'src="{html.escape(image_url, quote=True)}">'
        f'<a href="{html.escape(properties_url, quote=True)}" target="_blank" '
        f'rel="noopener">Explicit properties JSON</a></td>'
    )


def parse_arguments() -> argparse.Namespace:
    """Return command-line arguments for the native-text evaluator."""
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--input-root", type=Path, default=DEFAULT_INPUT_ROOT)
    parser.add_argument("--output-root", type=Path, default=DEFAULT_OUTPUT_ROOT)
    parser.add_argument("--target-language", default=DEFAULT_TARGET_LANGUAGE)
    return parser.parse_args()


def main() -> int:
    """Run the native PowerPoint text-layout evaluator."""
    arguments = parse_arguments()
    result = evaluate_text_replacement_examples(
        arguments.input_root, arguments.output_root, arguments.target_language
    )
    print(
        "Native text-layout evaluation complete: "
        f"{result.written_pages} page(s) written, {result.rendered_text_boxes} text box(es) "
        f"rendered, {result.skipped_presentations} presentation(s) skipped."
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
