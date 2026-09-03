"""Synthetic tests for the native PowerPoint text-layout evaluator."""

from __future__ import annotations

from collections.abc import Iterable
import json
from pathlib import Path
from tempfile import TemporaryDirectory
from typing import cast
import unittest
from unittest.mock import MagicMock, call, patch

from PIL import Image, ImageChops
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.shapes.autoshape import Shape
from pptx.shapes.base import BaseShape
from pptx.util import Inches, Pt

from scripts.text_replacement_evaluations import (
    ParagraphProperties,
    TextRunProperties,
    _dominant_run,
    _draw_style,
    _line_advance,
    _load_typefaces,
    _layout_lines,
    evaluate_text_replacement_examples,
)


class NativeTextLayoutEvaluationTests(unittest.TestCase):
    def setUp(self) -> None:
        argos_translation = patch(
            "pipeline.text_replacement_plugins.argos_translate.ArgosTranslateProvider._translate",
            return_value="translated",
        )
        argos_translation.start()
        self.addCleanup(argos_translation.stop)

    # Verifies FR-2026-08-03-13.
    def test_uses_committed_noto_mono_for_fixed_width_text(self) -> None:
        typefaces = _load_typefaces()

        self.assertEqual("Noto Sans Mono", typefaces["fixed-width"].getFamilyName())

    # Verifies FR-2026-08-03-13.
    def test_explicit_false_underline_does_not_create_an_underlined_draw_style(self) -> None:
        draw_style = _draw_style(
            TextRunProperties(
                text="plain",
                font_family=None,
                font_classification="sans-serif",
                font_size_points=18.0,
                bold=False,
                italic=False,
                underline="false",
                baseline=None,
            )
        )

        self.assertIsNone(draw_style.underline)

    # Verifies FR-2026-08-03-13.
    def test_preserves_empty_paragraph_line_advance(self) -> None:
        paragraph = ParagraphProperties(
            alignment=None,
            space_before_points=None,
            space_after_points=None,
            line_spacing=None,
            line_spacing_kind=None,
            level=0,
            margin_left_emu=None,
            indent_emu=None,
            bullet_kind=None,
            bullet_marker=None,
            empty_line_font_size_points=24.0,
            runs=(),
        )

        lines = _layout_lines((paragraph,), 300.0, _load_typefaces())

        self.assertEqual(1, len(lines))
        self.assertEqual(24.0 * 4.0 / 3.0 * 1.2, lines[0].height)

    # Verifies FR-2026-08-03-13.
    def test_rewraps_wide_character_text_without_spaces(self) -> None:
        paragraph = ParagraphProperties(
            alignment=None,
            space_before_points=None,
            space_after_points=None,
            line_spacing=None,
            line_spacing_kind=None,
            level=0,
            margin_left_emu=None,
            indent_emu=None,
            bullet_kind=None,
            bullet_marker=None,
            empty_line_font_size_points=None,
            runs=(
                TextRunProperties(
                    text="日本語の文字列は空白なしでも改行可能です",
                    font_family=None,
                    font_classification="sans-serif",
                    font_size_points=18.0,
                    bold=None,
                    italic=None,
                    underline=None,
                    baseline=None,
                ),
            ),
        )

        lines = _layout_lines((paragraph,), 80.0, _load_typefaces())

        self.assertGreater(len(lines), 1)
        self.assertTrue(all(line.width <= 80.0 for line in lines))

    # Verifies FR-2026-08-03-13.
    def test_rewraps_an_overwide_unbreakable_token_before_reducing_font_size(self) -> None:
        paragraph = ParagraphProperties(
            alignment=None,
            space_before_points=None,
            space_after_points=None,
            line_spacing=None,
            line_spacing_kind=None,
            level=0,
            margin_left_emu=None,
            indent_emu=None,
            bullet_kind=None,
            bullet_marker=None,
            empty_line_font_size_points=None,
            runs=(
                TextRunProperties(
                    text="############",
                    font_family=None,
                    font_classification="sans-serif",
                    font_size_points=18.0,
                    bold=None,
                    italic=None,
                    underline=None,
                    baseline=None,
                ),
            ),
        )

        lines = _layout_lines((paragraph,), 30.0, _load_typefaces())

        self.assertGreater(len(lines), 1)
        self.assertTrue(all(line.width <= 30.0 for line in lines))

    # Verifies FR-2026-08-03-13.
    def test_point_line_spacing_does_not_reduce_advance_below_glyph_height(self) -> None:
        paragraph = ParagraphProperties(
            alignment=None,
            space_before_points=None,
            space_after_points=None,
            line_spacing=7.0,
            line_spacing_kind="points",
            level=0,
            margin_left_emu=None,
            indent_emu=None,
            bullet_kind=None,
            bullet_marker=None,
            empty_line_font_size_points=None,
            runs=(
                TextRunProperties(
                    text="line height",
                    font_family=None,
                    font_classification="sans-serif",
                    font_size_points=18.0,
                    bold=None,
                    italic=None,
                    underline=None,
                    baseline=None,
                ),
            ),
        )

        line = _layout_lines((paragraph,), 300.0, _load_typefaces())[0]

        self.assertEqual(line.height, _line_advance(line))

    # Verifies FR-2026-08-03-13.
    def test_dominant_run_ignores_whitespace_when_selecting_replacement_style(self) -> None:
        first_run = TextRunProperties(
            text="a b",
            font_family="First",
            font_classification="sans-serif",
            font_size_points=10.0,
            bold=None,
            italic=None,
            underline=None,
            baseline=None,
        )
        second_run = TextRunProperties(
            text="ccc",
            font_family="Second",
            font_classification="serif",
            font_size_points=20.0,
            bold=None,
            italic=None,
            underline=None,
            baseline=None,
        )

        self.assertIs(second_run, _dominant_run((first_run, second_run)))

    # Verifies FR-2026-08-04-14.
    def test_reports_progress_for_each_source_language_directory(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"
            _write_presentation(input_root / "ja" / "first.pptx")
            _write_presentation(input_root / "sample_data" / "en" / "second.pptx")
            (input_root / "ja" / "~$locked.pptx").write_bytes(b"not a presentation")
            progress = MagicMock()
            progress_context = MagicMock()
            progress_context.__enter__.return_value = progress_context
            progress_context.start_current.return_value = progress

            with patch(
                "scripts.text_replacement_evaluations.LiveProgress",
                return_value=progress_context,
            ) as mocked_progress:
                evaluate_text_replacement_examples(input_root, root / "output")

            mocked_progress.assert_called_once_with()
            progress_context.start_overall.assert_called_once_with(2, "presentation")
            self.assertEqual(
                [call("ja", 1, "presentation"), call("sample_data/en", 1, "presentation")],
                progress_context.start_current.call_args_list,
            )
            self.assertEqual(
                [call("first.pptx"), call("second.pptx")],
                progress.set_postfix_str.call_args_list,
            )
            self.assertEqual([call(1), call(1)], progress.update.call_args_list)
            self.assertEqual([call(), call()], progress_context.advance_overall.call_args_list)

    # Verifies FR-2026-08-03-13, FR-2026-08-04-12, and FR-2026-08-22-09.
    def test_processes_the_configured_confidential_subtree_locally(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"
            confidential_root = input_root / "confidential" / "ja"
            _write_presentation(input_root / "ja" / "regular" / "layout.pptx")
            _write_presentation(confidential_root / "restricted.pptx")

            result = evaluate_text_replacement_examples(input_root, root / "output")

            self.assertEqual(2, result.processed_presentations)
            self.assertEqual(4, result.written_pages)
            self.assertTrue((root / "output" / "confidential" / "ja" / "restricted.html").is_file())
            self.assertTrue((root / "output" / "confidential" / "ja" / "restricted.sf.html").is_file())

    # Verifies FR-2026-08-03-13 and FR-2026-08-22-09.
    def test_ignores_powerpoint_temporary_lock_files(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"
            _write_presentation(input_root / "ja" / "layout.pptx")
            (input_root / "ja" / "~$layout.pptx").write_bytes(b"not a presentation")

            result = evaluate_text_replacement_examples(input_root, root / "output")

            self.assertEqual(1, result.processed_presentations)
            self.assertEqual(2, result.written_pages)

    # Verifies FR-2026-08-03-13 and FR-2026-08-22-09.
    def test_skips_text_frames_without_non_whitespace_run_text(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            source_path = root / "input" / "ja" / "empty-bullet.pptx"
            _write_empty_bullet_presentation(source_path)

            result = evaluate_text_replacement_examples(root / "input", root / "output")

            self.assertEqual(0, result.rendered_text_boxes)
            self.assertFalse(
                (
                    root
                    / "output"
                    / "ja"
                    / "empty-bullet.text-layout-artifacts"
                    / "text-box-0001.png"
                ).exists()
            )

    # Verifies FR-2026-08-03-13.
    def test_reduces_the_uniform_font_scale_before_rendering_overflowing_text(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            source_path = root / "input" / "ja" / "overflowing.pptx"
            _write_overflowing_presentation(source_path)

            result = evaluate_text_replacement_examples(root / "input", root / "output")

            self.assertEqual(1, result.rendered_text_boxes)
            properties_path = (
                root
                / "output"
                / "ja"
                / "overflowing.text-layout-artifacts"
                / "text-box-0001.json"
            )
            properties = json.loads(properties_path.read_text(encoding="utf-8"))
            self.assertEqual("fit", properties["rendering"]["fit_status"])
            self.assertLess(properties["rendering"]["font_scale"], 1.0)

    # Verifies FR-2026-08-04-04.
    def test_previews_no_autofit_replacements_using_derived_source_bounds(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            source_path = root / "input" / "ja" / "no-autofit.pptx"
            _write_no_autofit_bounds_presentation(source_path)

            evaluate_text_replacement_examples(root / "input", root / "output")

            artifacts = root / "output" / "ja" / "no-autofit.text-layout-artifacts"
            no_autofit = json.loads(
                (artifacts / "text-box-0001.provider-0001.explicit.json").read_text(
                    encoding="utf-8"
                )
            )
            text_to_fit_shape = json.loads(
                (artifacts / "text-box-0002.provider-0001.explicit.json").read_text(
                    encoding="utf-8"
                )
            )
            self.assertTrue(no_autofit["fitting"]["derived_from_source"])
            self.assertFalse(text_to_fit_shape["fitting"]["derived_from_source"])
            self.assertEqual(
                no_autofit["width_emu"],
                no_autofit["fitting"]["rectangle"]["width_emu"],
            )
            self.assertGreater(
                no_autofit["fitting"]["rectangle"]["height_emu"],
                no_autofit["height_emu"],
            )
            self.assertEqual(
                text_to_fit_shape["width_emu"],
                text_to_fit_shape["fitting"]["rectangle"]["width_emu"],
            )
            self.assertEqual(
                text_to_fit_shape["height_emu"],
                text_to_fit_shape["fitting"]["rectangle"]["height_emu"],
            )
            source_properties = json.loads(
                (artifacts / "text-box-0001.json").read_text(encoding="utf-8")
            )
            self.assertEqual(1.0, source_properties["rendering"]["font_scale"])
            self.assertNotEqual(
                (artifacts / "text-box-0001.provider-0001.png").read_bytes(),
                (artifacts / "text-box-0002.provider-0001.png").read_bytes(),
            )

    # Verifies FR-2026-08-03-13.
    def test_explicit_properties_resolve_an_inherited_body_bullet(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            source_path = root / "input" / "ja" / "inherited-bullet.pptx"
            _write_inherited_bullet_presentation(source_path)

            evaluate_text_replacement_examples(root / "input", root / "output")

            artifacts = root / "output" / "ja" / "inherited-bullet.text-layout-artifacts"
            source_properties = json.loads(
                (artifacts / "text-box-0001.json").read_text(encoding="utf-8")
            )
            explicit_properties = json.loads(
                (artifacts / "text-box-0001.explicit.json").read_text(encoding="utf-8")
            )
            self.assertIsNone(source_properties["paragraphs"][0]["bullet_kind"])
            self.assertEqual("character", explicit_properties["paragraphs"][0]["bullet_kind"])
            self.assertIsNotNone(explicit_properties["paragraphs"][0]["bullet_marker"])

    # Verifies FR-2026-08-03-13.
    def test_reports_and_renders_presentation_text_boxes_without_ocr(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            input_root = root / "input"
            output_root = root / "output"
            source_path = input_root / "ja" / "nested" / "layout.pptx"
            _write_presentation(source_path)

            with patch("pipeline.ocr.OcrProviderFactory.discover_default_plugins") as discover_ocr:
                result = evaluate_text_replacement_examples(input_root, output_root)
            discover_ocr.assert_not_called()
            second_output_root = root / "second-output"
            evaluate_text_replacement_examples(input_root, second_output_root)

            page_path = output_root / "ja" / "nested" / "layout.html"
            page = page_path.read_text(encoding="utf-8")
            source_font_page_path = output_root / "ja" / "nested" / "layout.sf.html"
            source_font_page = source_font_page_path.read_text(encoding="utf-8")
            artifacts = output_root / "ja" / "nested" / "layout.text-layout-artifacts"
            first_rendering = artifacts / "text-box-0001.png"
            first_properties = artifacts / "text-box-0001.json"
            first_explicit_properties = artifacts / "text-box-0001.explicit.json"
            first_provider_rendering = artifacts / "text-box-0001.provider-0001.png"
            first_provider_properties = artifacts / "text-box-0001.provider-0001.explicit.json"
            second_rendering = artifacts / "text-box-0002.png"
            source_font_artifacts = output_root / "ja" / "nested" / "layout.sf"
            source_font_provider_properties = (
                source_font_artifacts / "text-box-0001.provider-0001.explicit.json"
            )
            self.assertEqual(1, result.processed_presentations)
            self.assertEqual(0, result.skipped_presentations)
            self.assertEqual(2, result.written_pages)
            self.assertEqual(2, result.rendered_text_boxes)
            self.assertIn("Native text-layout evaluation", page)
            self.assertIn("Native source-font text-layout evaluation", source_font_page)
            self.assertIn("<th>Original</th><th>argos_translate</th>", page)
            self.assertIn("main { padding: 2rem; width: max-content; }", page)
            self.assertIn('href="layout.text-layout-artifacts/text-box-0001.json" target="_blank"', page)
            self.assertIn(
                'href="layout.text-layout-artifacts/text-box-0001.explicit.json" target="_blank"',
                page,
            )
            self.assertNotIn("max-height", page)
            self.assertIn(
                ".rendering img { background: white; border: 1px solid #dc2626; display: block; "
                "height: auto; }",
                page,
            )
            self.assertIn("slide 1, shape 2/1", page)
            self.assertTrue(first_rendering.is_file())
            self.assertTrue(first_properties.is_file())
            self.assertTrue(first_explicit_properties.is_file())
            self.assertTrue(first_provider_rendering.is_file())
            self.assertTrue(first_provider_properties.is_file())
            self.assertTrue(second_rendering.is_file())
            self.assertTrue(source_font_artifacts.is_dir())
            self.assertTrue(source_font_provider_properties.is_file())
            properties = json.loads(first_properties.read_text(encoding="utf-8"))
            self.assertEqual(30.0, properties["shape_rotation_degrees"])
            self.assertEqual(30.0, properties["effective_text_rotation_degrees"])
            self.assertEqual(91_440, properties["margin_left_emu"])
            self.assertEqual("text-to-fit-shape", properties["autofit_mode"])
            self.assertEqual(90_000, properties["autofit_font_scale"])
            self.assertEqual(10_000, properties["autofit_line_spacing_reduction"])
            self.assertEqual("character", properties["paragraphs"][0]["bullet_kind"])
            self.assertEqual("•", properties["paragraphs"][0]["bullet_marker"])
            self.assertEqual("true", properties["paragraphs"][0]["runs"][0]["underline"])
            self.assertEqual(
                24.0, properties["paragraphs"][-1]["empty_line_font_size_points"]
            )
            self.assertEqual("fit", properties["rendering"]["fit_status"])
            self.assertEqual(1.0, properties["rendering"]["font_scale"])
            explicit_properties = json.loads(first_explicit_properties.read_text(encoding="utf-8"))
            self.assertEqual("none", explicit_properties["autofit_mode"])
            self.assertIsNone(explicit_properties["autofit_font_scale"])
            self.assertEqual(
                "Noto Serif JP", explicit_properties["paragraphs"][0]["runs"][0]["font_family"]
            )
            self.assertEqual(
                16.0, explicit_properties["paragraphs"][0]["runs"][0]["font_size_points"]
            )
            provider_properties = json.loads(first_provider_properties.read_text(encoding="utf-8"))
            source_font_provider = json.loads(
                source_font_provider_properties.read_text(encoding="utf-8")
            )
            self.assertEqual("argos_translate", provider_properties["replacement"]["provider"])
            self.assertEqual("ja", provider_properties["replacement"]["source_language"])
            self.assertEqual("en", provider_properties["replacement"]["target_language"])
            self.assertEqual("paragraph", provider_properties["replacement"]["unit"])
            self.assertEqual(
                "Noto Serif JP", provider_properties["paragraphs"][0]["runs"][0]["font_family"]
            )
            self.assertEqual(1, len(provider_properties["paragraphs"][0]["runs"]))
            self.assertIn("measurement_faces", source_font_provider)
            self.assertTrue(source_font_provider["measurement_faces"])
            self.assertEqual(
                "noto-fallback", source_font_provider["measurement_faces"][0]["source"]
            )
            self.assertEqual(
                "Example Serif", source_font_provider["measurement_faces"][0]["original_reference"]
            )
            self.assertEqual(
                "Example Serif", source_font_provider["measurement_faces"][0]["resolved_family"]
            )
            self.assertEqual(
                "latin", source_font_provider["measurement_faces"][0]["script"]
            )
            run_classifications = [
                run["font_classification"] for run in properties["paragraphs"][0]["runs"]
            ]
            self.assertEqual(["serif", "fixed-width"], run_classifications)
            self.assertEqual(
                first_rendering.read_bytes(),
                (
                    second_output_root
                    / "ja"
                    / "nested"
                    / "layout.text-layout-artifacts"
                    / "text-box-0001.png"
                ).read_bytes(),
            )
            with Image.open(first_rendering) as rendered:
                self.assertEqual((192, 96), rendered.size)
                self.assertIsNotNone(
                    ImageChops.difference(
                        rendered.convert("RGB"), Image.new("RGB", rendered.size, "white")
                    ).getbbox()
                )


def _write_presentation(source_path: Path) -> None:
    source_path.parent.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
    text_box.rotation = 30
    text_frame = text_box.text_frame
    text_frame.margin_left = Inches(0.1)
    text_frame.margin_right = Inches(0.1)
    text_frame.margin_top = Inches(0.05)
    text_frame.margin_bottom = Inches(0.05)
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    autofit = text_frame._element.bodyPr.find(qn("a:normAutofit"))
    assert autofit is not None
    autofit.set("fontScale", "90000")
    autofit.set("lnSpcReduction", "10000")
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.space_before = Pt(2)
    paragraph.space_after = Pt(3)
    paragraph_properties = paragraph._p.get_or_add_pPr()
    paragraph_properties.set("marL", "342900")
    paragraph_properties.set("indent", "-285750")
    bullet = OxmlElement("a:buChar")
    bullet.set("char", "•")
    paragraph_properties.append(bullet)
    first_run = paragraph.add_run()
    first_run.text = "Serif"
    first_run.font.name = "Example Serif"
    first_run.font.size = Pt(16)
    first_run.font.bold = True
    first_run.font.underline = True
    second_run = paragraph.add_run()
    second_run.text = " Mono"
    second_run.font.name = "Example Mono"
    second_run.font.size = Pt(12)
    second_run.font.italic = True
    second_run.font._element.set("baseline", "20000")
    empty_paragraph = text_frame.add_paragraph()
    empty_paragraph._p.get_or_add_endParaRPr().set("sz", "2400")

    group = slide.shapes.add_group_shape()
    grouped_text_box = group.shapes.add_textbox(Inches(0), Inches(0), Inches(1), Inches(0.5))
    grouped_text_box.text = "Grouped text"
    presentation.save(str(source_path))


def _write_empty_bullet_presentation(source_path: Path) -> None:
    source_path.parent.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    text_frame = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1)).text_frame
    paragraph_properties = text_frame.paragraphs[0]._p.get_or_add_pPr()
    bullet = OxmlElement("a:buChar")
    bullet.set("char", "•")
    paragraph_properties.append(bullet)
    text_frame.paragraphs[0]._p.get_or_add_endParaRPr().set("sz", "2400")
    presentation.save(str(source_path))


def _write_overflowing_presentation(source_path: Path) -> None:
    source_path.parent.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    text_frame = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(1.5), Inches(0.5)).text_frame
    text_frame.word_wrap = True
    paragraph = text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = "layout fitting uses a single scale for all source runs"
    run.font.size = Pt(24)
    presentation.save(str(source_path))


def _write_no_autofit_bounds_presentation(source_path: Path) -> None:
    source_path.parent.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    no_autofit = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(0.2))
    no_autofit.text = "Source text " * 40
    no_autofit.text_frame.paragraphs[0].runs[0].font.size = Pt(24)
    no_autofit.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    text_to_fit_shape = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(0.2))
    text_to_fit_shape.text = "Source text " * 40
    text_to_fit_shape.text_frame.paragraphs[0].runs[0].font.size = Pt(24)
    text_to_fit_shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    presentation.save(str(source_path))


def _write_inherited_bullet_presentation(source_path: Path) -> None:
    source_path.parent.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()
    layout = presentation.slide_layouts[2]
    layout_body_placeholder = cast(
        Shape,
        next(
            shape
            for shape in cast(Iterable[BaseShape], layout.placeholders)
            if shape.placeholder_format.type == PP_PLACEHOLDER.BODY
        ),
    )
    list_style = layout_body_placeholder.text_frame._element.find(qn("a:lstStyle"))
    assert list_style is not None
    level_properties = list_style.find(qn("a:lvl1pPr"))
    if level_properties is None:
        level_properties = OxmlElement("a:lvl1pPr")
        list_style.append(level_properties)
    no_bullet = level_properties.find(qn("a:buNone"))
    if no_bullet is not None:
        level_properties.remove(no_bullet)
    bullet = OxmlElement("a:buChar")
    bullet.set("char", "•")
    level_properties.append(bullet)
    slide = presentation.slides.add_slide(layout)
    body_placeholder = next(
        shape
        for shape in slide.placeholders
        if shape.placeholder_format.type == PP_PLACEHOLDER.BODY
    )
    body_placeholder.text = "Inherited list formatting"
    presentation.save(str(source_path))
